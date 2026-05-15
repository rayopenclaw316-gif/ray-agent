#!/usr/bin/env python3
"""繼續製作 nihms-986955 / J Neural Eng 2018 的 PPT，新增結果、討論、結論投影片"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy, os

SRC = "/Users/rayopenclaw/Downloads/發展用於靜默語音辨識的表面肌.pptx"
DST = "/Users/rayopenclaw/Downloads/發展用於靜默語音辨識的表面肌.pptx"

FONT_SIZE = Pt(24)  # 304800 EMU = 24pt

def add_slide(prs, title_text, body_lines):
    """用 '標題及內容' layout 新增一張投影片"""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title_ph = slide.shapes.title
    body_ph = slide.placeholders[1]

    title_ph.text = title_text

    tf = body_ph.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        if para.runs:
            para.runs[0].font.size = FONT_SIZE

    return slide


def update_toc(prs):
    """填入目錄投影片（第2張，index=1）"""
    slide = prs.slides[1]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name != "標題 1":
            tf = shape.text_frame
            tf.clear()
            sections = [
                "1. 先前研究",
                "2. 方法",
                "   2-1 受試者與設備",
                "   2-2 電極貼片位置",
                "   2-3 語料庫蒐集",
                "   2-4 語音活動偵測",
                "   2-5 特徵選擇",
                "   2-6 語法模型選擇",
                "3. 結果",
                "   3-1 孤立詞辨識",
                "   3-2 音素辨識調參",
                "   3-3 最終系統",
                "4. 討論",
                "5. 結論",
            ]
            for i, sec in enumerate(sections):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = sec
                if para.runs:
                    para.runs[0].font.size = Pt(20)
            break


prs = Presentation(SRC)

# 更新目錄
update_toc(prs)

# ── 結果：孤立詞辨識 ──────────────────────────────────────────────
add_slide(prs, "結果－孤立詞辨識", [
    "測試對象：語料庫1（65個孤立詞，9位受試者）",
    "",
    "各特徵平均詞錯誤率（WER）比較：",
    "• MFCC（梅爾頻率倒譜系數）：9.6%  ← 最佳",
    "• 共激活指數（Co-activation Index）：41.2%（MFCC的4倍以上）",
    "• 均方根RMS、主頻成分、自協方差振幅範圍等：均不及MFCC",
    "",
    "結論：MFCC 是 sEMG 靜默語音辨識最適合的特徵",
    "→ 沿用於後續所有辨識實驗",
])

# ── 結果：音素辨識調參 ────────────────────────────────────────────
add_slide(prs, "結果－音素辨識調參", [
    "測試對象：語料庫2（202個詞彙，4位受試者）",
    "",
    "▸ 調整高斯混合數量（Gaussian Mixtures per HMM State）",
    "  4混合 → WER ≈ 24%",
    "  8混合 → WER ≈ 15%  （大幅改善）",
    "  16混合 → WER = 11.3%  ← 最佳，沿用",
    "",
    "▸ HLDA 降維（從154維特徵壓縮）",
    "  預設154維 → WER 11.3%",
    "  40維        → WER 7.3%  ← 最佳，沿用",
    "  維度繼續縮減不再帶來改善",
    "",
    "→ 最終設定：16混合 + 40維HLDA，WER = 7.3%",
])

# ── 結果：最終系統 ────────────────────────────────────────────────
add_slide(prs, "結果－最終系統", [
    "測試對象：語料庫3（2200詞彙，1200句，6位受試者）",
    "",
    "最終演算法組合（從 HTK 遷移至 KALDI）：",
    "• 三音素模型（Triphone）：共享HMM參數以辨識未見三音素組合",
    "• MLLR（最大似然線性回歸）：對GMM均值做線性轉換壓縮參數",
    "• SGMM（子空間高斯混合模型）：所有音素狀態共用參數子空間",
    "• HLDA 40維 + MFCC特徵",
    "• 感測器從11個縮減至8個（不顯著影響準確率）",
    "",
    "最終成績：",
    "  平均 WER = 8.9%（辨識率 91.1%）",
    "  最佳受試者 WER = 1.4%",
    "  特殊操作語料：7.5%　常用短句語料：5.1%",
])

# ── 討論 ─────────────────────────────────────────────────────────
add_slide(prs, "討論", [
    "▸ 特徵選擇",
    "  MFCC 原為聲學語音設計，卻也最適合 sEMG 訊號",
    "  → 暗示 sEMG 訊號與聲學訊號共享頻譜結構",
    "",
    "▸ 語法模型",
    "  NL等效語法 在準確率（WER 6.8%）與彈性之間取得最佳平衡",
    "  → 可適用於更廣泛的日常語句，不限特定句型",
    "",
    "▸ 音素辨識的突破",
    "  從詞彙層級 → 音素層級，可辨識訓練時未見過的詞彙",
    "  → 大幅降低使用者的資料蒐集負擔",
    "",
    "▸ 限制與未來方向",
    "  仍需受試者特定訓練資料（hours級別）",
    "  未來：深度學習（DNN/LSTM）+ 大規模跨受試者資料",
    "  → 可用少量個人資料微調大型基礎模型（如聲學ASR的做法）",
])

# ── 結論 ─────────────────────────────────────────────────────────
add_slide(prs, "結論", [
    "本研究建立了完整的 sEMG 靜默語音辨識（SSR）系統，達成三項核心突破：",
    "",
    "1. 特徵選擇：MFCC 在孤立詞辨識中 WER = 9.6%，遠優於其他特徵",
    "",
    "2. 語法模型：NL等效語法 WER = 6.8%，兼顧準確率與語句彈性",
    "",
    "3. 音素辨識：三音素 + MLLR + SGMM + HLDA",
    "   → 2200詞彙，1200連續短句，WER = 8.9%（最佳受試者1.4%）",
    "",
    "應用潛力：",
    "  • 喉切除術患者的輔助溝通裝置",
    "  • 軍事人員的免持隱密通訊",
    "  • 公共場所的私人語音控制介面",
])

prs.save(DST)
print(f"完成！共 {len(prs.slides)} 張投影片，已存回 {DST}")
