#!/usr/bin/env python3
"""
重建 nihms-986955 PPT：
- 保留原有 slides 1-13（只給 12、13 加圖）
- 刪除舊的 slides 14-16
- 新增 14a-14g + 15討論 + 16結論
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, os

SRC = "/Users/rayopenclaw/Downloads/發展用於靜默語音辨識的表面肌.pptx"
DST = "/Users/rayopenclaw/Downloads/發展用於靜默語音辨識的表面肌.pptx"
FIGS = "/Users/rayopenclaw/ray-agent/papers/nihms_figs"

FONT_SIZE = Pt(18)
TITLE_FONT = Pt(28)

# ── 輔助函式 ────────────────────────────────────────────────────────

def add_slide(prs, title_text, body_lines, img_path=None, img_top_ratio=0.55):
    """新增 '標題及內容' 投影片，可選擇在下半部插入圖片"""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title_ph = slide.shapes.title
    body_ph = slide.placeholders[1]

    title_ph.text = title_text
    for run in title_ph.text_frame.paragraphs[0].runs:
        run.font.size = TITLE_FONT

    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, line in enumerate(body_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        if para.runs:
            para.runs[0].font.size = FONT_SIZE

    if img_path and os.path.exists(img_path):
        sw = prs.slide_width
        sh = prs.slide_height
        top = Emu(int(sh * img_top_ratio))
        left = Inches(0.5)
        width = sw - Inches(1.0)
        height = sh - top - Inches(0.2)
        slide.shapes.add_picture(img_path, left, top, width, height)

        # 縮短文字佔位框高度，讓圖片不被遮住
        body_ph.height = top - Inches(0.1)

    return slide


def add_slide_with_image_right(prs, title_text, body_lines, img_path):
    """文字左半、圖片右半的版面"""
    layout = prs.slide_layouts[5]  # 只有標題
    slide = prs.slides.add_slide(layout)

    sw = prs.slide_width
    sh = prs.slide_height

    title_ph = slide.shapes.title
    title_ph.text = title_text
    for run in title_ph.text_frame.paragraphs[0].runs:
        run.font.size = TITLE_FONT

    # 文字方塊（左半）
    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(1.4), Inches(5.8), sh - Inches(1.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        if para.runs:
            para.runs[0].font.size = FONT_SIZE

    # 圖片（右半）
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(6.3), Inches(1.4),
            Inches(6.5), sh - Inches(1.8)
        )

    return slide


def add_slide_full(prs, title_text, body_lines):
    """純文字，不含圖片"""
    return add_slide(prs, title_text, body_lines, img_path=None)


def delete_slides_from(prs, start_idx):
    """刪除從 start_idx 開始的所有投影片"""
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sid in slide_ids[start_idx:]:
        xml_slides.remove(sid)


# ── 主程式 ──────────────────────────────────────────────────────────

prs = Presentation(SRC)

# 刪除舊的 slide 14、15、16（index 13、14、15）
delete_slides_from(prs, 13)

# ── 給 slide 12（孤立詞辨識, index 11）加圖 ─────────────────────────
slide12 = prs.slides[11]
sw = prs.slide_width
sh = prs.slide_height
img_path = f"{FIGS}/fig2_crop.png"
if os.path.exists(img_path):
    top = Inches(3.8)
    slide12.shapes.add_picture(img_path, Inches(1.5), top, sw - Inches(3.0), sh - top - Inches(0.2))

# ── 給 slide 13（音素調參, index 12）加兩張圖 ────────────────────────
slide13 = prs.slides[12]
img4 = f"{FIGS}/fig4_crop.png"
img5 = f"{FIGS}/fig5_crop.png"
top13 = Inches(3.8)
half_w = (sw - Inches(1.5)) / 2
if os.path.exists(img4):
    slide13.shapes.add_picture(img4, Inches(0.5), top13, half_w, sh - top13 - Inches(0.2))
if os.path.exists(img5):
    slide13.shapes.add_picture(img5, Inches(0.5) + half_w + Inches(0.3), top13, half_w, sh - top13 - Inches(0.2))

# ── 新增 slide 14a：最終演算法整體流程 ──────────────────────────────
add_slide(prs, "最終演算法整體流程", [
    "【輸入】11通道 sEMG 訊號",
    "    ↓",
    "【特徵提取】MFCC → 154維特徵向量（11通道 × 7MFCC + 7delta）",
    "    ↓",
    "【降維】HLDA → 壓縮至 40 維",
    "    ↓",
    "【第一步建模】單音素 HMM（Monophone HMM）",
    "    ↓",
    "【強制對齊】Forced Alignment → 找出每個音素的精確時間邊界",
    "    ↓",
    "【加強模型】逐步增加 Gaussian Mixtures（最多16混合）",
    "    ↓",
    "【遷移工具】HTK → KALDI",
    "    ↓",
    "【升級模型】三音素 HMM（Triphone）+ Decision Tree Clustering",
    "    ↓",
    "【壓縮參數】MLLR 線性轉換",
    "    ↓",
    "【共享結構】SGMM 子空間高斯混合",
    "    ↓",
    "【輸出】辨識結果（2200詞彙，WER 8.9%）",
    "",
    "感測器：11個 → 8個，準確率無顯著下降",
], img_path=f"{FIGS}/fig1_crop.png", img_top_ratio=0.62)

# ── 新增 slide 14b：Forced Alignment ────────────────────────────────
add_slide_full(prs, "Forced Alignment 強制對齊", [
    "▌ 是什麼",
    "自動標記技術：告訴系統「這句話裡每個音素從幾秒到幾秒」",
    "",
    "▌ 為什麼需要",
    "HMM 初次訓練時只知道整句話的文字，不知道每個音素的確切時間邊界",
    "邊界模糊 → 訓練資料不精確 → 模型效果差",
    "",
    "▌ 怎麼做（Viterbi 演算法）",
    "1. 給定詞序列，例如 'STOP RIGHT NOW'",
    "2. 轉成音素序列：/s/-/t/-/ɒ/-/p/-/r/-/aɪ/-/t/-/n/-/aʊ/",
    "3. Viterbi 在整段 sEMG 訊號上找最大機率路徑：",
    "   δₜ(i) = max P(O₁...Oₜ , 狀態=i)",
    "   δ：到時間 t、走到狀態 i 的最大機率",
    "   → 從頭掃到尾，記錄每步最優來源，回溯得到完整路徑",
    "4. 結果：每個音素得到精確起止時間標記",
    "",
    "▌ 效果",
    "精確標記 → 重新訓練 HMM → 模型更準確",
])

# ── 新增 slide 14c：Gaussian Mixtures ───────────────────────────────
add_slide(prs, "Gaussian Mixtures 高斯混合", [
    "▌ 問題：HMM 每個狀態要描述「此狀態下特徵向量的分布」",
    "",
    "▌ 單一高斯（Single Gaussian）",
    "只能描述單峰（橢圓形）分布",
    "同一音素在不同語境下 sEMG 形狀差異很大（多峰）→ 一個高斯不夠",
    "",
    "▌ 高斯混合模型（GMM）：用 K 個高斯加權疊加",
    "   p(x) = Σₖ wₖ × N(x ; μₖ , Σₖ)",
    "   wₖ：第k個高斯的權重（所有wₖ加總=1）",
    "   μₖ：第k個高斯的中心位置（均值向量）",
    "   Σₖ：第k個高斯的形狀（協方差矩陣）",
    "",
    "▌ 混合數越多越好？",
    "多 → 描述更複雜分布，但需要更多訓練資料",
    "本研究：4混合WER 24% → 8混合 15% → 16混合 11.3%（最佳）",
], img_path=f"{FIGS}/fig4_crop.png", img_top_ratio=0.60)

# ── 新增 slide 14d：Triphone + KALDI ─────────────────────────────────
add_slide_full(prs, "Triphone 三音素模型 + KALDI", [
    "▌ 為什麼要升級到三音素",
    "單音素 HMM 假設：同一音素不管前後文都一樣 ← 錯誤",
    "例：/t/ 在「s-t-op」vs「t-op」的臉部肌電訊號完全不同",
    "三音素：/s-t+ɒ/ = 「前面是s，當前是t，後面是ɒ」",
    "",
    "▌ 三音素組合爆炸的問題",
    "37個音素 × 37 × 37 ≈ 50,000 種組合，訓練資料根本不夠",
    "",
    "▌ 解法：Decision Tree Clustering（決策樹分群）",
    "用語音學問題分群：「前面音素是不是鼻音？」「當前是不是母音？」",
    "相似的三音素合併共享 HMM 參數",
    "未見過的三音素 → 自動分配到最相似的已知三音素",
    "",
    "▌ 為什麼從 HTK 遷移到 KALDI",
    "HTK：不支援 Triphone + Decision Tree + SGMM",
    "KALDI：Johns Hopkins University 開源工具包",
    "    支援 Triphone、MLLR、SGMM 等全套現代方法",
    "    腳本化 pipeline，可重現、易擴充",
])

# ── 新增 slide 14e：MLLR ──────────────────────────────────────────────
add_slide_full(prs, "MLLR 最大似然線性回歸", [
    "▌ 問題",
    "11通道 sEMG × MFCC 特徵 → GMM 每個狀態有大量參數",
    "單一高斯全協方差矩陣 → >10,000 個參數，訓練資料不夠",
    "",
    "▌ MLLR 做什麼",
    "對 GMM 的均值（μ）做線性變換，用少量矩陣取代大量個別參數",
    "",
    "▌ 公式",
    "   μ̂ = W × [μ ; 1]",
    "   μ：原始 GMM 均值向量",
    "   W：共用變換矩陣（用最大似然法從資料估算）",
    "   μ̂：變換後的新均值，替換原始均值",
    "   [μ ; 1]：在向量後加1，讓 W 同時包含旋轉與平移",
    "",
    "▌ 為什麼有效",
    "不同通道的 sEMG 訊號存在相關性（嘴角動，臉頰也動）",
    "線性變換能捕捉跨通道關係，用更少參數表達同樣資訊",
    "",
    "▌ 好處",
    "參數量大幅減少，小量訓練資料也能訓練穩健模型",
])

# ── 新增 slide 14f：SGMM ──────────────────────────────────────────────
add_slide_full(prs, "SGMM 子空間高斯混合模型", [
    "▌ 問題：就算用了 MLLR，每個音素狀態仍各自有一套 GMM，總參數仍龐大",
    "",
    "▌ SGMM 核心想法：所有音素狀態共用一個高斯混合池",
    "每個狀態只描述自己在公共空間中的「偏移量」",
    "",
    "▌ 類比",
    "全校學生共用同一本字典（公共子空間）",
    "每個人只需額外記下「我和字典的差異」，不用各自從頭背一本字典",
    "",
    "▌ 公式",
    "   μⱼₘ = M × vⱼ",
    "   M：共用子空間矩陣（所有音素狀態共享）",
    "   vⱼ：第 j 個狀態的低維「個性向量」",
    "   μⱼₘ：第 j 個狀態、第 m 個高斯的均值",
    "",
    "▌ 訓練流程",
    "1. 用標準 GMM 初始化",
    "2. EM 演算法交替更新 M（公共結構）與 vⱼ（各狀態個性）",
    "3. 收斂後：每個狀態均值 = 公共矩陣 × 低維向量",
    "",
    "▌ 好處",
    "總參數 = M（大但共用）+ 少量vⱼ × 狀態數，遠小於各自獨立的 GMM",
    "資料少也能訓練，同時保持各狀態的獨特性",
])

# ── 新增 slide 14g：最終成績 ──────────────────────────────────────────
add_slide(prs, "最終成績", [
    "最終系統：Triphone + HLDA(40維) + 16混合 + MLLR + SGMM + 8感測器",
    "",
    "測試：語料庫3，2200詞彙，超過1200個連續短句",
    "",
    "★ 平均 WER = 8.9%（辨識率 91.1%）",
    "★ 最佳受試者 WER = 1.4%",
    "",
    "各語料集成績：",
    "   數字語料      WER 13.8%",
    "   簡訊語料      WER  9.0%",
    "   特殊操作語料  WER  7.5%",
    "   常用短句語料  WER  5.1%",
    "",
    "比較：前人最佳成績 WER 15.3%（本研究比前人好 42%）",
], img_path=f"{FIGS}/table2_crop.png", img_top_ratio=0.58)

# ── 新增 slide 15：討論 ───────────────────────────────────────────────
add_slide_full(prs, "討論", [
    "▌ 特徵選擇",
    "MFCC 原為聲學語音設計，卻也最適合 sEMG 訊號",
    "→ 暗示 sEMG 與聲學訊號共享頻譜結構",
    "",
    "▌ 語法模型",
    "NL等效語法在準確率（WER 6.8%）與彈性之間取得最佳平衡",
    "→ 可適用於更廣泛的日常語句，不限特定句型",
    "",
    "▌ 音素辨識的突破",
    "詞彙層級 → 音素層級，可辨識訓練時未見過的詞彙",
    "→ 大幅降低使用者的資料蒐集負擔",
    "",
    "▌ 限制與未來方向",
    "仍需受試者特定訓練資料（hours 級別）",
    "跨受試者模型需要大量多樣化資料",
    "未來：深度學習（DNN / LSTM）+ 大規模跨受試者訓練資料",
    "→ 少量個人資料微調大型基礎模型（類似聲學 ASR 的做法）",
])

# ── 新增 slide 16：結論 ───────────────────────────────────────────────
add_slide_full(prs, "結論", [
    "本研究建立了完整的 sEMG 靜默語音辨識（SSR）系統，達成三項核心突破：",
    "",
    "1. 特徵選擇",
    "   MFCC 在孤立詞辨識中 WER = 9.6%，遠優於其他特徵",
    "",
    "2. 語法模型",
    "   NL 等效語法 WER = 6.8%，兼顧準確率與語句彈性",
    "",
    "3. 音素辨識（最終系統）",
    "   Triphone + MLLR + SGMM + HLDA",
    "   → 2200詞彙，1200連續短句，WER = 8.9%（最佳 1.4%）",
    "",
    "應用潛力：",
    "   • 喉切除術患者的輔助溝通裝置",
    "   • 軍事人員的免持隱密通訊",
    "   • 公共場所的私人語音控制介面",
])

prs.save(DST)
print(f"完成！共 {len(prs.slides)} 張投影片，已存回：{DST}")
