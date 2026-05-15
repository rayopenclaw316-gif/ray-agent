#!/usr/bin/env python3
"""
SSR sEMG Paper PPT - Biomedical Signal Processing and Control 2024
Style: white background, #0E2841 titles, red technical terms, Times New Roman + 標楷體
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import re, os

IMG_DIR = os.path.join(os.path.dirname(__file__), 'papers', 'S1746809424001101_imgs')

TITLE_COLOR = RGBColor(0x0E, 0x28, 0x41)
RED_COLOR   = RGBColor(0xFF, 0x00, 0x00)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY        = RGBColor(0x66, 0x66, 0x66)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE  = RGBColor(0xF0, 0xF4, 0xF8)

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)
TITLE_L, TITLE_T, TITLE_W, TITLE_H         = Cm(2.3), Cm(1.0), Cm(29.2), Cm(3.7)
CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H = Cm(2.3), Cm(5.1), Cm(29.2), Cm(12.1)
PGNUM_L, PGNUM_T, PGNUM_W, PGNUM_H         = Cm(23.9), Cm(17.7), Cm(3.0), Cm(1.0)

FONT_EN = "Times New Roman"
FONT_ZH = "標楷體"

RED_TERMS = sorted([
    "sEMG", "EMG", "EEG", "GRU", "LSTM", "BiGRU", "CNN", "RNN", "MLP", "ResNet",
    "SSR", "SSI", "Transformer", "STFT", "Fbank", "fbank",
    "Silent Speech", "silent speech", "Transfer Learning", "transfer learning",
    "Data Augmentation", "data augmentation", "Speed Perturbation", "speed perturbation",
    "Time-Frequency Masking", "time-frequency masking",
    "Zero-Crossing Rate", "zero-crossing rate", "ZCR",
    "Filter Bank", "filter bank", "Mel", "ParallelNet",
], key=len, reverse=True)
RED_SET = set(RED_TERMS)


def _is_zh(c):
    return '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f'


def _set_fonts(run, font_name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
        el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', font_name)


def _add_run(para, text, font_name, size, bold=False, color=BLACK):
    r = para.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_fonts(r, font_name)
    return r


def _split_lang(text):
    if not text:
        return []
    segs, cur, cur_zh = [], text[0], _is_zh(text[0])
    for ch in text[1:]:
        zh = _is_zh(ch)
        if zh == cur_zh:
            cur += ch
        else:
            segs.append((cur, cur_zh))
            cur, cur_zh = ch, zh
    segs.append((cur, cur_zh))
    return segs


def add_text(para, text, size=18, bold=False, color=None, highlight=True):
    if color is not None:
        for seg, is_zh in _split_lang(text):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, color)
        return
    if not highlight:
        for seg, is_zh in _split_lang(text):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, BLACK)
        return
    pattern = '(' + '|'.join(re.escape(t) for t in RED_TERMS) + ')'
    for part in re.split(pattern, text):
        if not part:
            continue
        c = RED_COLOR if part in RED_SET else BLACK
        for seg, is_zh in _split_lang(part):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, c)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_title(slide, text, size=28):
    tf = _tb(slide, TITLE_L, TITLE_T, TITLE_W, TITLE_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_text(p, text, size, bold=True, color=TITLE_COLOR)


def add_pgnum(slide, n):
    tf = _tb(slide, PGNUM_L, PGNUM_T, PGNUM_W, PGNUM_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, str(n), FONT_EN, 16, color=BLACK)


def add_content(slide, items, base=18,
                l=CONTENT_L, t=CONTENT_T, w=CONTENT_W, h=CONTENT_H):
    tf = _tb(slide, l, t, w, h)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(item, str):
            p.space_before = Pt(5)
            add_text(p, item, base)
            continue
        text   = item.get('text', '')
        size   = item.get('size', base)
        bold   = item.get('bold', False)
        color  = item.get('color', None)
        bullet = item.get('bullet', True)
        sp     = item.get('space_before', 5)
        p.space_before = Pt(sp)
        prefix = ''
        if bullet and text and text[:1] not in ('•', '①','②','③','④','⑤','—','│',' '):
            prefix = '• '
        if prefix:
            _add_run(p, prefix, FONT_EN, size, False, BLACK)
        if color is not None:
            add_text(p, text, size, bold, color=color)
        else:
            add_text(p, text, size, bold)


def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)


def make_slide(prs, num, title, items, title_size=28, base=18):
    s = _blank(prs)
    add_title(s, title, title_size)
    add_content(s, items, base)
    add_pgnum(s, num)
    return s


def make_split_slide(prs, num, title, items, img_path, base=17,
                     img_w=Cm(14.0), title_size=28):
    """左文右圖 split layout"""
    s = _blank(prs)
    add_title(s, title, title_size)
    text_w = CONTENT_W - img_w - Cm(0.5)
    add_content(s, items, base, l=CONTENT_L, t=CONTENT_T, w=text_w, h=CONTENT_H)
    img_l = CONTENT_L + text_w + Cm(0.5)
    img_h = CONTENT_H
    add_image(s, img_path, img_l, CONTENT_T, img_w, img_h)
    add_pgnum(s, num)
    return s


def make_fullimg_slide(prs, num, title, img_path, caption='', title_size=28):
    """標題下方全寬圖片"""
    s = _blank(prs)
    add_title(s, title, title_size)
    img_h = Cm(11.5)
    img_t = Cm(5.2)
    add_image(s, img_path, CONTENT_L, img_t, CONTENT_W, img_h)
    if caption:
        tf = _tb(s, CONTENT_L, img_t + img_h + Cm(0.1), CONTENT_W, Cm(0.8))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, caption, FONT_ZH if any(_is_zh(c) for c in caption) else FONT_EN,
                 14, color=GRAY)
    add_pgnum(s, num)
    return s


# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: 封面 ─────────────────────────────────────────────────────
    s1 = _blank(prs)
    tf1 = _tb(s1, Cm(2.3), Cm(1.8), Cm(29.2), Cm(8.5))
    for i, (line, sz, bd, col) in enumerate([
        ("Design and Implementation of a Silent Speech Recognition",  28, True, TITLE_COLOR),
        ("System Based on sEMG Signals: A Neural Network Approach",   28, True, TITLE_COLOR),
        ("基於表面肌電圖訊號的無聲語音辨識系統設計與實作：神經網路方法", 20, True, TITLE_COLOR),
    ]):
        p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6 if i == 2 else 0)
        fn = FONT_ZH if any(_is_zh(c) for c in line) else FONT_EN
        _add_run(p, line, fn, sz, bd, col)

    tf2 = _tb(s1, Cm(2.3), Cm(11.0), Cm(29.2), Cm(3.5))
    for i, (line, sz) in enumerate([
        ("Bokai Huang, Yizi Shao, Hao Zhang et al.", 16),
        ("Aerospace Information Research Institute, Chinese Academy of Sciences", 14),
        ("Biomedical Signal Processing and Control  |  Vol. 92, June 2024  |  106052", 14),
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        _add_run(p, line, FONT_EN, sz, False, BLACK)

    tf3 = _tb(s1, Cm(2.3), Cm(15.2), Cm(29.2), Cm(1.5))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(6)
    add_text(p3, "報告人：陳睿莆　|　國立虎尾科技大學　|　DOI: 10.1016/j.bspc.2024.106052", 14, color=GRAY)
    add_pgnum(s1, 1)

    # ── Slide 2: 目錄 ─────────────────────────────────────────────────────
    s2 = _blank(prs)
    add_title(s2, "目錄 / Table of Contents")
    toc = [
        ("研究背景與動機",                          3),
        ("系統整體流程概觀",                         4),
        ("語料庫設計與資料採集",                     5),
        ("電極配置與訊號採集裝置",                   6),
        ("訊號前處理",                              7),
        ("時頻混合特徵提取",                         8),
        ("資料增強策略",                             9),
        ("神經網路模型架構比較",                     10),
        ("GRU 模型詳解",                            11),
        ("實驗設定與評估方法",                       12),
        ("實驗結果：多受試者混合分類",               13),
        ("實驗結果：單一受試者分類 + 混淆矩陣",       14),
        ("跨受試者與跨期問題",                       15),
        ("結論與對本研究的啟發",                     16),
    ]
    tf2 = _tb(s2, CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H)
    first = True
    for label, pg in toc:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_before = Pt(5)
        add_text(p, f"{label}", 17)
        _add_run(p, f"  {'·'*3}  {pg}", FONT_EN, 17, color=GRAY)
    add_pgnum(s2, 2)

    # ── Slide 3: 研究背景與動機 ───────────────────────────────────────────
    make_slide(prs, 3, "研究背景與動機", [
        {"text": "什麼是 Silent Speech Recognition（SSR）？", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "不依賴聲帶振動，透過分析說話時產生的生理訊號，重建語音內容", "size": 17, "space_before": 3},
        {"text": "應用場景：語言障礙輔助溝通、嘈雜/安靜環境通訊、人機互動、物聯網", "size": 17, "space_before": 3},
        {"text": "為何選擇 sEMG？", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "非侵入性、高時間解析度、使用方便、技術成熟", "size": 17, "space_before": 3},
        {"text": "相比 ECoG（侵入腦表面）、磁感測器（口腔內置磁鐵）、唇語攝影（需固定角度）更實用", "size": 17, "space_before": 3},
        {"text": "現有挑戰：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "早期研究語料小（10–33 類），無法滿足日常溝通需求", "size": 17, "space_before": 3},
        {"text": "sEMG 訊號長度不固定，傳統機器學習方法難以有效處理時序資訊", "size": 17, "space_before": 3},
        {"text": "跨受試者與跨期個體差異大，模型泛化能力不足", "size": 17, "space_before": 3},
        {"text": "本文貢獻：設計 135 類中文語料庫，提出時頻混合特徵 + 資料增強 + GRU 深度學習系統",
         "bold": True, "bullet": False, "size": 17, "space_before": 10, "color": RED_COLOR},
    ])

    # ── Slide 4: 系統流程概觀（全圖） ─────────────────────────────────────
    make_fullimg_slide(prs, 4, "系統整體流程概觀",
                       os.path.join(IMG_DIR, 'p03_img00.jpeg'),
                       caption="Fig. 1　資料採集 → 訊號前處理 → 特徵提取 → 資料增強（×4）→ 神經網路模型 → 辨識結果")

    # ── Slide 5: 語料庫設計 ──────────────────────────────────────────────
    make_slide(prs, 5, "語料庫設計與資料採集", [
        {"text": "語料庫設計（最大規模中文 SSR 語料庫）：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "共 135 類詞語與短句，涵蓋 11 個語義類別", "size": 17, "space_before": 3},
        {"text": "類別：描述詞、程度詞、身體部位、否定詞、問候語、請求、回應、疑問、物品、位置、時間",
         "size": 16, "space_before": 2},
        {"text": "大多數類別為 1–2 個漢字，最長不超過 4 個漢字（符合口語需求）", "size": 17, "space_before": 3},
        {"text": "受試者與採集規格：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "12 位受試者（7 男 5 女，年齡 24.3 ± 1.2 歲），母語均為普通話", "size": 17, "space_before": 3},
        {"text": "每個詞語每位受試者錄製 20 次 → 總計 12 × 135 × 20 = 32,400 筆資料", "size": 17, "space_before": 3},
        {"text": "分兩天收集（每天兩段），避免肌肉疲勞影響訊號品質", "size": 17, "space_before": 3},
        {"text": "採集裝置：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "8 通道 sEMG，24-bit ADC（ADS1298IPAGR，Texas Instruments）", "size": 17, "space_before": 3},
        {"text": "取樣率 1000 Hz，SNR 38–50 dB，背景雜訊 2.4–4.7 μV", "size": 17, "space_before": 3},
        {"text": "採集方式：受試者默念詞語，以 Enter 鍵觸發切換，系統自動記錄起止時間戳記", "size": 17, "space_before": 3},
    ])

    # ── Slide 6: 電極配置（左文右解剖圖） ────────────────────────────────
    make_split_slide(prs, 6, "電極配置與肌肉選取", [
        {"text": "7 個肌群位置（右圖對應編號）：", "bold": True, "bullet": False, "size": 17, "space_before": 4},
        {"text": "① Zygomaticus minor（顴小肌）", "size": 16, "bullet": False, "space_before": 4},
        {"text": "② Orbicularis oris（口輪匝肌）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "③ Levator anguli oris（提口角肌）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "④ Mentalis（頦肌）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "⑤ Depressor anguli oris（降口角肌）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "⑥ Masseter（咬肌）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "⑦ Mylohyoideus（下頜舌骨肌，左右各一）", "size": 16, "bullet": False, "space_before": 3},
        {"text": "特殊設計：", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "刻意避開喉部區域（無聲狀態無肌肉活動 + 吞嚥動作干擾）", "size": 16, "space_before": 3},
        {"text": "驅動電極貼前額正中，參考電極貼右耳後乳突（mastoid process）", "size": 16, "space_before": 3},
    ],
    os.path.join(IMG_DIR, 'p04_img01.jpeg'),
    img_w=Cm(12.0))

    # ── Slide 7: 訊號前處理（左文右訊號圖） ──────────────────────────────
    make_split_slide(prs, 7, "訊號前處理", [
        {"text": "三步驟前處理流程：", "bold": True, "bullet": False, "size": 17, "space_before": 4},
        {"text": "① 高通濾波（High-pass Filter）：截止頻率 0.05 Hz，移除基線漂移（Baseline Drift）",
         "size": 16, "space_before": 5},
        {"text": "② 陷波濾波（Notch Filter）：50 Hz，移除電源頻率干擾（Power Frequency Noise）",
         "size": 16, "space_before": 4},
        {"text": "③ Z-score 標準化：每通道獨立進行，統一訊號量綱",
         "size": 16, "space_before": 4},
        {"text": "右圖說明：", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "8 通道前處理後訊號（範例詞語「我沒聽清楚」）", "size": 16, "space_before": 3},
        {"text": "Channel 2、3 訊號最強，對應顴小肌與口輪匝肌（說話主力肌群）", "size": 16, "space_before": 3},
        {"text": "各通道訊號形態有差異，提供互補資訊供模型學習", "size": 16, "space_before": 3},
    ],
    os.path.join(IMG_DIR, 'p04_img00.jpeg'),
    img_w=Cm(15.0))

    # ── Slide 8: 特徵提取 ────────────────────────────────────────────────
    make_slide(prs, 8, "時頻混合特徵提取", [
        {"text": "訊號分幀（Framing）：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "視窗長度 25 ms、步進 10 ms（確保短時平穩性且保留時間資訊）", "size": 17, "space_before": 3},
        {"text": "時域特徵（5 維/通道）：", "bold": True, "bullet": False, "size": 18, "space_before": 7},
        {"text": "短時能量（Short-Time Energy）：每幀訊號能量強度", "size": 17, "space_before": 3},
        {"text": "Zero-Crossing Rate（ZCR，過零率）：訊號頻率特性指標", "size": 17, "space_before": 3},
        {"text": "頻域特徵（40 維/通道）：", "bold": True, "bullet": False, "size": 18, "space_before": 7},
        {"text": "Fbank（Filter Bank）特徵：借用語音識別中的 Mel 濾波器組", "size": 17, "space_before": 3},
        {"text": "流程：分幀 → 加窗 → STFT → Mel 濾波 → 對數壓縮，映射非線性低頻能量", "size": 17, "space_before": 3},
        {"text": "整合後：(5 + 40) × 8 通道 = 360 維/幀，輸入形狀為 [360, T]（T 隨詞長不定）",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    # ── Slide 9: 資料增強（左文右四格圖） ────────────────────────────────
    make_split_slide(prs, 9, "資料增強策略（Data Augmentation）", [
        {"text": "目的：擴充訓練資料 4 倍，提升模型泛化能力", "bold": True, "bullet": False, "size": 17, "space_before": 4},
        {"text": "策略一：Speed Perturbation（速度擾動）", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "調快（Faster）：時間軸壓縮，模擬說話較快", "size": 16, "space_before": 3},
        {"text": "調慢（Slower）：時間軸延伸，模擬說話較慢", "size": 16, "space_before": 3},
        {"text": "策略二：Time-Frequency Masking（時頻遮蔽）", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "隨機遮蔽特徵圖的時間段或頻率段（仿 SpecAugment）", "size": 16, "space_before": 3},
        {"text": "提升模型對部分訊號缺失的魯棒性（Robustness）", "size": 16, "space_before": 3},
        {"text": "效果：資料增強使準確率提升 +1.7%（消融實驗驗證）",
         "bold": True, "bullet": False, "size": 16, "space_before": 8, "color": RED_COLOR},
    ],
    os.path.join(IMG_DIR, 'p06_img00.jpeg'),
    img_w=Cm(14.5))

    # ── Slide 10: 神經網路模型架構比較（全圖） ────────────────────────────
    make_fullimg_slide(prs, 10, "CNN Block 架構比較（六種設計）",
                       os.path.join(IMG_DIR, 'p08_img01.jpeg'),
                       caption="CNN / MultiScaled CNN / ResNet / Inception / Xception / ParallelNet（逐通道並行提取 + 全局融合）")

    # ── Slide 11: GRU 模型詳解（左文右架構圖） ────────────────────────────
    make_split_slide(prs, 11, "GRU 模型與 CNN+BiGRU 架構", [
        {"text": "為何 GRU 表現最佳？", "bold": True, "bullet": False, "size": 17, "space_before": 4},
        {"text": "sEMG 特徵圖本質是時序資料，時間依賴性比空間特徵更重要", "size": 16, "space_before": 3},
        {"text": "CNN 的局部性與平移不變性假設對抽象特徵圖不完全適用", "size": 16, "space_before": 3},
        {"text": "GRU（Gated Recurrent Unit）架構：", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "Update Gate（更新閘）：控制保留多少過去記憶", "size": 16, "space_before": 3},
        {"text": "Reset Gate（重置閘）：控制忽略多少過去資訊", "size": 16, "space_before": 3},
        {"text": "右圖為 CNN + BiGRU 架構：CNN 提取每幀空間特徵 → BiGRU 雙向捕捉時序依賴 → Max-Pool → MLP → 輸出",
         "size": 16, "space_before": 3},
        {"text": "模型輸入：360 維 × T 幀特徵圖（T 因詞長不同而異，GRU 天然支援可變長度）",
         "bold": True, "bullet": False, "size": 16, "space_before": 8, "color": RED_COLOR},
    ],
    os.path.join(IMG_DIR, 'p07_img00.jpeg'),
    img_w=Cm(14.5))

    # ── Slide 12: 實驗設定 ───────────────────────────────────────────────
    make_slide(prs, 12, "實驗設定與評估方法", [
        {"text": "評估方式：5-Fold Cross-Validation（五折交叉驗證）", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "兩項分類任務：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "① 多受試者混合分類（Multi-Subject Mixed）：12 位受試者資料混合訓練/驗證，135 類",
         "size": 17, "space_before": 3},
        {"text": "② 單一受試者分類（Single-Subject）：每位受試者獨立訓練/驗證，評估個人化性能",
         "size": 17, "space_before": 3},
        {"text": "比較模型（7 種）：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "GRU、CNN+GRU、CNN+LSTM、CNN+BiLSTM、ResNet+GRU、Inception+GRU、Xception+GRU、ParallelNet+GRU",
         "size": 17, "space_before": 3},
        {"text": "消融實驗（Ablation Study）：驗證 Fbank 特徵有效性與資料增強貢獻", "size": 17, "space_before": 3},
        {"text": "追加實驗：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "跨受試者：7 人訓練 → 1 人驗證（Leave-One-Out）", "size": 17, "space_before": 3},
        {"text": "跨期：同一受試者，Dataset 1（2700 筆）訓練 → Dataset 2（675 筆，相隔 4 天）驗證",
         "size": 17, "space_before": 3},
    ])

    # ── Slide 13: 結果 多受試者（全圖：左柱狀圖+右折線圖） ─────────────
    make_fullimg_slide(prs, 13, "實驗結果：多受試者混合分類",
                       os.path.join(IMG_DIR, 'p08_img00.jpeg'),
                       caption="(a) 各模型平均準確率：GRU 最佳 88.01%　　(b) 各受試者準確率折線：GRU（藍色）全程領先")

    # ── Slide 14: 結果 單一受試者 + 混淆矩陣（左圖） ────────────────────
    s14 = _blank(prs)
    add_title(s14, "實驗結果：單一受試者分類與混淆矩陣")
    # 左側文字
    add_content(s14, [
        {"text": "單一受試者最佳準確率：GRU 達 97.19%", "bold": True, "bullet": False,
         "size": 17, "color": RED_COLOR, "space_before": 4},
        {"text": "大部分受試者準確率介於 83–97%，個體差異明顯", "size": 16, "space_before": 4},
        {"text": "混淆矩陣（右圖）：", "bold": True, "bullet": False, "size": 17, "space_before": 8},
        {"text": "對角線深藍 → 正確率高，偶發混淆集中在語義相似的詞語", "size": 16, "space_before": 3},
        {"text": "單字詞（1 個漢字）普遍比多字詞準確率低", "size": 16, "space_before": 3},
        {"text": "中文同音字造成訊號相似，為 sEMG-based SSR 的固有挑戰", "size": 16, "space_before": 3},
    ],
    l=CONTENT_L, t=CONTENT_T, w=Cm(13.5), h=CONTENT_H)
    add_image(s14, os.path.join(IMG_DIR, 'p09_img00.jpeg'),
              CONTENT_L + Cm(14.0), CONTENT_T, Cm(15.5), CONTENT_H)
    add_pgnum(s14, 14)

    # ── Slide 15: 跨受試者與跨期問題 ─────────────────────────────────────
    make_slide(prs, 15, "跨受試者與跨期問題", [
        {"text": "跨受試者實驗（Cross-Subject）：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "訓練：7 人資料　→　驗證：1 人資料（Leave-One-Out）", "size": 17, "space_before": 3},
        {"text": "GRU 準確率從 88.01% 暴降至 47.22%（↓ 40.79%）",
         "bold": True, "bullet": False, "size": 17, "space_before": 3, "color": RED_COLOR},
        {"text": "跨期實驗（Cross-Session）：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "同一受試者，間隔 4 天的兩次採集（Dataset 1：2700 筆；Dataset 2：675 筆）", "size": 17, "space_before": 3},
        {"text": "同期：95.70%　→　跨期：78.67%（↓ 17.03%）",
         "bold": True, "bullet": False, "size": 17, "space_before": 3, "color": RED_COLOR},
        {"text": "根本原因：電極位置微小偏移、肌肉疲勞狀態、皮膚阻抗變化，使 sEMG 訊號跨期不一致",
         "size": 17, "space_before": 3},
        {"text": "固有限制：中文同音字 + 聲調資訊缺失 → 短字詞辨識天花板低", "size": 17, "space_before": 3},
        {"text": "未來方向：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "Transfer Learning（遷移學習）：解決跨受試者 / 跨期問題的核心方向", "size": 17, "space_before": 3},
        {"text": "連續辨識（Continuous Recognition）：目前為離散詞類，未來朝 token-level 輸出發展",
         "size": 17, "space_before": 3},
    ])

    # ── Slide 16: 結論與啟發 ──────────────────────────────────────────────
    make_slide(prs, 16, "結論與對本研究的啟發", [
        {"text": "主要貢獻：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "建立最大規模中文 SSR 語料庫（135 類），8 通道 sEMG，12 位受試者", "size": 17, "space_before": 3},
        {"text": "時頻混合特徵（Fbank + ZCR + 短時能量），有效捕捉 sEMG 短時資訊", "size": 17, "space_before": 3},
        {"text": "Speed Perturbation + Time-Frequency Masking 資料增強，準確率提升 +1.7%", "size": 17, "space_before": 3},
        {"text": "GRU 在所有模型中表現最佳：混合 88.01%、單一受試者最高 97.19%", "size": 17, "space_before": 3},
        {"text": "與陳睿莆研究的關聯：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "時頻混合特徵策略可延伸至 CSL-EMG 資料集，替代純 MFCC 作為輸入表示", "size": 17, "space_before": 3},
        {"text": "GRU 輕量化、原生支援可變長序列，適合本地端推論（Edge Computing）部署",
         "size": 17, "space_before": 3},
        {"text": "跨期問題與本研究一致 → Transfer Learning 為必要研究方向，可提升實用性",
         "size": 17, "space_before": 3},
        {"text": "資料增強策略可直接套用於 CSL-EMG 資料集，解決資料量不足問題",
         "size": 17, "space_before": 3},
        {"text": "核心啟示：sEMG 是時序訊號，RNN/GRU 架構比純 CNN 更適合作為 backbone",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    out = '/Users/rayopenclaw/ray-agent/pptx/BSPC2024_SSR_sEMG.pptx'
    prs.save(out)
    print(f"✓ 簡報已儲存：{out}（共 16 張投影片）")


if __name__ == '__main__':
    build()
