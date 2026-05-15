#!/usr/bin/env python3
"""
Sensors 2025 SSR Paper Presentation
Style: pure white background, #0E2841 titles, red technical terms, Times New Roman + 標楷體
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import re

TITLE_COLOR = RGBColor(0x0E, 0x28, 0x41)
RED_COLOR   = RGBColor(0xFF, 0x00, 0x00)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY        = RGBColor(0x66, 0x66, 0x66)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE  = RGBColor(0xF0, 0xF4, 0xF8)

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)
TITLE_L, TITLE_T, TITLE_W, TITLE_H             = Cm(2.3), Cm(1.0), Cm(29.2), Cm(3.7)
CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H     = Cm(2.3), Cm(5.1), Cm(29.2), Cm(12.1)
PGNUM_L, PGNUM_T, PGNUM_W, PGNUM_H             = Cm(23.9), Cm(17.7), Cm(3.0), Cm(1.0)

FONT_EN = "Times New Roman"
FONT_ZH = "標楷體"

RED_TERMS = sorted([
    "EMG", "EEG", "sEMG", "SNN", "CNN", "1D-CNN", "BLE",
    "SSR", "SSI", "Siamese", "Transformer", "GRU", "LSTM",
    "Few-shot", "few-shot", "Trigram", "trigram",
    "Moving Window", "moving window", "PMLSF",
    "Silent Speech", "silent speech",
    "Sensor Fusion", "sensor fusion",
    "Language Model", "language model",
], key=len, reverse=True)
RED_SET = set(RED_TERMS)


def _is_zh(c):
    return '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f'


def _set_fonts(run, font_name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
        el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', font_name)


def _add_run(para, text, font_name, size, bold=False, color=BLACK):
    r = para.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_fonts(r, font_name)
    return r


def _split_lang(text):
    if not text:
        return []
    segs, cur, cur_zh = [], text[0], _is_zh(text[0])
    for ch in text[1:]:
        zh = _is_zh(ch)
        if zh == cur_zh:
            cur += ch
        else:
            segs.append((cur, cur_zh))
            cur, cur_zh = ch, zh
    segs.append((cur, cur_zh))
    return segs


def add_text(para, text, size=18, bold=False, color=None, highlight=True):
    if color is not None:
        for seg, is_zh in _split_lang(text):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, color)
        return
    if not highlight:
        for seg, is_zh in _split_lang(text):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, BLACK)
        return
    pattern = '(' + '|'.join(re.escape(t) for t in RED_TERMS) + ')'
    for part in re.split(pattern, text):
        if not part:
            continue
        c = RED_COLOR if part in RED_SET else BLACK
        for seg, is_zh in _split_lang(part):
            _add_run(para, seg, FONT_ZH if is_zh else FONT_EN, size, bold, c)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_title(slide, text, size=28):
    tf = _textbox(slide, TITLE_L, TITLE_T, TITLE_W, TITLE_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_text(p, text, size, bold=True, color=TITLE_COLOR)


def add_pgnum(slide, n):
    tf = _textbox(slide, PGNUM_L, PGNUM_T, PGNUM_W, PGNUM_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, str(n), FONT_EN, 16, color=BLACK)


def add_content(slide, items, base=18):
    tf = _textbox(slide, CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H)
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if isinstance(item, str):
            p.space_before = Pt(5)
            add_text(p, item, base)
            continue

        text   = item.get('text', '')
        size   = item.get('size', base)
        bold   = item.get('bold', False)
        color  = item.get('color', None)
        bullet = item.get('bullet', True)
        sp     = item.get('space_before', 5)
        p.space_before = Pt(sp)

        prefix = ''
        if bullet and text and not text[:1] in ('•', '①', '②', '③', '④', '⑤', '—', '│', ' '):
            prefix = '• '

        if prefix:
            _add_run(p, prefix, FONT_EN, size, False, BLACK)
        if color is not None:
            add_text(p, text, size, bold, color=color)
        else:
            add_text(p, text, size, bold)


def make_slide(prs, num, title, items, title_size=28, base=18):
    s = _blank_slide(prs)
    add_title(s, title, title_size)
    add_content(s, items, base)
    add_pgnum(s, num)
    return s


# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ────────────────────────────────────────────────────
    s1 = _blank_slide(prs)
    tf1 = _textbox(s1, Cm(2.3), Cm(2.8), Cm(29.2), Cm(7.0))
    for i, (line, sz, bd) in enumerate([
        ("Sentence-Level Silent Speech Recognition", 30, True),
        ("Using a Wearable EMG/EEG Sensor System", 30, True),
        ("with AI-Driven Sensor Fusion and Language Model", 26, True),
    ]):
        p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4 if i else 0)
        _add_run(p, line, FONT_EN, sz, bd, TITLE_COLOR)

    tf2 = _textbox(s1, Cm(2.3), Cm(11.0), Cm(29.2), Cm(2.5))
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, "Satterlee, Zuo, Moon, Lee, Peterson, Kang  ·  SDSU  ·  MDPI Sensors  ·  2025",
             FONT_EN, 17, False, BLACK)
    p2 = tf2.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    add_text(p2, "報告人：陳睿莆　|　國立虎尾科技大學", 17, color=BLACK)

    tf3 = _textbox(s1, Cm(2.3), Cm(14.5), Cm(29.2), Cm(1.0))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    _add_run(p3, "DOI: 10.3390/s25196168", FONT_EN, 13, False, GRAY)
    add_pgnum(s1, 1)

    # ── Slide 2: 目錄 ─────────────────────────────────────────────────────
    make_slide(prs, 2, "目錄 / Table of Contents", [
        {"text": "研究背景與動機",                              "bullet": True, "size": 20, "space_before": 8},
        {"text": "系統架構概觀",                                "bullet": True, "size": 20, "space_before": 6},
        {"text": "硬體設計：EMG / EEG 感測器系統",               "bullet": True, "size": 20, "space_before": 6},
        {"text": "資料收集與語料庫",                            "bullet": True, "size": 20, "space_before": 6},
        {"text": "SNN 分類方法（1D-CNN Siamese Network）",       "bullet": True, "size": 20, "space_before": 6},
        {"text": "滑動視窗（Moving Window）分割策略",            "bullet": True, "size": 20, "space_before": 6},
        {"text": "感測器融合（Sensor Fusion）",                  "bullet": True, "size": 20, "space_before": 6},
        {"text": "Trigram 語言模型後處理",                       "bullet": True, "size": 20, "space_before": 6},
        {"text": "實驗結果與分析",                              "bullet": True, "size": 20, "space_before": 6},
        {"text": "討論、結論與研究關聯",                         "bullet": True, "size": 20, "space_before": 6},
    ], base=20)

    # ── Slide 3: 摘要 ─────────────────────────────────────────────────────
    make_slide(prs, 3, "摘要 / Abstract", [
        {"text": "研究目標：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "開發穿戴式句子級無聲語音辨識（Silent Speech Recognition, SSR）系統", "size": 17, "space_before": 3},
        {"text": "使用單通道 EMG + 單通道 EEG 訊號，透過藍牙 BLE 無線傳輸", "size": 17, "space_before": 3},
        {"text": "辨識四個軍事指令句（涵蓋 10 個關鍵詞）", "size": 17, "space_before": 3},
        {"text": "核心技術：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "Siamese Neural Network（SNN）以 1D-CNN 為骨幹，Few-shot 分類", "size": 17, "space_before": 3},
        {"text": "平行雙流感測器融合（Parallel Multi-Layer Sensor Fusion，PMLSF）", "size": 17, "space_before": 3},
        {"text": "Trigram Language Model 進行句子級語義校正", "size": 17, "space_before": 3},
        {"text": "最佳成果：句子正確率 95.25%（EMG+EEG 融合 + Language Model）",
         "bold": True, "bullet": False, "size": 18, "space_before": 10, "color": RED_COLOR},
    ])

    # ── Slide 4: 研究背景 ──────────────────────────────────────────────────
    make_slide(prs, 4, "研究背景與動機", [
        {"text": "什麼是 Silent Speech Interface（SSI）？", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "不依賴聲帶振動的語音溝通介面，捕捉發音動作產生的生理訊號", "size": 17, "space_before": 3},
        {"text": "應用場景：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "嘈雜戰場環境下的靜默通訊（本文主要場景）", "size": 17, "space_before": 3},
        {"text": "喉部手術後喪失發聲能力的患者輔助溝通", "size": 17, "space_before": 3},
        {"text": "潛水、太空等無法正常發聲的特殊環境", "size": 17, "space_before": 3},
        {"text": "現有挑戰：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "EMG 易受肌肉疲勞、電極位移影響，個體間差異大", "size": 17, "space_before": 3},
        {"text": "從孤立詞提升至連續句子辨識，詞彙邊界偵測困難", "size": 17, "space_before": 3},
        {"text": "EEG 訊號雜訊多，解析困難，但攜帶語言意圖資訊", "size": 17, "space_before": 3},
        {"text": "本研究貢獻：EMG+EEG 多模態融合，首次驗證句子級 SSR 可行性",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    # ── Slide 5: 系統架構 ──────────────────────────────────────────────────
    make_slide(prs, 5, "系統架構概觀", [
        {"text": "端到端五階段 Pipeline：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "① 穿戴式感測器採集 EMG + EEG 原始訊號（1000 Hz / 16 bit）", "size": 17, "space_before": 6},
        {"text": "② Moving Window 切割連續訊號（視窗 1000 ms，步進 100 ms）", "size": 17, "space_before": 4},
        {"text": "③ SNN（1D-CNN Siamese）對每個視窗進行 11 類 Few-shot 分類", "size": 17, "space_before": 4},
        {"text": "④ 平行雙流 Sensor Fusion（EMG stream + EEG stream）加權融合", "size": 17, "space_before": 4},
        {"text": "⑤ Trigram Language Model 重排候選詞序，輸出最終句子", "size": 17, "space_before": 4},
        {"text": "設計哲學：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "Few-shot 設計：每類只需 5 個支持樣本，快速個人化部署", "size": 17, "space_before": 3},
        {"text": "Siamese Network 直接比對特徵距離，無需大量重新訓練", "size": 17, "space_before": 3},
        {"text": "語言先驗知識補充感測器誤差，提升句子語義合理性", "size": 17, "space_before": 3},
    ])

    # ── Slide 6: 硬體設計 ──────────────────────────────────────────────────
    make_slide(prs, 6, "硬體設計：EMG / EEG 感測器系統", [
        {"text": "EMG 感測器：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "單通道 sEMG，電極貼附於下頜肌群（臉頰/下巴區域）", "size": 17, "space_before": 3},
        {"text": "記錄無聲默念時的下頜動作與唇部肌肉收縮訊號", "size": 17, "space_before": 3},
        {"text": "EEG 感測器：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "單通道腦電圖，電極貼附於前額（Fp1/Fp2 附近）", "size": 17, "space_before": 3},
        {"text": "捕捉語言意圖相關的神經皮質活動（Broca's area 活動）", "size": 17, "space_before": 3},
        {"text": "無線傳輸設計：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "藍牙低功耗（BLE）無線模組，取代 USB 有線連接", "size": 17, "space_before": 3},
        {"text": "關鍵優勢：避免電源線引入 60 Hz 工頻雜訊，提升訊號純淨度", "size": 17, "space_before": 3},
        {"text": "電池供電，提升穿戴舒適性與訊號穩定性", "size": 17, "space_before": 3},
        {"text": "規格：取樣率 1000 Hz　|　解析度 16 bit　|　單通道 EMG + 單通道 EEG",
         "bold": True, "bullet": False, "size": 16, "space_before": 8},
    ])

    # ── Slide 7: 資料收集 ──────────────────────────────────────────────────
    make_slide(prs, 7, "資料收集與語料庫設計", [
        {"text": "語料庫 A — 孤立詞辨識（Isolated Word）：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "10 個目標詞：Need, Medical, Assistance, Land, Here, Do, Not, Pick, Us, Up",
         "size": 17, "space_before": 3},
        {"text": "每詞 150 次 × 10 詞 = 1,500 個樣本（單一受試者）", "size": 17, "space_before": 3},
        {"text": "語料庫 B — 句子級辨識（Sentence-Level）：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "四個軍事指令句：", "size": 17, "space_before": 3, "bullet": False},
        {"text": " ①  Need Medical Assistance       ②  Land Here", "size": 16, "bullet": False, "space_before": 2},
        {"text": " ③  Do Not Land Here               ④  Pick Us Up", "size": 16, "bullet": False, "space_before": 2},
        {"text": "每句 100 次 × 4 句 = 400 個句子樣本", "size": 17, "space_before": 4},
        {"text": "收集協議：受試者靜坐，以 Silent Articulation（無聲默念）方式錄製", "size": 17, "space_before": 6},
        {"text": "每詞平均時長約 500–1500 ms，句子間有靜默間隔（Null class）", "size": 17, "space_before": 3},
    ])

    # ── Slide 8: SNN 分類方法 ──────────────────────────────────────────────
    make_slide(prs, 8, "SNN 分類方法：1D-CNN Siamese Network", [
        {"text": "Siamese Network 核心概念：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "兩條「權重共享」的 CNN 分支，同時處理「查詢樣本」與「支持集樣本」",
         "size": 17, "space_before": 3},
        {"text": "計算兩分支輸出特徵向量的 L1 距離，判斷是否屬於相同類別", "size": 17, "space_before": 3},
        {"text": "Few-shot 設定：11-way 5-shot（10 個詞 + Null 靜默類，每類 5 個參考樣本）",
         "size": 17, "space_before": 3},
        {"text": "1D-CNN 骨幹架構：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "輸入：1000 點時序訊號（1000 ms 視窗 × 1000 Hz 取樣率）", "size": 17, "space_before": 3},
        {"text": "結構：Conv1D → BatchNorm → MaxPool → Dropout → FC → 特徵向量",
         "size": 17, "space_before": 3},
        {"text": "訓練目標：Binary Cross-Entropy（相同類別輸出 1，不同類別輸出 0）",
         "size": 17, "space_before": 3},
        {"text": "優勢：新使用者只需提供 5 個樣本即可完成分類，無需重新訓練整個模型",
         "bold": True, "bullet": False, "size": 17, "space_before": 8},
    ])

    # ── Slide 9: 滑動視窗 ──────────────────────────────────────────────────
    make_slide(prs, 9, "滑動視窗（Moving Window）分割策略", [
        {"text": "核心問題：連續訊號中如何偵測詞彙邊界？", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "Moving Window 參數：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "視窗大小（Window Size）：1000 ms", "size": 17, "space_before": 3},
        {"text": "步進量（Step Size）：100 ms（90% 重疊）", "size": 17, "space_before": 3},
        {"text": "每個視窗輸入 SNN 進行 11 類分類（10 詞 + Null 靜默類）", "size": 17, "space_before": 3},
        {"text": "後處理邏輯：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "連續相同預測 → 合併為一個詞彙事件（Majority Vote）", "size": 17, "space_before": 3},
        {"text": "Null 視窗 → 詞彙間的靜默間隔，用於切分詞彙序列", "size": 17, "space_before": 3},
        {"text": "候選詞序列 → 輸出至 Trigram Language Model 進行句子重排", "size": 17, "space_before": 3},
        {"text": "設計優勢：全自動端到端處理，無需人工標記詞彙邊界時間戳記",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    # ── Slide 10: 感測器融合 ───────────────────────────────────────────────
    make_slide(prs, 10, "感測器融合（Sensor Fusion）", [
        {"text": "融合策略：平行多層感測器融合（PMLSF）", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "架構說明：", "bold": True, "bullet": False, "size": 18, "space_before": 6},
        {"text": "EMG stream：獨立 1D-CNN SNN 提取肌肉動作特徵向量", "size": 17, "space_before": 3},
        {"text": "EEG stream：獨立 1D-CNN SNN 提取神經語言意圖特徵向量", "size": 17, "space_before": 3},
        {"text": "兩串流的分類置信分數（Confidence Scores）加權線性融合", "size": 17, "space_before": 3},
        {"text": "融合公式：Score_fusion = α × Score_EMG + (1 − α) × Score_EEG",
         "size": 17, "space_before": 3, "bullet": False},
        {"text": "融合效果（數據支持）：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "孤立詞準確率：EMG 90.33%　→　融合 98.67%　（↑ 8.34%）", "size": 17, "space_before": 3},
        {"text": "句子正確率：　 EMG 73.41%　→　融合 93.30%　（↑ 19.89%）", "size": 17, "space_before": 3},
        {"text": "結論：EMG（動作精確度）與 EEG（語言意圖）形成互補，融合效益顯著",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    # ── Slide 11: 語言模型後處理 ──────────────────────────────────────────
    make_slide(prs, 11, "Trigram 語言模型後處理", [
        {"text": "動機：單詞辨識誤差在句子層面累積，需語義約束校正",
         "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "Trigram Language Model 原理：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "基於三連詞統計，估計條件概率 P(w₃ | w₁, w₂)", "size": 17, "space_before": 3},
        {"text": "訓練語料：4 個目標句子的詞序排列組合（小型特化語料庫）", "size": 17, "space_before": 3},
        {"text": "對 SNN 輸出的候選詞序列重新計算句子概率分數", "size": 17, "space_before": 3},
        {"text": "後處理流程：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "① SNN 輸出 Top-K 候選詞序列（每個視窗保留 Top-3 候選）", "size": 17, "space_before": 3},
        {"text": "② Trigram LM 評分，懲罰語義不合理的詞序組合", "size": 17, "space_before": 3},
        {"text": "③ 選擇概率最高的完整句子作為最終輸出", "size": 17, "space_before": 3},
        {"text": "效果：Config B 融合 93.30% → Config C 加入 LM 後 95.25%（↑ 1.95%）",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": RED_COLOR},
    ])

    # ── Slide 12: 實驗結果摘要 ────────────────────────────────────────────
    make_slide(prs, 12, "實驗結果摘要", [
        {"text": "三種實驗配置：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "Config A — 孤立詞辨識（孤立詞語料庫 A，1500 樣本）", "size": 17, "space_before": 3},
        {"text": "Config B — 句子內詞彙辨識（句子語料庫 B，含 Moving Window）", "size": 17, "space_before": 3},
        {"text": "Config C — 完整系統（Config B + Trigram LM 後處理）", "size": 17, "space_before": 3},
        {"text": "關鍵數字：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "Config A — 孤立詞：EMG 90.33%　|　EEG 82.67%　|　融合 98.67%",
         "size": 17, "space_before": 3},
        {"text": "Config B — 句子：  EMG 73.41%　|　EEG 52.16%　|　融合 93.30%",
         "size": 17, "space_before": 3},
        {"text": "Config C — 句子：  EMG 93.33%　|　EEG 87.83%　|　融合+LM 95.25%",
         "bold": True, "size": 17, "space_before": 3, "color": RED_COLOR},
        {"text": "融合相比單 EMG 提升最多 +19.89%，語言模型再貢獻 +1.95%",
         "bold": True, "bullet": False, "size": 17, "space_before": 8},
    ])

    # ── Slide 13: 完整數據表 ──────────────────────────────────────────────
    s13 = _blank_slide(prs)
    add_title(s13, "完整實驗數據表")

    rows, cols = 10, 5
    tbl = s13.shapes.add_table(rows, cols, CONTENT_L, CONTENT_T, CONTENT_W, Cm(11.5)).table

    headers = ["配置", "訊號模態", "孤立詞準確率", "句中詞準確率", "句子準確率"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(15)
                run.font.name = FONT_ZH
                run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR

    data = [
        ("A", "EMG",             "90.33%", "29.15%", "—"),
        ("A", "EEG",             "82.67%", "17.10%", "—"),
        ("A", "EMG + EEG 融合",  "98.67%", "26.36%", "—"),
        ("B", "EMG",             "—",      "85.07%", "73.41%"),
        ("B", "EEG",             "—",      "80.06%", "52.16%"),
        ("B", "EMG + EEG 融合",  "—",      "91.24%", "93.30%"),
        ("C", "EMG",             "—",      "—",      "93.33%"),
        ("C", "EEG",             "—",      "—",      "87.83%"),
        ("C", "EMG + EEG + LM",  "—",      "—",      "95.25% ★"),
    ]
    for ri, row in enumerate(data):
        best = (row[4] == "95.25% ★")
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(14)
                    run.font.name = FONT_ZH if any(_is_zh(c) for c in val) else FONT_EN
                    run.font.bold = best
                    run.font.color.rgb = RED_COLOR if best else BLACK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
    add_pgnum(s13, 13)

    # ── Slide 14: 討論 ────────────────────────────────────────────────────
    make_slide(prs, 14, "討論", [
        {"text": "EMG vs EEG 的互補性：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "EMG 動作精確度優於 EEG（90.33% vs 82.67%），但兩者攜帶不重疊資訊",
         "size": 17, "space_before": 3},
        {"text": "融合後孤立詞達 98.67%，說明多模態整合有效", "size": 17, "space_before": 3},
        {"text": "句子辨識的瓶頸：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "孤立詞準確高（98.67%）→ 句子辨識下降原因：Moving Window 詞彙邊界偵測誤差累積",
         "size": 17, "space_before": 3},
        {"text": "步進 100 ms 可能遺漏短詞或誤判靜默段，影響詞序完整性", "size": 17, "space_before": 3},
        {"text": "語言模型的效益與限制：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "Trigram LM 僅有 4 句訓練語料，已帶來 +1.95% 提升（93.30% → 95.25%）",
         "size": 17, "space_before": 3},
        {"text": "更大詞彙量場景下，語言模型貢獻預期更為顯著", "size": 17, "space_before": 3},
        {"text": "主要限制：單一受試者（泛化能力待驗證）；小詞彙量（10 詞，4 句）",
         "bold": True, "bullet": False, "size": 17, "space_before": 8, "color": GRAY},
    ])

    # ── Slide 15: 結論與研究關聯 ──────────────────────────────────────────
    make_slide(prs, 15, "結論與對本研究的啟發", [
        {"text": "主要貢獻：", "bold": True, "bullet": False, "size": 18, "space_before": 4},
        {"text": "首次驗證 EMG+EEG 融合用於句子級 Silent Speech Recognition 可行性",
         "size": 17, "space_before": 3},
        {"text": "SNN（1D-CNN）Few-shot 架構有效降低資料收集門檻", "size": 17, "space_before": 3},
        {"text": "三階段系統（Moving Window + Sensor Fusion + Language Model）達 95.25%",
         "size": 17, "space_before": 3},
        {"text": "技術意義：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "穿戴式、BLE 無線、低複雜度 → 實際部署可行性高", "size": 17, "space_before": 3},
        {"text": "Few-shot 設計 → 適合快速個人化，無需大量重新訓練", "size": 17, "space_before": 3},
        {"text": "與陳睿莆研究的關聯：", "bold": True, "bullet": False, "size": 18, "space_before": 8},
        {"text": "EMG+EEG 融合策略可作為 CSL-EMG 資料集研究的延伸方向", "size": 17, "space_before": 3},
        {"text": "SNN 少樣本架構可與 Transformer Encoder 結合，提升跨用戶泛化能力",
         "size": 17, "space_before": 3},
        {"text": "BLE 無線設計與電池供電理念與義肢控制系統硬體設計一致", "size": 17, "space_before": 3},
    ])

    out = '/Users/rayopenclaw/ray-agent/pptx/Sensors2025_SSR.pptx'
    prs.save(out)
    print(f"✓ 簡報已儲存：{out}（共 15 張投影片）")


if __name__ == '__main__':
    build()
