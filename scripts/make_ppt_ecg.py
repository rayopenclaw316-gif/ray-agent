#!/usr/bin/env python3
"""
ECG 心律不整分類 PPT 產生器
使用 python-pptx，遵循 ray-agent 風格規範
"""
import os, re
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pptx.util import Emu

BASE       = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR    = os.path.join(BASE, 'output', 'ecg')
OUT_PATH   = os.path.join(BASE, 'ECG_CNN_Report.pptx')

# ── 顏色常數 ────────────────────────────────────────────
DARK_BLUE = RGBColor(0x0E, 0x28, 0x41)
RED       = RGBColor(0xFF, 0x00, 0x00)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# 技術術語紅色高亮清單
RED_TERMS = [
    '1D CNN', 'CNN', 'Conv1D', 'MaxPool1D', 'ReLU', 'Dropout', 'CrossEntropyLoss',
    'Adam', 'Softmax', 'Linear', 'Flatten', 'Batch Size', 'Epoch',
    'MIT-BIH', 'ECG', 'Holter', 'Arrhythmia',
    'Normal', 'Supraventricular', 'Ventricular', 'Fusion', 'Unknown',
    'Kaggle', 'PyTorch', 'MPS',
    '反向傳播', '梯度下降', '卷積', '池化', '全連接層', '過擬合',
]

# ── 字型工具 ────────────────────────────────────────────
def _set_fonts(run, font_name):
    rPr = run._r.get_or_add_rPr()
    for tag in ('latin', 'ea', 'cs'):
        el = rPr.find(qn(f'a:{tag}'))
        if el is None:
            el = etree.SubElement(rPr, qn(f'a:{tag}'))
        el.set('typeface', font_name)

def _is_chinese(ch):
    return '\u4e00' <= ch <= '\u9fff' or '\u3000' <= ch <= '\u303f'

def _split_lang(text):
    """拆分成 (segment, is_chinese) 清單"""
    if not text:
        return []
    segments, buf, cur_cn = [], '', None
    for ch in text:
        cn = _is_chinese(ch)
        if cur_cn is None:
            cur_cn = cn
        if cn != cur_cn:
            segments.append((buf, cur_cn))
            buf, cur_cn = ch, cn
        else:
            buf += ch
    if buf:
        segments.append((buf, cur_cn))
    return segments

def _match_red(text):
    """回傳 text 中與 RED_TERMS 相符的片段（貪婪，最長優先）"""
    sorted_terms = sorted(RED_TERMS, key=len, reverse=True)
    pattern = '|'.join(re.escape(t) for t in sorted_terms)
    return list(re.finditer(pattern, text))

def add_text(para, text, size=18, bold=False, color=BLACK, align=None):
    """將含中英文混合 + 紅色術語的文字加入段落"""
    if align:
        para.alignment = align

    # 將 text 按紅色術語切分
    sorted_terms = sorted(RED_TERMS, key=len, reverse=True)
    pattern = '(' + '|'.join(re.escape(t) for t in sorted_terms) + ')'
    parts = re.split(pattern, text)

    for part in parts:
        if not part:
            continue
        is_red = part in RED_TERMS
        segs = _split_lang(part)
        for seg_text, is_cn in segs:
            run = para.add_run()
            run.text = seg_text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RED if is_red else color
            _set_fonts(run, '標楷體' if is_cn else 'Times New Roman')


# ── 版面工具 ────────────────────────────────────────────
def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 純白背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def add_title(slide, text, top_cm=1.0):
    txBox = slide.shapes.add_textbox(Cm(2.3), Cm(top_cm), Cm(29.2), Cm(3.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text(p, text, size=29, bold=True, color=DARK_BLUE)

def add_content_box(slide, top_cm=5.1, height_cm=12.1):
    txBox = slide.shapes.add_textbox(Cm(2.3), Cm(top_cm), Cm(29.2), Cm(height_cm))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_bullet(tf, text, level=0, size=17, bold=False):
    """加一行 bullet 文字（自動處理混合語言 + 紅色術語）"""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    indent = '    ' * level
    add_text(p, indent + text, size=size, bold=bold)

def add_page_num(slide, num, total):
    txBox = slide.shapes.add_textbox(Cm(23.9), Cm(17.7), Cm(3.0), Cm(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_text(p, f'{num} / {total}', size=13, color=RGBColor(0x88, 0x88, 0x88))

def add_image(slide, path, left_cm, top_cm, width_cm=None, height_cm=None):
    kwargs = {}
    if width_cm:  kwargs['width']  = Cm(width_cm)
    if height_cm: kwargs['height'] = Cm(height_cm)
    slide.shapes.add_picture(path, Cm(left_cm), Cm(top_cm), **kwargs)


# ── 建立 PPT ────────────────────────────────────────────
def build_ppt():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    TOTAL = 17

    # ══════════════════════════════════════════════════
    # Slide 1 — 封面
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    tb = s.shapes.add_textbox(Cm(2.3), Cm(3.5), Cm(29.2), Cm(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_text(p, 'ECG Heartbeat Classification Using 1D CNN', size=28, bold=True, color=DARK_BLUE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    add_text(p2, '利用一維卷積神經網路進行心電圖心律不整分類', size=22, bold=True, color=DARK_BLUE)

    tb2 = s.shapes.add_textbox(Cm(2.3), Cm(10.0), Cm(29.2), Cm(4.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for line in [
        '深度學習課程期末報告',
        '資料集：MIT-BIH Arrhythmia Database（Kaggle）',
        '模型：1D CNN（PyTorch）',
        '報告者：陳睿莆　國立虎尾科技大學　機械與電腦輔助工程系',
        '日期：2026-04-19',
    ]:
        p = tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        add_text(p, line, size=16)
    add_page_num(s, 1, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 2 — 目錄
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '目錄 Table of Contents')
    tf = add_content_box(s)
    toc = [
        ('1.  心律不整簡介', 3),
        ('2.  資料集如何取得與收集', 4),
        ('3.  資料集概覽', 6),
        ('4.  使用模型：1D CNN', 7),
        ('5.  模型架構詳解', 8),
        ('6.  逐層訓練說明', 9),
        ('7.  程式碼解說', 12),
        ('8.  訓練結果', 15),
        ('9.  結論', 17),
    ]
    for item, page in toc:
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        add_text(p, f'{item}{"·" * (52 - len(item))} {page}', size=17)
    add_page_num(s, 2, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 3 — 心律不整簡介
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '心律不整簡介 Arrhythmia Introduction')
    tf = add_content_box(s)
    bullets = [
        ('什麼是心律不整？', 0, True),
        ('心臟正常跳動依賴心臟傳導系統（sinoatrial node → atrioventricular node → His bundle → Purkinje fibers）依序發出電訊號。', 1, False),
        ('當電訊號傳導異常，導致心跳過快、過慢或不規則，即稱為 Arrhythmia（心律不整）。', 1, False),
        ('', 0, False),
        ('臨床上的五種常見類別（本研究分類目標）：', 0, True),
        ('Normal（正常）：規律竇性節律，P-QRS-T 波形完整。', 1, False),
        ('Supraventricular（上心室早期收縮）：起源於心室以上，QRS 波形較窄。', 1, False),
        ('Ventricular（心室早期收縮）：起源於心室，QRS 波形寬且異常，危險性較高。', 1, False),
        ('Fusion（融合搏動）：正常與異常搏動同時發生，波形介於兩者之間。', 1, False),
        ('Unknown（未分類）：無法歸入上述四類的心搏，通常由雜訊或罕見波形造成。', 1, False),
        ('', 0, False),
        ('臨床意義：心律不整若未即時診斷，可能導致心臟衰竭、中風，甚至猝死。自動化偵測對提升醫療效率至關重要。', 0, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 3, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 4 — 資料收集方式
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '資料如何收集？How Was the Data Collected?')
    tf = add_content_box(s, top_cm=5.1, height_cm=7.0)
    bullets = [
        ('MIT-BIH Arrhythmia Database（原始資料來源）', 0, True),
        ('由美國麻省理工學院（MIT）與波士頓貝斯以色列醫院（BIH）聯合建立，1975–1979 年收集。', 1, False),
        ('受試者：47 位患者，年齡 23–89 歲，男女各半，門診與住院患者皆有。', 1, False),
        ('收集方式：受試者佩戴 Holter Monitor（攜帶式心電圖記錄儀），連續記錄 24–48 小時的 ECG 訊號。', 1, False),
        ('採樣率：360 Hz，每次記錄使用 2 個導聯（lead）。', 1, False),
        ('', 0, False),
        ('Holter Monitor 工作原理', 0, True),
        ('電極貼片黏貼於胸部，偵測心肌電位變化，轉為 ECG 波形（P, QRS, T 波）。', 1, False),
        ('每個心搏（heartbeat）切割為 187 個時間點的等長片段，以 R 波（心室最大收縮波）為中心對齊。', 1, False),
        ('', 0, False),
        ('資料標注：兩位心臟科醫師獨立標注每個心搏的類別，有分歧時由第三位仲裁。', 0, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)

    # Kaggle 取得方式
    tf2 = s.shapes.add_textbox(Cm(2.3), Cm(13.2), Cm(29.2), Cm(4.5))
    tf2.text_frame.word_wrap = True
    bullets2 = [
        ('本研究資料取得方式', 0, True),
        ('原始資料由 Shayan Fazeli 預處理後上傳至 Kaggle（shayanfazeli/heartbeat）。', 1, False),
        ('預處理步驟：切割心搏 → 補零至 187 點 → 正規化至 [0, 1] → 儲存為 CSV。', 1, False),
        ('下載指令：kaggle datasets download -d shayanfazeli/heartbeat', 1, False),
    ]
    for text, level, bold in bullets2:
        add_bullet(tf2.text_frame, text, level=level, bold=bold)
    add_page_num(s, 4, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 5 — ECG 訊號樣本
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, 'ECG 訊號樣本 Signal Examples')
    tf = add_content_box(s, top_cm=5.0, height_cm=2.5)
    add_bullet(tf, '每筆資料為一個心搏片段，包含 187 個採樣點（約 0.52 秒 @ 360 Hz），以 R 波為中心對齊。', size=16)
    add_bullet(tf, '不同類別的 ECG 波形在 QRS 寬度、振幅、形狀上有明顯差異，是 1D CNN 學習的依據。', size=16)
    img_path = os.path.join(IMG_DIR, 'ecg_samples.png')
    if os.path.exists(img_path):
        add_image(s, img_path, left_cm=2.0, top_cm=7.5, width_cm=29.8)
    add_page_num(s, 5, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 6 — 資料集概覽
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '資料集概覽 Dataset Overview')
    tf = add_content_box(s, top_cm=5.0, height_cm=5.5)
    add_bullet(tf, '訓練集：87,554 筆｜測試集：21,892 筆｜特徵維度：187（時間點）', bold=True)
    add_bullet(tf, '類別分布（訓練集）：', bold=True)
    dist = [
        ('Class 0 — Normal：72,471 筆（82.8%）', 1),
        ('Class 2 — Ventricular：5,788 筆（6.6%）', 1),
        ('Class 4 — Unknown：6,431 筆（7.3%）', 1),
        ('Class 1 — Supraventricular：2,223 筆（2.5%）', 1),
        ('Class 3 — Fusion：641 筆（0.7%）', 1),
    ]
    for text, level in dist:
        add_bullet(tf, text, level=level)
    add_bullet(tf, '注意：資料集嚴重不平衡（Normal 佔 82.8%），訓練時模型可能對少數類別表現較差。', bold=False)
    img_path = os.path.join(IMG_DIR, 'class_distribution.png')
    if os.path.exists(img_path):
        add_image(s, img_path, left_cm=5.0, top_cm=11.0, width_cm=24.0)
    add_page_num(s, 6, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 7 — 為什麼選 1D CNN
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '使用模型：1D CNN 為何選擇此架構？')
    tf = add_content_box(s)
    bullets = [
        ('ECG 訊號是一維時間序列', 0, True),
        ('每筆資料為 187 個連續時間點的振幅值，本質是「時間 → 振幅」的一維序列。', 1, False),
        ('傳統 2D CNN 適用於影像（二維空間），對一維訊號效率較低。', 1, False),
        ('1D CNN 的核（kernel）沿時間軸滑動，直接捕捉局部波形特徵（QRS 寬度、振幅變化）。', 1, False),
        ('', 0, False),
        ('1D CNN 的優勢', 0, True),
        ('局部感知（Local Receptive Field）：每個 kernel 只看連續幾個時間點，符合 ECG 波形的局部特性。', 1, False),
        ('參數共享（Parameter Sharing）：同一個 kernel 掃描整條訊號，大幅減少參數量。', 1, False),
        ('平移不變性：即使 QRS 波形在時間軸上稍有偏移，CNN 仍能辨識。', 1, False),
        ('計算效率高：比 RNN/LSTM 收斂更快，適合即時分析。', 1, False),
        ('', 0, False),
        ('與其他方法比較', 0, True),
        ('SVM / 傳統 ML：需手動設計特徵（頻域、時域），費時費力。', 1, False),
        ('LSTM / Transformer：可學習長期依賴，但訓練成本高，對本任務未必必要。', 1, False),
        ('1D CNN：自動學習特徵 + 計算快 + 準確率高，是 ECG 分類的主流方法。', 1, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 7, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 8 — 模型架構概覽
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '模型架構 Model Architecture')
    tf = add_content_box(s, top_cm=5.0, height_cm=3.5)
    add_bullet(tf, '輸入：(Batch, 1, 187)　→　3 個卷積塊（Conv1D + ReLU + MaxPool1D）　→　全連接分類器　→　5 類輸出', bold=True)
    add_bullet(tf, '共 3 個卷積塊：逐步提取從低階（局部波形）到高階（抽象心律特徵）的特徵。', size=16)

    # 用矩形繪製架構流程圖

    boxes = [
        ('Input\n(B,1,187)',   2.3,  9.5, 3.8, 2.5),
        ('Conv1D(1→32,k=5)\nReLU\nMaxPool1D\n(B,32,91)',   6.8,  9.5, 4.5, 2.5),
        ('Conv1D(32→64,k=5)\nReLU\nMaxPool1D\n(B,64,43)',  12.0, 9.5, 4.5, 2.5),
        ('Conv1D(64→128,k=3)\nReLU\nMaxPool1D\n(B,128,20)', 17.2, 9.5, 4.5, 2.5),
        ('Flatten\n(B,2560)',  22.4, 9.5, 3.0, 2.5),
        ('FC 2560→256\nReLU\nDropout\nFC 256→64→5', 26.1, 9.5, 5.4, 2.5),
    ]
    box_colors = [
        RGBColor(0xE3,0xF2,0xFD),
        RGBColor(0xBB,0xDE,0xFB),
        RGBColor(0x90,0xCA,0xF9),
        RGBColor(0x64,0xB5,0xF6),
        RGBColor(0xE8,0xF5,0xE9),
        RGBColor(0xC8,0xE6,0xC9),
    ]
    for i, (label, l, t, w, h) in enumerate(boxes):
        shape = s.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = box_colors[i]
        shape.line.color.rgb = DARK_BLUE
        tf2 = shape.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = DARK_BLUE
        _set_fonts(run, 'Times New Roman')

    # 箭頭（用文字方塊模擬）
    arrow_positions = [(6.0, 10.5), (11.1, 10.5), (16.3, 10.5), (21.5, 10.5), (25.3, 10.5)]
    for lc, tc in arrow_positions:
        tb = s.shapes.add_textbox(Cm(lc), Cm(tc), Cm(0.7), Cm(0.5))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = '→'
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_BLUE

    add_page_num(s, 8, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 9 — 逐層解說：卷積層
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '逐層解說（一）Conv1D 卷積層')
    tf = add_content_box(s)
    bullets = [
        ('Conv1D 的運作原理', 0, True),
        ('一個「卷積核（kernel）」是一組可學習的權重，長度為 kernel_size（如 5）。', 1, False),
        ('核沿時間軸從左到右滑動，每次與對應的 5 個時間點做內積（加權求和），輸出一個數值。', 1, False),
        ('重複滑動後得到一條新的時間序列（Feature Map），代表「此核所偵測到的局部特徵」的強度。', 1, False),
        ('使用多個核（out_channels）可同時偵測多種不同特徵（如 QRS 上升沿、下降沿、頂點）。', 1, False),
        ('', 0, False),
        ('本模型三個卷積塊的維度變化：', 0, True),
        ('卷積塊 1：輸入 (B, 1, 187)  →  Conv1D(1, 32, k=5) → (B, 32, 183)  →  MaxPool1D(2) → (B, 32, 91)', 1, False),
        ('卷積塊 2：輸入 (B, 32, 91)  →  Conv1D(32, 64, k=5) → (B, 64, 87)  →  MaxPool1D(2) → (B, 64, 43)', 1, False),
        ('卷積塊 3：輸入 (B, 64, 43)  →  Conv1D(64, 128, k=3) → (B, 128, 41) →  MaxPool1D(2) → (B, 128, 20)', 1, False),
        ('', 0, False),
        ('輸出長度公式：L_out = (L_in − kernel_size) / stride + 1　（padding=0, stride=1）', 0, False),
        ('MaxPool1D 作用：取每 2 個相鄰點的最大值，長度減半，保留最顯著特徵、降低計算量。', 0, False),
        ('ReLU 作用：f(x) = max(0, x)，抑制負值（無意義的特徵），引入非線性。', 0, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 9, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 10 — 逐層解說：Flatten + 全連接層
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '逐層解說（二）Flatten + 全連接層')
    tf = add_content_box(s)
    bullets = [
        ('Flatten — 攤平操作', 0, True),
        ('三個卷積塊後，輸出形狀為 (B, 128, 20)。', 1, False),
        ('Flatten 將 128 個 channel × 20 個時間點 攤平為一條 2560 維的向量：(B, 128, 20) → (B, 2560)。', 1, False),
        ('這個向量是整條 ECG 訊號的高階特徵表示。', 1, False),
        ('', 0, False),
        ('Linear（全連接層）', 0, True),
        ('FC1：Linear(2560 → 256)　每個輸入神經元與每個輸出神經元都有一條可學習的連接（權重）。', 1, False),
        ('FC1 後接 ReLU + Dropout(0.5)：訓練時隨機丟棄 50% 的神經元，防止過擬合。', 1, False),
        ('FC2：Linear(256 → 64)　進一步壓縮特徵維度。', 1, False),
        ('FC3：Linear(64 → 5)　輸出 5 個數值（logits），分別對應 5 個心律類別的未歸一化分數。', 1, False),
        ('', 0, False),
        ('損失函數：CrossEntropyLoss', 0, True),
        ('= Softmax + Negative Log-Likelihood。先將 5 個 logits 轉為機率分布，再計算預測機率與真實標籤的差距。', 1, False),
        ('公式：L = -log(p_true)　其中 p_true 是模型對正確類別預測的機率，越高損失越低。', 1, False),
        ('', 0, False),
        ('優化器：Adam（Adaptive Moment Estimation）', 0, True),
        ('為每個參數自動調整學習率，結合動量（Momentum）與 RMSprop 的優點，收斂速度快。', 1, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 10, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 11 — 訓練流程
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '逐層解說（三）訓練流程與反向傳播')
    tf = add_content_box(s)
    bullets = [
        ('一次訓練迭代（iteration）的完整流程：', 0, True),
        ('① 前向傳播（Forward Pass）：輸入 (B,1,187) 依序流過所有層，得到輸出 logits (B,5)。', 1, False),
        ('② 計算損失（Loss）：CrossEntropyLoss 比較 logits 與真實標籤，產生一個純量損失值。', 1, False),
        ('③ 反向傳播（Backward Pass）：對損失值求偏微分，用 Chain Rule 逐層計算每個參數的梯度。', 1, False),
        ('④ 參數更新（Optimizer Step）：Adam 根據梯度與歷史梯度調整每個參數的值（往損失下降方向移動）。', 1, False),
        ('⑤ 清除梯度（zero_grad）：PyTorch 預設累積梯度，每次迭代前必須清除。', 1, False),
        ('', 0, False),
        ('Epoch vs. Iteration', 0, True),
        ('1 Epoch = 完整訓練資料跑過一遍（87554 筆 ÷ Batch Size 256 ≈ 342 iterations）。', 1, False),
        ('本次訓練共 20 Epochs，總計約 6840 次參數更新。', 1, False),
        ('', 0, False),
        ('Dropout 的訓練/推論差異', 0, True),
        ('訓練時（model.train()）：隨機丟棄 50% 神經元，強迫模型不依賴單一特徵，增加泛化能力。', 1, False),
        ('推論時（model.eval()）：全部神經元啟用，但輸出乘以 (1-p) 補償訓練時的期望值。', 1, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 11, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 12 — 程式碼：資料載入
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '程式碼解說（一）資料載入 ECGDataset')
    tf = add_content_box(s)
    code_lines = [
        ('class ECGDataset(Dataset):', False),
        ('    def __init__(self, csv_path):', False),
        ('        df = pd.read_csv(csv_path, header=None)', False),
        ('        # df.shape = (N, 188)：前 187 欄為訊號值，最後 1 欄為標籤', True),
        ('        self.X = torch.tensor(df.iloc[:, :187].values, dtype=torch.float32)', False),
        ('        # X.shape = (N, 187)  ← N 筆樣本，每筆 187 個時間點', True),
        ('        self.y = torch.tensor(df.iloc[:, -1].values, dtype=torch.long)', False),
        ('        # y.shape = (N,)  ← 每筆對應一個類別標籤 0~4', True),
        ('', False),
        ('    def __getitem__(self, idx):', False),
        ('        return self.X[idx].unsqueeze(0), self.y[idx]', False),
        ('        # unsqueeze(0) 把 (187,) 變成 (1, 187)', True),
        ('        # Conv1d 需要 (batch, channel, length) 格式', True),
        ('        # 回傳：x.shape = (1, 187)，y = 純量標籤', True),
    ]
    for line, is_comment in code_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(1)
        run = p.add_run()
        if line == '':
            run.text = ''
        else:
            run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) if is_comment else RGBColor(0x1A, 0x1A, 0x1A)
        run.font.bold = is_comment
        _set_fonts(run, 'Courier New')
    add_page_num(s, 12, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 13 — 程式碼：模型定義
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '程式碼解說（二）模型定義 ECG1DCNN')
    tf = add_content_box(s)
    code_lines = [
        ('self.conv1 = nn.Conv1d(in_channels=1, out_channels=32, kernel_size=5)', False),
        ('# 輸入 (B,1,187) → 輸出 (B,32,183)  長度: 187-5+1=183', True),
        ('self.pool1 = nn.MaxPool1d(kernel_size=2)', False),
        ('# 輸出 (B,32,91)  長度: 183//2=91', True),
        ('', False),
        ('self.conv2 = nn.Conv1d(32, 64, kernel_size=5)', False),
        ('# 輸入 (B,32,91) → 輸出 (B,64,87)  長度: 91-5+1=87', True),
        ('self.pool2 = nn.MaxPool1d(2)  # → (B,64,43)', False),
        ('', False),
        ('self.conv3 = nn.Conv1d(64, 128, kernel_size=3)', False),
        ('# 輸入 (B,64,43) → 輸出 (B,128,41)  長度: 43-3+1=41', True),
        ('self.pool3 = nn.MaxPool1d(2)  # → (B,128,20)', False),
        ('', False),
        ('self.fc1 = nn.Linear(128 * 20, 256)  # 2560 → 256', False),
        ('# Flatten 後：128 channels × 20 time steps = 2560 維特徵向量', True),
        ('self.dropout = nn.Dropout(p=0.5)     # 隨機關閉 50% 神經元，防過擬合', False),
        ('self.fc2 = nn.Linear(256, 64)        # 256 → 64', False),
        ('self.fc3 = nn.Linear(64, 5)          # 64 → 5 類別 logits', False),
    ]
    for line, is_comment in code_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) if is_comment else RGBColor(0x1A, 0x1A, 0x1A)
        run.font.bold = is_comment
        _set_fonts(run, 'Courier New')
    add_page_num(s, 13, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 14 — 程式碼：訓練迴圈
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '程式碼解說（三）訓練迴圈')
    tf = add_content_box(s)
    code_lines = [
        ('def train_epoch(model, loader, criterion, optimizer):', False),
        ('    model.train()                    # 開啟訓練模式，Dropout 生效', True),
        ('    for X_batch, y_batch in loader:', False),
        ('        # X_batch: (B,1,187)  y_batch: (B,)', True),
        ('        optimizer.zero_grad()        # ① 清除上次迭代的梯度', True),
        ('        logits = model(X_batch)      # ② 前向傳播 → logits: (B,5)', True),
        ('        loss = criterion(logits, y_batch)  # ③ 計算 CrossEntropyLoss', True),
        ('        loss.backward()              # ④ 反向傳播，計算所有參數梯度', True),
        ('        optimizer.step()             # ⑤ Adam 更新參數權重', True),
        ('', False),
        ('# 評估模式（不更新參數）', True),
        ('def eval_epoch(model, loader, criterion):', False),
        ('    model.eval()                     # 關閉 Dropout', True),
        ('    with torch.no_grad():            # 不計算梯度，節省記憶體', True),
        ('        logits = model(X_batch)', False),
        ('        # argmax(1) 取各樣本機率最高的類別索引', True),
        ('        correct += (logits.argmax(1) == y_batch).sum().item()', False),
    ]
    for line, is_comment in code_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0x00, 0x70, 0xC0) if is_comment else RGBColor(0x1A, 0x1A, 0x1A)
        run.font.bold = is_comment
        _set_fonts(run, 'Courier New')
    add_page_num(s, 14, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 15 — 訓練曲線
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '訓練結果（一）Loss / Accuracy 曲線')
    tf = add_content_box(s, top_cm=5.0, height_cm=3.2)
    add_bullet(tf, '訓練設定：Batch Size=256，Epochs=20，Optimizer=Adam（lr=1e-3），裝置=Apple MPS GPU', bold=True, size=16)
    add_bullet(tf, 'Loss 曲線：前 5 個 Epoch 急速下降，之後趨於平緩，Train/Val 曲線接近，無明顯過擬合。', size=16)
    add_bullet(tf, 'Accuracy 曲線：第 1 Epoch 即達 94.2%（Val），最終收斂至 98.36%（Val）。', size=16)
    img_path = os.path.join(IMG_DIR, 'training_curve.png')
    if os.path.exists(img_path):
        add_image(s, img_path, left_cm=1.5, top_cm=8.5, width_cm=30.5)
    add_page_num(s, 15, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 16 — 混淆矩陣 + 分類報告
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '訓練結果（二）混淆矩陣與分類報告')
    tf = add_content_box(s, top_cm=5.0, height_cm=4.0)
    add_bullet(tf, '測試集整體 Accuracy：98.36%（21892 筆）', bold=True)
    add_bullet(tf, 'Normal F1=0.99　Ventricular F1=0.95　Unknown F1=0.99　Supraventricular F1=0.82　Fusion F1=0.77', size=16)
    add_bullet(tf, '少數類別（Supraventricular, Fusion）F1 較低，主因訓練樣本極少（2,223 / 641 筆），Recall 不足。', size=16)
    add_bullet(tf, '改善方向：Oversampling（如 SMOTE）或調整 Class Weight 可提升少數類別辨識率。', size=16)
    img_path = os.path.join(IMG_DIR, 'confusion_matrix.png')
    if os.path.exists(img_path):
        add_image(s, img_path, left_cm=7.0, top_cm=9.0, width_cm=20.0)
    add_page_num(s, 16, TOTAL)

    # ══════════════════════════════════════════════════
    # Slide 17 — 結論
    # ══════════════════════════════════════════════════
    s = blank_slide(prs)
    add_title(s, '結論 Conclusion')
    tf = add_content_box(s)
    bullets = [
        ('本次報告完整實作一套 ECG 心律不整自動分類系統，涵蓋從資料取得到模型部署的完整流程。', 0, False),
        ('', 0, False),
        ('主要成果', 0, True),
        ('資料集：MIT-BIH Arrhythmia Database，共 109,446 筆心搏，透過 Kaggle 取得預處理 CSV。', 1, False),
        ('模型：1D CNN（3 卷積塊 + 2 全連接層），參數量小、訓練快，適合即時 ECG 分析。', 1, False),
        ('性能：測試集 Accuracy 98.36%，Normal / Ventricular 類別 F1 ≥ 0.95。', 1, False),
        ('', 0, False),
        ('關鍵技術重點', 0, True),
        ('Conv1D 卷積核沿時間軸提取 ECG 局部波形特徵（QRS 寬度、振幅、形狀）。', 1, False),
        ('MaxPool1D 降維保留最顯著特徵，同時減少計算量。', 1, False),
        ('CrossEntropyLoss + Adam 優化組合使模型快速收斂，20 Epochs 即可達到接近最優性能。', 1, False),
        ('Dropout 有效防止過擬合，Train/Val Accuracy 差距維持在 0.5% 以內。', 1, False),
        ('', 0, False),
        ('未來改進方向', 0, True),
        ('處理類別不平衡（SMOTE / Class Weight）以提升 Supraventricular、Fusion 的 Recall。', 1, False),
        ('引入 Transformer 或 LSTM 以捕捉更長範圍的心律依賴關係。', 1, False),
        ('結合多導聯 ECG（12-lead）提升診斷準確率。', 1, False),
    ]
    for text, level, bold in bullets:
        add_bullet(tf, text, level=level, bold=bold)
    add_page_num(s, 17, TOTAL)

    # ── 儲存 ──
    prs.save(OUT_PATH)
    print(f'PPT 已儲存：{OUT_PATH}')


if __name__ == '__main__':
    build_ppt()
