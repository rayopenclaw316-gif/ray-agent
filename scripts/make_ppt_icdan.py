#!/usr/bin/env python3
"""ICDAN 論文 PPT — 依 PDF 全文重製，含原圖"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image as PILImage
import re, os

# ── 顏色 ──────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x0E, 0x28, 0x41)
RED       = RGBColor(0xFF, 0x00, 0x00)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY      = RGBColor(0x77, 0x77, 0x77)

IMG = "/Users/rayopenclaw/ray-agent/papers/pdf_imgs"
TOTAL = 12

RED_PAT = re.compile(
    r'(?<!\w)(sEMG|EMG|EEG|1D-CNN|CNN|SNN|SSR|SSI|CDAN|ICDAN|MMD|'
    r'Domain Adaptation|GAN|GRL|DANN|DDC|DAN|Softmax|Transformer|GRU|LSTM|'
    r'Sensor Fusion|Language Model|Moving Window|BN|FC)(?!\w)'
)
CJK = re.compile(r'[一-鿿　-〿＀-￯]')

# ── 字型工具 ──────────────────────────────────────────────────
def _set_fonts(run):
    rPr = run._r.get_or_add_rPr()
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for tag, name in [('latin','Times New Roman'),('ea','標楷體'),('cs','標楷體')]:
        el = rPr.find(f'{{{ns}}}{tag}')
        if el is None: el = etree.SubElement(rPr, f'{{{ns}}}{tag}')
        el.set('typeface', name)

def _split(text):
    segs, buf, mode = [], '', None
    for ch in text:
        m = bool(CJK.match(ch))
        if mode is None: mode = m
        if m == mode: buf += ch
        else: segs.append((buf, mode)); buf, mode = ch, m
    if buf: segs.append((buf, mode))
    return segs or [(text, False)]

def _run(para, text, size=Pt(17), bold=False, color=BLACK, italic=False):
    for chunk, _ in _split(text):
        r = para.add_run(); r.text = chunk
        r.font.size = size; r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
        _set_fonts(r)

def txt(para, text, size=Pt(17), bold=False, color=BLACK, italic=False):
    parts, last = [], 0
    for m in RED_PAT.finditer(text):
        if m.start() > last: parts.append((text[last:m.start()], False))
        parts.append((m.group(), True)); last = m.end()
    if last < len(text): parts.append((text[last:], False))
    for t, red in parts:
        _run(para, t, size=size, bold=bold, color=RED if red else color, italic=italic)

# ── 投影片工廠 ────────────────────────────────────────────────
def ns(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def title(slide, text, size=Pt(28)):
    b = slide.shapes.add_textbox(Cm(2.3), Cm(1.0), Cm(29.2), Cm(3.7))
    b.text_frame.word_wrap = True
    txt(b.text_frame.paragraphs[0], text, size=size, bold=True, color=DARK_BLUE)

def pnum(slide, n):
    b = slide.shapes.add_textbox(Cm(23.9), Cm(17.7), Cm(3.0), Cm(1.0))
    p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f'{n} / {TOTAL}'
    r.font.size = Pt(11); r.font.color.rgb = GRAY; _set_fonts(r)

def cp(tf, first=False): return tf.paragraphs[0] if first else tf.add_paragraph()

def content_tf(slide, l=2.3, t=5.1, w=29.2, h=12.1):
    b = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    b.text_frame.word_wrap = True; return b.text_frame

def hline(slide, l, t, w):
    s = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BLUE; s.line.fill.background()

def img(slide, fname, l, t, max_w, max_h):
    path = os.path.join(IMG, fname)
    im = PILImage.open(path); iw, ih = im.size; r = iw/ih
    w = min(max_w, max_h * r); h = w / r
    slide.shapes.add_picture(path, Cm(l), Cm(t), Cm(w), Cm(h))
    return w, h

def lines(tf, data, default=Pt(17)):
    first = True
    for item in data:
        if isinstance(item, tuple):
            line, size, bold, color = item
        else:
            line, size, bold, color = item, default, False, BLACK
            if item.startswith('【'): bold, color = True, DARK_BLUE
            elif item.startswith('  '): size = Pt(15)
            elif item.startswith('    '): size = Pt(14)
        p = cp(tf, first); first = False
        if line == '': p.add_run(); continue
        txt(p, line, size=size, bold=bold, color=color)

# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)

# ─── Slide 1：封面 ────────────────────────────────────────────
sl = ns(prs)
b = sl.shapes.add_textbox(Cm(2.3), Cm(1.6), Cm(29.2), Cm(4.5))
b.text_frame.word_wrap = True
p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
txt(p, 'EMG-Based Cross-Subject Silent Speech Recognition\n'
       'Using Conditional Domain Adversarial Network (ICDAN)',
    size=Pt(24), bold=True, color=DARK_BLUE)

b2 = sl.shapes.add_textbox(Cm(2.3), Cm(6.6), Cm(29.2), Cm(1.5))
b2.text_frame.word_wrap = True
p2 = b2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
txt(p2, '基於 sEMG 的跨受試者無聲語音辨識：條件域對抗網路（ICDAN）方法',
    size=Pt(17), bold=True, color=DARK_BLUE)

hline(sl, 5.0, 8.5, 24.0)

for y, t, sz, c in [
    (8.9, 'IEEE Transactions on Cognitive and Developmental Systems, Vol. 15, No. 4, December 2023', 13, DARK_BLUE),
    (10.1, 'Yakun Zhang*, Huihui Cai*, Jinghan Wu, Liang Xie, Minpeng Xu, Dong Ming, Ye Yan, Erwei Yin†', 12, GRAY),
    (11.2, 'National Innovation Institute of Defense Technology · Tianjin University · Tianjin AI Innovation Center', 11, GRAY),
    (12.2, 'DOI: 10.1109/TCDS.2023.3316701  ·  (* co-first authors  † corresponding author)', 10, GRAY),
]:
    b3 = sl.shapes.add_textbox(Cm(2.3), Cm(y), Cm(29.2), Cm(1.1))
    p3 = b3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    _run(p3, t, size=Pt(sz), color=c)
pnum(sl, 1)

# ─── Slide 2：目錄 ────────────────────────────────────────────
sl = ns(prs); title(sl, '目錄'); tf = content_tf(sl)
toc = [
    ('1.  研究背景與動機',                   ' 3'),
    ('2.  資料集',                            ' 4'),
    ('3.  訊號前處理與時序特徵',              ' 5'),
    ('4.  1D-CNN 特徵提取架構',               ' 6'),
    ('5.  ICDAN 完整模型架構',                ' 7'),
    ('6.  損失函數詳解',                      ' 8'),
    ('7.  多資料跨人辨識實驗結果',            ' 9'),
    ('8.  少樣本跨人辨識結果',               '10'),
    ('9.  模型比較與分析',                   '11'),
    ('10. 結論與研究啟發',                   '12'),
]
first = True
for item, pg in toc:
    p = cp(tf, first); first = False
    txt(p, item + '  ' + '·'*52 + '  ' + pg, size=Pt(17))
pnum(sl, 2)

# ─── Slide 3：研究背景 + Fig.1 ───────────────────────────────
sl = ns(prs); title(sl, '研究背景與動機')
tf = content_tf(sl, l=2.3, t=5.1, w=14.5, h=12.1)
lines(tf, [
    '【為何需要跨受試者 SSR？】',
    '  • sEMG 捕捉臉部肌肉收縮，無需發聲即可解碼語音意圖',
    '  • 同受試者（within-subject）辨識率已達高水準',
    '  • 跨受試者因個體差異受關注不足',
    '',
    '【三大跨受試者挑戰】',
    '  • 生理差異：臉部肌肉解剖結構因人而異',
    '  • 電極偏差：每次貼附位置輕微偏移',
    '  • 健康狀態：說話習慣、皮膚阻抗不同',
    '',
    '【現有研究的侷限】',
    '  • 多數研究受試者 ≤ 10 人、標籤 ≤ 2 類',
    '  • 缺乏有效特徵提取與跨受試者辨識方法',
    '',
    '【本論文貢獻】',
    '  • 建立 70 人、101 類中文 EMG 資料集',
    '  • 1D-CNN 時序特徵學習',
    '  • ICDAN = CDAN + MMD 改善少樣本跨人辨識',
])
img(sl, 'raw_p02_i01.jpeg', l=17.2, t=5.0, max_w=15.0, max_h=12.2)
pnum(sl, 3)

# ─── Slide 4：資料集 + Fig.2 ─────────────────────────────────
sl = ns(prs); title(sl, '資料集')
tf = content_tf(sl, l=2.3, t=5.1, w=15.0, h=12.1)
lines(tf, [
    '【受試者】',
    '  • 70 名普通話母語者：50 男 / 20 女',
    '  • 年齡：20–35 歲（平均 27.5 歲）',
    '  • 倫理審查：天津大學 IRB（No. TJUE-2021-138）',
    '',
    '【語料】',
    '  • 101 類常用生活短語與人機互動指令',
    '  • 每句 3–5 個詞，例：「我想吃水果」',
    '',
    '【採集規格】',
    '  • 無線 EMG 設備，取樣率 1000 Hz',
    '  • 6 通道（6 個臉頸肌群）：',
    '    Ch1 Mentalis（頤肌）',
    '    Ch2 Risorius（笑肌）',
    '    Ch3 Levator labii superioris（上唇提肌）',
    '    Ch4 Anterior belly of digastric（二腹肌前腹）',
    '    Ch5 Mylohyoid（下頜舌骨肌）',
    '    Ch6 Platysma（頸闊肌）',
    '',
    '  • 每人 10 sessions，每次默讀全部短語（只張嘴，不發聲）',
    '  • 有效樣本：69,875 筆（訓練 60,026 / 測試 9,849）',
])
img(sl, 'raw_p03_i01.jpeg', l=17.5, t=5.2, max_w=14.5, max_h=12.0)
pnum(sl, 4)

# ─── Slide 5：前處理 + Fig.3 + Fig.4 ────────────────────────
sl = ns(prs); title(sl, '訊號前處理與時序特徵')
tf = content_tf(sl, l=2.3, t=5.1, w=14.5, h=12.1)
lines(tf, [
    '【原始 EMG 的三種雜訊】',
    '  1. 工業頻率干擾：50 Hz 和 150 Hz（交流電源）',
    '  2. 直流偏移（DC Bias）：造成基線任意偏移',
    '  3. 低頻基線漂移（0–10 Hz）：數百微伏漂移',
    '',
    '【前處理步驟】',
    '  1. 4 階 Butterworth 帶通濾波（10–400 Hz）',
    '     → 去除 DC 偏移與低頻雜訊',
    '  2. 陷波濾波（Notch filter）',
    '     → 消除 50 Hz / 150 Hz 工頻干擾',
    '',
    '【時序特徵（Time-Series Features）】',
    '  TS(x) = {x₁, x₂, ..., x_t}，t = 2000',
    '',
    '  優勢：保留更豐富的原始 EMG 分類資訊',
    '  不同受試者相同指令的訊號包絡相似',
    '（如右圖：Ch1、Ch5 為兩人的主導通道）',
])
# Fig.3 上方
img(sl, 'raw_p04_i01.jpeg', l=17.2, t=5.0, max_w=14.5, max_h=6.5)
# Fig.4 下方
img(sl, 'raw_p04_i02.jpeg', l=17.2, t=11.8, max_w=14.5, max_h=4.5)
pnum(sl, 5)

# ─── Slide 6：1D-CNN 架構 + Fig.5（上半部）────────────────────
sl = ns(prs); title(sl, '1D-CNN 特徵提取架構')
tf = content_tf(sl, l=2.3, t=5.1, w=14.5, h=12.1)
lines(tf, [
    '【為何用 1D-CNN？】',
    '  • 傳統方法手動特徵提取 → 遺失部分資訊',
    '  • 1D-CNN 自動學習時序局部特徵，更適合 EMG',
    '',
    '【架構細節（Table I）】',
    ('  輸入：6 × 2000（6 通道 × 2000 時間點）', Pt(15), False, BLACK),
    ('  卷積核大小：1 × 3', Pt(15), False, BLACK),
    ('  MaxPool：1 × 2（第一層）、1 × 3（其餘層）', Pt(15), False, BLACK),
    ('  通道數：64 → 128 → 256 → 512（逐層加深）', Pt(15), False, BLACK),
    '',
    '【訓練策略】',
    '  • Dropout = 0.5（pooling 後）',
    '  • Batch Normalization（BN）：conv 層後，改善梯度',
    '  • L2 正規化：FC 層，防止過擬合',
    '  • Adam 優化器：lr = 0.0001，weight decay = 0.0005',
    '  • Batch size = 16',
    '',
    '【單受試者內效能】',
    '  • 無跨受試者：94.45%（上界）',
    '  • 跨受試者（60 人訓練）：87.80%',
])
img(sl, 'fig5_clean.png', l=17.2, t=5.0, max_w=14.5, max_h=12.0)
pnum(sl, 6)

# ─── Slide 7：ICDAN 完整架構 + Fig.5 大圖 ────────────────────
sl = ns(prs); title(sl, 'ICDAN 完整模型架構（Fig. 5）')
# 圖放大佔大部分版面
img(sl, 'fig5_clean.png', l=2.3, t=4.8, max_w=29.2, max_h=11.0)
# 底部說明文字
b = sl.shapes.add_textbox(Cm(2.3), Cm(16.0), Cm(29.2), Cm(1.5))
b.text_frame.word_wrap = True
p = b.text_frame.paragraphs[0]
txt(p, 'ICDAN 三路損失並行：'
       '① Source Data → 1D-CNN → FC → Classifier → ℒ_cls（分類損失）  '
       '② Source + Target → MMD Loss（特徵均值對齊）  '
       '③ Source + Target → GRL → Discriminator → CDAN Loss（對抗域判別）',
    size=Pt(13), color=DARK_BLUE)
pnum(sl, 7)

# ─── Slide 8：損失函數 ────────────────────────────────────────
sl = ns(prs); title(sl, '損失函數詳解')
tf = content_tf(sl)
loss_data = [
    ('【完整優化目標】', Pt(17), True, DARK_BLUE),
    ('', Pt(17), False, BLACK),
    ('  min_{G,T}  max_D  [ LY(T(G(x)), Yˢ)  −  λ · ω · LD(M(hˢ·ᵗ), Y^{s,t}_D)  +  λ · LM(fˢ, fᵗ) ]',
     Pt(15), True, DARK_BLUE),
    ('', Pt(17), False, BLACK),
    ('【① 分類損失 ℒ_cls（Eq.2）】', Pt(17), True, DARK_BLUE),
    ('  LY = −(1/nˢ) Σᵢ Σ_c  yˢᵢ,c · log(gˢᵢ,c)', Pt(15), False, BLACK),
    ('  nˢ：source 資料數；C：標籤類別數（101類）', Pt(14), False, GRAY),
    ('', Pt(17), False, BLACK),
    ('【② MMD 統計對齊損失 ℒ_MMD（Eq.3）】', Pt(17), True, DARK_BLUE),
    ('  LM = ‖ (1/|Fˢ|) Σ φ(fˢ)  −  (1/|Fᵗ|) Σ φ(fᵗ) ‖²', Pt(15), False, BLACK),
    ('  φ：核函數映射（RKHS）；LM→0 表示兩域特徵分佈完全對齊', Pt(14), False, GRAY),
    ('', Pt(17), False, BLACK),
    ('【③ CDAN 對抗損失 ℒ_adv（Eq.4）】', Pt(17), True, DARK_BLUE),
    ('  LD = −(1/nˢ) Σ log D(gˢᵢ)  −  (1/nᵗ) Σ log(1 − D(gᵗⱼ))', Pt(15), False, BLACK),
    ('  D：域判別器；gˢ/gᵗ：source/target 特徵；GRL 使反向梯度反轉', Pt(14), False, GRAY),
    ('', Pt(17), False, BLACK),
    ('【④ CDAN 條件化（Eq.6）】', Pt(17), True, DARK_BLUE),
    ('  M(h) = (1/√d) · R_ff ⊗ R_gg（隨機多線性映射）', Pt(15), False, BLACK),
    ('  h = T(G(x))（特徵 × 分類器預測的外積）；d：特徵維度', Pt(14), False, GRAY),
    ('  ω：樣本熵值權重；λ：控制 source/target 對齊比例的超參數', Pt(14), False, GRAY),
]
first = True
for item in loss_data:
    if isinstance(item, tuple) and len(item) == 4:
        line, size, bold, color = item
    else:
        line, size, bold, color = item, Pt(17), False, BLACK
    p = cp(tf, first); first = False
    if line == '': p.add_run(); continue
    txt(p, line, size=size, bold=bold, color=color)
pnum(sl, 8)

# ─── Slide 9：多資料跨人辨識實驗結果 ─────────────────────────
sl = ns(prs); title(sl, '多資料跨人辨識實驗結果')
tf = content_tf(sl, l=2.3, t=5.1, w=29.2, h=12.1)
lines(tf, [
    '【特徵比較實驗（Table II，VGGNet 骨幹）】',
    ('  特徵           模型      訓練集準確率   跨受試者準確率   SD', Pt(14), True, DARK_BLUE),
    ('  TC-4（時域組合）   VGGNet      —            83.61%         —', Pt(14), False, BLACK),
    ('  FC-4（頻域組合）   VGGNet      —            80.4x%         —', Pt(14), False, BLACK),
    ('  MFCC             VGGNet      —            79.4x%         —', Pt(14), False, BLACK),
    ('  時序特徵（本文）   VGGNet      —            83.61%         —', Pt(14), False, BLACK),
    ('  時序特徵（本文）   1D-CNN    94.45%         87.80%        6.57  ← 最佳', Pt(14), True, DARK_BLUE),
    '',
    '【五折交叉驗證（Table III，資料以 6:1 分割）】',
    ('  Group 1: 84.65%   Group 2: 79.20%   Group 3: 87.30%', Pt(14), False, BLACK),
    ('  Group 4: 83.50%   Group 5: 87.56%   平均：84.54%（3/5 組 > 85%）', Pt(14), False, BLACK),
    '',
    '【1D-CNN vs 傳統方法的優勢】',
    '  • 相較 VGGNet，1D-CNN 跨受試者準確率提升約 5%',
    '  • 架構更簡單，避免過擬合，訓練速度更快',
    '  • 時序特徵比傳統時域/頻域特徵更保留 EMG 分類資訊',
    '',
    '【不同受試者的個體差異】',
    ('  • 最佳受試者（Sub.9）：86.28%；最差（Sub.10，臉部鬍鬚+厚脂肪層）：30.04%', Pt(14), False, BLACK),
    ('  • 差距高達 56.24% → 跨受試者挑戰的核心原因', Pt(14), False, BLACK),
])
pnum(sl, 9)

# ─── Slide 10：少樣本跨人辨識 + Fig.6 ────────────────────────
sl = ns(prs); title(sl, '少樣本跨人辨識結果')
tf = content_tf(sl, l=2.3, t=5.1, w=15.5, h=12.1)
lines(tf, [
    '【實驗設置（少樣本場景）】',
    '  • Source domain：20 人有標記資料（訓練）',
    '  • Target domain：10 人無標記資料',
    '    → 隨機 1:1 分成訓練集與測試集',
    '',
    '【不同訓練集規模的影響（Table IV）】',
    ('  訓練人數   訓練集準確率   跨受試者準確率   SD', Pt(14), True, DARK_BLUE),
    ('  60 人      94.45%         87.80%          6.37', Pt(14), False, BLACK),
    ('  40 人      94.52%         82.74%          7.49', Pt(14), False, BLACK),
    ('  20 人      94.56%         70.68%         15.75  ← 大幅下降', Pt(14), False, BLACK),
    '',
    '  訓練資料從 60,000 降至 20,000 → 準確率下降 17.12%',
    '',
    '【ICDAN 的改善效果（Fig. 6）】',
    '  • 1D-CNN 基準（20 人訓練）：71.42%',
    ('  • ICDAN（加入遷移學習）：+14.88% → 86.30%', Pt(15), True, DARK_BLUE),
    '  • 藍色柱＝1D-CNN；橘紅色柱＝ICDAN',
    '  • Sub.10 因個體差異特殊，1D-CNN 僅 29%',
])
img(sl, 'fig6_clean.png', l=18.2, t=5.5, max_w=13.5, max_h=8.0)
# 圖說
b = sl.shapes.add_textbox(Cm(18.2), Cm(13.8), Cm(13.5), Cm(0.8))
p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, 'Fig. 6  十位受試者跨人辨識準確率（藍=1D-CNN，橘紅=ICDAN）',
     size=Pt(10), color=GRAY)
pnum(sl, 10)

# ─── Slide 11：模型比較與分析 ─────────────────────────────────
sl = ns(prs); title(sl, '模型比較與分析')
tf = content_tf(sl, l=2.3, t=5.1, w=15.5, h=12.1)
lines(tf, [
    '【遷移學習方法比較（Table V，20 人訓練）】',
    ('  方法        平均準確率   SD     特性', Pt(14), True, DARK_BLUE),
    ('  1D-CNN      71.42%      14.23  基準（無遷移）', Pt(13), False, BLACK),
    ('  DDC         76.17%      14.25  MMD 對齊特徵', Pt(13), False, BLACK),
    ('  DAN         75.47%      13.63  深度自適應網路', Pt(13), False, BLACK),
    ('  DANN        76.74%      13.83  基本對抗訓練', Pt(13), False, BLACK),
    ('  CDAN        85.38%      11.83  條件域對抗', Pt(13), False, BLACK),
    ('  ICDAN       86.30%      11.31  CDAN + MMD  ← 最佳', Pt(13), True, DARK_BLUE),
    '',
    '【關鍵觀察】',
    '  • DDC/DAN/DANN 提升有限（約 +5%）',
    '  • CDAN 大幅躍升（+14%）：EMG 多模態特性需條件化對齊',
    '  • ICDAN 再提升 ~1%：MMD 補充統計均值對齊',
    '  • ICDAN 的 SD 最小（11.31）→ 受試者間穩定性最高',
    '',
    '【ROC / PR 曲線（Fig. 7）】',
    ('  1D-CNN  AUC = 0.98，AP = 0.63', Pt(14), False, BLACK),
    ('  CDAN    AUC = 0.99，AP = 0.83  → 精準率-召回率大幅改善', Pt(14), False, BLACK),
])
img(sl, 'raw_p08_i01.jpeg', l=18.2, t=5.2, max_w=13.5, max_h=11.0)
b = sl.shapes.add_textbox(Cm(18.2), Cm(16.5), Cm(13.5), Cm(0.8))
p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
_run(p, 'Fig. 9  不同訓練集規模下三種模型的 Box Plot（藍=1D-CNN，橘=CDAN，灰=ICDAN）',
     size=Pt(9), color=GRAY)
pnum(sl, 11)

# ─── Slide 12：結論 ───────────────────────────────────────────
sl = ns(prs); title(sl, '結論與研究啟發')
tf = content_tf(sl)
lines(tf, [
    '【論文核心貢獻（直接引用結論）】',
    '  1. 建立首個 70 人、101 類中文 sEMG 跨受試者語音資料集',
    '  2. 時序特徵 + 1D-CNN：跨受試者準確率達 87.80%（60 人訓練）',
    '  3. ICDAN（CDAN + MMD）：少樣本場景下準確率提升 14.88%',
    '  4. 首次將 CDAN 成功應用於 EMG 跨受試者解碼，突破多模態遷移',
    '',
    '【侷限與未來工作（論文原文）】',
    '  • 目前仍需少量 Target 無標記資料（非零樣本）',
    '  • 未來目標：突破零樣本跨受試者 SSR',
    '',
    '【對 CSL-EMG 研究的直接啟發】',
    '  • ICDAN 架構可直接移植至中文 sEMG 跨人辨識場景',
    '  • 條件化對抗訓練（CDAN）對多模態 EMG 至關重要，DANN 不夠',
    '  • 可嘗試以 Transformer 或 GRU 替換 1D-CNN 作為更強特徵提取骨幹',
    '  • 建立大規模多人資料集是提升跨受試者效能的根本途徑',
    '',
    ('論文連結：https://ieeexplore.ieee.org/document/10254583/',
     Pt(12), False, GRAY),
    ('DOI: 10.1109/TCDS.2023.3316701  ·  本簡報依 PDF 全文製作',
     Pt(12), False, GRAY),
])
pnum(sl, 12)

# ─── 儲存 ─────────────────────────────────────────────────────
out = '/Users/rayopenclaw/ray-agent/pptx/ICDAN_CrossSubject_SSR.pptx'
prs.save(out)
print(f'Done → {out}')
