#!/usr/bin/env python3
"""
從 PDF 提取圖片 + 建立 matplotlib 圖表，重新生成帶圖的 PPT
"""
import os, io
import fitz  # PyMuPDF
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Emu

# ── 路徑設定 ─────────────────────────────────────────
PDF_PATH = '/Users/rayopenclaw/Downloads/1-s2.0-S0167639325000457-main.pdf'
IMG_DIR  = '/Users/rayopenclaw/ray-agent/papers/ppt_imgs'
OUT_PATH = '/Users/rayopenclaw/ray-agent/papers/AuxCEMGR_解析.pptx'
os.makedirs(IMG_DIR, exist_ok=True)

# ── Step 1：從 PDF 提取圖片 ──────────────────────────
print('=== 從 PDF 提取圖片 ===')
doc = fitz.open(PDF_PATH)
img_map = {}  # xref -> 儲存路徑

for page_num in range(len(doc)):
    page = doc[page_num]
    img_list = page.get_images(full=True)
    for idx, img_info in enumerate(img_list):
        xref = img_info[0]
        if xref in img_map:
            continue
        base_img = doc.extract_image(xref)
        img_bytes = base_img['image']
        ext = base_img['ext']
        w, h = base_img['width'], base_img['height']
        # 過濾太小的圖（通常是裝飾用圖示）
        if w < 100 or h < 100:
            continue
        fname = f'{IMG_DIR}/p{page_num+1}_x{xref}.{ext}'
        with open(fname, 'wb') as f:
            f.write(img_bytes)
        img_map[xref] = {'path': fname, 'page': page_num+1, 'w': w, 'h': h}
        print(f'  [p{page_num+1}] xref={xref} {w}×{h} → {os.path.basename(fname)}')

doc.close()

# 整理成 list，依頁面排序
img_list_sorted = sorted(img_map.values(), key=lambda x: x['page'])
print(f'共提取 {len(img_list_sorted)} 張圖片')
for i, img in enumerate(img_list_sorted):
    print(f'  [{i}] p{img["page"]} {img["w"]}x{img["h"]} {os.path.basename(img["path"])}')

# ── Step 2：建立 matplotlib 圖表 ────────────────────

FONT_CHN = 'Heiti TC'  # macOS 中文字體

def setup_font():
    import matplotlib.font_manager as fm
    fonts = [f.name for f in fm.fontManager.ttflist]
    for candidate in ['Heiti TC','PingFang TC','STHeiti','Arial Unicode MS']:
        if candidate in fonts:
            return candidate
    return None

CHN_FONT = setup_font() or 'DejaVu Sans'

plt.rcParams['font.family'] = CHN_FONT
plt.rcParams['axes.unicode_minus'] = False

# ── 圖表 A：主要實驗結果（橫條圖）────────────────────
def make_main_results_chart():
    fig, ax = plt.subplots(figsize=(9, 5))
    models = [
        'Wadkins 2019\n(LSTM+CTC)',
        'Baseline\n(無增強)',
        'Baseline\n(完整增強)',
        'AuxCEMGR\n(無增強)',
        'AuxCEMGR\n(完整增強)',
    ]
    cer_values = [66.8, 66.5, 44.5, 62.5, 38.0]
    colors = ['#999999', '#E07070', '#F0A060', '#80B0E0', '#2E75B6']
    bars = ax.barh(models, cer_values, color=colors, edgecolor='white',
                   height=0.55)
    for bar, val in zip(bars, cer_values):
        ax.text(val + 0.8, bar.get_y() + bar.get_height()/2,
                f'{val}%', va='center', fontsize=13, fontweight='bold',
                color='#222222')
    ax.set_xlabel('字元錯誤率 CER (%)', fontsize=13)
    ax.set_title('主要實驗結果比較（測試集，越低越好）',
                 fontsize=14, fontweight='bold', pad=12)
    ax.set_xlim(0, 80)
    ax.axvline(38.0, color='#1F3864', linestyle='--', linewidth=1.5,
               label='最佳 38.0%')
    ax.legend(fontsize=11)
    ax.invert_yaxis()
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    out = f'{IMG_DIR}/chart_main_results.png'
    plt.savefig(out, dpi=150, bbox_inches='tight')
    plt.close()
    print(f'  [圖表A] {out}')
    return out

# ── 圖表 B：資料增強效果（分組柱狀圖）───────────────
def make_augmentation_chart():
    fig, ax = plt.subplots(figsize=(8, 4.5))
    x = np.arange(4)
    width = 0.35
    baseline = [66.5, 53.9, 49.6, 44.5]
    auxcemgr = [62.5, 45.8, 44.5, 38.0]
    labels = ['無增強', '-頻譜相減\n→+9.4%', '-Mixup\n→+5.1%', '完整增強']
    b1 = ax.bar(x - width/2, baseline, width, label='Baseline',
                color='#F0A060', edgecolor='white')
    b2 = ax.bar(x + width/2, auxcemgr, width, label='AuxCEMGR',
                color='#2E75B6', edgecolor='white')
    for bar in list(b1) + list(b2):
        h = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2, h + 0.5,
                f'{h}', ha='center', va='bottom', fontsize=10,
                fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylabel('CER (%)', fontsize=12)
    ax.set_title('三種資料增強對 CER 的影響（測試集）',
                 fontsize=13, fontweight='bold', pad=10)
    ax.set_ylim(0, 80)
    ax.legend(fontsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    out = f'{IMG_DIR}/chart_augmentation.png'
    plt.savefig(out, dpi=150, bbox_inches='tight')
    plt.close()
    print(f'  [圖表B] {out}')
    return out

# ── 圖表 C：電極通道貢獻（模擬數值，由論文圖 8 讀取）
def make_channel_chart():
    fig, ax = plt.subplots(figsize=(8, 4.5))
    channels = [f'CH{i}' for i in range(1, 9)]
    # 從論文 Fig.8 估讀的近似值
    baseline_cer  = [88.5, 85.5, 88.0, 86.5, 91.5, 91.0, 90.0, 89.0]
    auxcemgr_cer  = [87.0, 86.5, 87.5, 84.5, 91.0, 90.5, 88.5, 87.5]
    x = np.arange(len(channels))
    w = 0.35
    ax.bar(x - w/2, baseline_cer, w, label='Baseline',
           color='#F0A060', edgecolor='white')
    ax.bar(x + w/2, auxcemgr_cer, w, label='AuxCEMGR',
           color='#2E75B6', edgecolor='white')
    ax.set_xticks(x)
    ax.set_xticklabels(channels, fontsize=12)
    ax.set_ylabel('CER (%)', fontsize=12)
    ax.set_ylim(82, 94)
    ax.set_title('各電極通道單獨使用的 CER（越低越好）\n（CH4 最重要；CH5、CH6 喉部靜默時效果最差）',
                 fontsize=12, fontweight='bold', pad=10)
    ax.legend(fontsize=11)
    # 標記 CH4
    ax.annotate('CH4 最重要', xy=(3, 84.5), xytext=(3.8, 83.2),
                fontsize=10, color='#C00000',
                arrowprops=dict(arrowstyle='->', color='#C00000'))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    out = f'{IMG_DIR}/chart_channels.png'
    plt.savefig(out, dpi=150, bbox_inches='tight')
    plt.close()
    print(f'  [圖表C] {out}')
    return out

# ── 圖表 D：模型架構示意圖 ──────────────────────────
def make_architecture_diagram():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis('off')

    def box(cx, cy, w, h, text, fc, tc='white', fs=10):
        rect = mpatches.FancyBboxPatch(
            (cx - w/2, cy - h/2), w, h,
            boxstyle='round,pad=0.08',
            facecolor=fc, edgecolor='white', linewidth=1.5)
        ax.add_patch(rect)
        ax.text(cx, cy, text, ha='center', va='center',
                fontsize=fs, color=tc, fontweight='bold',
                wrap=True, multialignment='center')

    def arrow(x1, y1, x2, y2, color='#555555'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=color,
                                   lw=1.8))

    # 主流程
    box(0.9, 3.0, 1.5, 1.0, '臉部\nsEMG\n(8 通道)', '#666666', fs=9)
    arrow(1.65, 3.0, 2.1, 3.0)
    box(2.7,  3.0, 1.1, 1.0, 'CNN\n(2層)', '#2E75B6', fs=10)
    arrow(3.25, 3.0, 3.7, 3.0)
    box(4.5,  3.0, 1.5, 1.0, 'Transformer\n編碼器\n(6層)', '#1F3864', fs=9)
    arrow(5.25, 3.0, 5.7, 3.0)
    box(6.4,  3.0, 1.2, 1.0, 'CTC\n解碼器', '#2E75B6', fs=10)
    arrow(7.0, 3.0, 7.45, 3.0)
    box(8.2,  3.0, 1.4, 1.0, '中文\n字元序列', '#006B00', fs=9)

    # 輔助任務 1：拼音
    ax.annotate('', xy=(6.4, 4.7), xytext=(4.5, 3.5),
                arrowprops=dict(arrowstyle='->', color='#C56000', lw=1.5))
    box(6.4, 5.1, 1.5, 0.8, '拼音 CTC\n（輔助任務1）', '#C56000', fs=9)
    ax.text(7.35, 5.1, '→ 拼音序列', va='center', fontsize=9,
            color='#C56000')

    # 輔助任務 2：Session
    ax.annotate('', xy=(6.4, 1.3), xytext=(4.5, 2.5),
                arrowprops=dict(arrowstyle='->', color='#C00000', lw=1.5))
    box(6.4, 0.9, 1.8, 0.8, 'GRL + Session 分類\n（輔助任務2）', '#C00000', fs=8.5)
    ax.text(7.45, 0.9, '→ 消除電極偏差', va='center', fontsize=9,
            color='#C00000')

    # 說明文字
    ax.text(0.1, 5.7, 'AuxCEMGR 模型架構示意圖',
            fontsize=13, fontweight='bold', color='#1F3864')
    ax.text(4.5, 1.85, 'MFSC\n特徵', ha='center', fontsize=8.5,
            color='#555555', style='italic')

    plt.tight_layout(pad=0.3)
    out = f'{IMG_DIR}/chart_architecture.png'
    plt.savefig(out, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f'  [圖表D] {out}')
    return out

print('\n=== 建立 matplotlib 圖表 ===')
chart_main   = make_main_results_chart()
chart_aug    = make_augmentation_chart()
chart_ch     = make_channel_chart()
chart_arch   = make_architecture_diagram()

# ── Step 3：整頁渲染 PDF，取得所有圖（含向量圖）───────
# 論文圖在哪些頁（0-indexed）：
# Fig1(p2=idx1), Fig2(p3=idx2), Fig3(p4=idx3), Fig4(p5=idx4),
# Fig5~9(p7=idx6), Fig10(p8=idx7), Fig11~12(p9=idx8)
print('\n=== 整頁渲染 PDF 圖頁 ===')
RENDER_DPI = 150  # 渲染解析度
RENDER_PAGES = {
    'fig1':  (1, None),      # page idx=1 (p2)  整頁
    'fig2':  (2, None),      # page idx=2 (p3)  整頁
    'fig3':  (3, None),      # page idx=3 (p4)  整頁
    'fig4':  (4, None),      # page idx=4 (p5)  整頁
    'fig5to9': (6, None),    # page idx=6 (p7)  整頁
    'fig10': (7, None),      # page idx=7 (p8)  整頁
    'fig11_12': (8, None),   # page idx=8 (p9)  整頁
}

doc2 = fitz.open(PDF_PATH)
rendered = {}
mat = fitz.Matrix(RENDER_DPI/72, RENDER_DPI/72)  # scale factor

for name, (page_idx, clip) in RENDER_PAGES.items():
    if page_idx >= len(doc2):
        rendered[name] = None
        continue
    page = doc2[page_idx]
    pix = page.get_pixmap(matrix=mat, alpha=False)
    out_path = f'{IMG_DIR}/render_{name}.png'
    pix.save(out_path)
    rendered[name] = out_path
    print(f'  [{name}] p{page_idx+1} → {os.path.basename(out_path)} ({pix.width}×{pix.height})')

doc2.close()

# 論文圖對應（用渲染頁）
fig1_path   = rendered.get('fig1')
fig2_path   = rendered.get('fig2')
fig3_path   = rendered.get('fig3')
fig4_path   = rendered.get('fig4')
fig5_path   = rendered.get('fig5to9')
fig6_path   = rendered.get('fig5to9')   # 同一頁
fig8_path   = rendered.get('fig5to9')
fig9_path   = rendered.get('fig5to9')
fig10_path  = rendered.get('fig10')
fig11_path  = rendered.get('fig11_12')
fig12_path  = rendered.get('fig11_12')

print(f'\n論文圖對應：')
for name, path in [('Fig1',fig1_path),('Fig2',fig2_path),('Fig3',fig3_path),
                   ('Fig4',fig4_path),('Fig5',fig5_path),('Fig10',fig10_path),
                   ('Fig11/12',fig11_path)]:
    print(f'  {name}: {os.path.basename(path) if path else "未找到"}')

# ── Step 4：建立 PPT ─────────────────────────────────
print('\n=== 建立 PPT ===')

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)

C_DARK  = RGBColor(0x1F, 0x38, 0x64)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_ORNG  = RGBColor(0xC5, 0x60, 0x00)
C_GRN   = RGBColor(0x00, 0x6B, 0x00)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LIGHT = RGBColor(0xD6, 0xE4, 0xF7)

FONT  = '標楷體'
blank = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txtbox(slide, l, t, w, h, items):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        if item.get('sb', 0): p.space_before = Pt(item['sb'])
        run = p.add_run()
        indent = '    ' * item.get('indent', 0)
        blt    = item.get('blt', '')
        run.text = indent + (blt + ' ' if blt else '') + item.get('text', '')
        run.font.name  = item.get('font', FONT)
        run.font.size  = item.get('size', Pt(20))
        run.font.bold  = item.get('bold', False)
        run.font.color.rgb = item.get('color', C_BLACK)
    return box

def hdr(slide, title):
    rect(slide, Cm(0), Cm(0), Cm(33.87), Cm(2.3), C_DARK)
    txtbox(slide, Cm(0.5), Cm(0.25), Cm(33), Cm(1.85), [
        {'text': title, 'size': Pt(26), 'bold': True, 'color': C_WHITE}
    ])

def add_img(slide, img_path, l, t, w, h=None):
    """插入圖片，若 h=None 則等比例縮放"""
    if not img_path or not os.path.exists(img_path):
        return None
    if h is None:
        return slide.shapes.add_picture(img_path, l, t, width=w)
    return slide.shapes.add_picture(img_path, l, t, width=w, height=h)

def add_caption(slide, text, l, t, w):
    txtbox(slide, l, t, w, Cm(0.7), [
        {'text': text, 'size': Pt(14), 'color': C_GRAY,
         'align': PP_ALIGN.CENTER}
    ])

# ── 條目快捷函數 ──────────────────────────────────────
def b0(text, bold=False, color=C_BLACK, size=Pt(20), sb=7):
    return {'text':text,'blt':'•','indent':0,'bold':bold,'color':color,'size':size,'sb':sb}
def b1(text, bold=False, color=C_GRAY, size=Pt(18), sb=3):
    return {'text':text,'blt':'─','indent':1,'bold':bold,'color':color,'size':size,'sb':sb}
def b2(text, bold=False, color=C_GRAY, size=Pt(17), sb=2):
    return {'text':text,'blt':'·','indent':2,'bold':bold,'color':color,'size':size,'sb':sb}
def h(text, color=C_MID, size=Pt(21), sb=10):
    return {'text':text,'blt':'','indent':0,'bold':True,'color':color,'size':size,'sb':sb}
def t(text='', bold=False, color=C_BLACK, size=Pt(20), sb=0, align=PP_ALIGN.LEFT):
    return {'text':text,'blt':'','indent':0,'bold':bold,'color':color,'size':size,'sb':sb,'align':align}

def mk(title, items):
    s = new_slide(); hdr(s, title)
    txtbox(s, Cm(0.8), Cm(2.6), Cm(32.2), Cm(15.7), items)
    return s

def mk2(title, items, img_path, img_left=Cm(17), img_top=Cm(2.7),
        img_w=Cm(15.5), caption=''):
    """左文字右圖片的版面"""
    s = new_slide(); hdr(s, title)
    txtbox(s, Cm(0.8), Cm(2.6), Cm(15.8), Cm(15.7), items)
    if img_path and os.path.exists(img_path):
        pic = slide_img(s, img_path, img_left, img_top, img_w)
        if caption and pic:
            cap_top = img_top + pic.height + Cm(0.15)
            add_caption(s, caption, img_left, cap_top, img_w)
    return s

def slide_img(slide, img_path, l, t, w):
    if img_path and os.path.exists(img_path):
        return slide.shapes.add_picture(img_path, l, t, width=w)
    return None

# ══════════════════════════════════════════════════════
# 投影片生成
# ══════════════════════════════════════════════════════

# 1 封面
s1 = new_slide()
rect(s1, Cm(0), Cm(0),   Cm(33.87), Cm(19.05), C_DARK)
rect(s1, Cm(0), Cm(5.5), Cm(33.87), Cm(9),     C_MID)
txtbox(s1, Cm(1), Cm(6.0), Cm(32), Cm(3.5), [
    {'text':'從臉部肌電訊號到中文文字','size':Pt(38),'bold':True,
     'color':C_WHITE,'align':PP_ALIGN.CENTER}])
txtbox(s1, Cm(1), Cm(9.5), Cm(32), Cm(2.5), [
    {'text':'AuxCEMGR 模型深度解析','size':Pt(28),'bold':True,
     'color':C_LIGHT,'align':PP_ALIGN.CENTER}])
txtbox(s1, Cm(1), Cm(12.5), Cm(32), Cm(3.5), [
    {'text':'Neural Chinese Silent Speech Recognition with Facial Electromyography',
     'size':Pt(17),'color':C_LIGHT,'align':PP_ALIGN.CENTER},
    {'text':'Xie et al.  |  Speech Communication 171 (2025) 103230',
     'size':Pt(17),'color':C_LIGHT,'align':PP_ALIGN.CENTER,'sb':8},
    {'text':'國防創新研究院 / 北京大學 / 哈工大（深圳）',
     'size':Pt(16),'color':C_LIGHT,'align':PP_ALIGN.CENTER,'sb':6}])

# 2 研究背景(1)
mk('研究背景 (1)：語音辨識的成功與限制', [
    h('自動語音辨識（ASR）已非常成熟'),
    b0('智慧音箱、自動字幕、客服機器人、文字聽打等廣泛應用'),
    b0('2017 年 Transformer 架構讓 ASR 性能大幅提升'),
    t(),
    h('但 ASR 有一個根本限制'),
    b0('必須有「聲音輸入」才能運作'),
    b0('沒有聲音？→ 傳統 ASR 完全失效', color=C_RED, bold=True),
    t(),
    h('→ 無聲語音辨識（Silent Speech Recognition, SSR）因此誕生'),
    b0('目標：不需要聲音，直接從肌肉動作還原語意'),
    b0('本論文使用的訊號：臉部表面肌電圖（Facial sEMG）'),
    t(),
    h('本論文的歷史性地位'),
    b0('全球第一篇中文端對端 EMG-SSR 研究', color=C_RED, bold=True),
    b0('過去研究幾乎全是英語，中文是聲調語言，挑戰更大'),
])

# 3 應用場景
mk('研究背景 (2)：什麼時候需要無聲語音辨識？', [
    h('場景 A：醫療需求'),
    b0('喉切除術（laryngectomy）患者：聲帶移除，無法發聲'),
    b0('語音復健：協助患者學習無聲溝通'),
    t(),
    h('場景 B：噪音 / 保密需求'),
    b0('工廠、戰場等高噪音環境，麥克風無法正常使用'),
    b0('AlterEgo 裝置（MIT, Kapur 2018）：靜默下達指令'),
    t(),
    h('場景 C：多模態語音強化'),
    b0('嘈雜環境下用 EMG 輔助聲學訊號，提升辨識準確率'),
    t(),
    h('現有研究的空缺'),
    b0('幾乎所有 EMG-SSR 研究都聚焦英文', color=C_ORNG),
    b0('中文（普通話）是聲調語言，字元數遠多於英文', color=C_ORNG),
    b0('中文端對端 EMG-SSR → 本論文是全球第一篇！', color=C_RED, bold=True),
])

# 4 sEMG 原理
mk('sEMG 是什麼？（從零開始的外行說明）', [
    h('EMG = 表面肌電圖（Surface Electromyography）'),
    b0('肌肉收縮時，神經細胞發出電訊號（動作電位）給肌纖維'),
    b0('動作電位傳導到皮膚表面 → 微弱電壓變化（單位：微伏特 μV）'),
    b0('「表面」電極貼在皮膚外側即可偵測，不需穿刺 → 非侵入式', color=C_GRN),
    t(),
    h('與語音辨識的關聯'),
    b0('說話 = 大腦控制臉部 / 喉部肌肉 → 產生聲音'),
    b0('默唸時：這些肌肉仍然有微弱的收縮動作（EMG 訊號存在！）',
       color=C_GRN, bold=True),
    b0('結論：即使完全靜默，臉部 EMG 仍帶有「你在說什麼」的資訊'),
    t(),
    h('直觀類比'),
    b0('想像自己在跑步 → 腿部肌肉也有微弱電訊號'),
    b0('默唸「籃球」時，下頜 / 嘴唇的肌肉仍做出對應動作，但幅度較小',
       color=C_ORNG),
])

# 5 臉部肌群 + Fig1（EMG→文字範例）
s5 = new_slide(); hdr(s5, '說話時臉部肌肉的運動分析')
txtbox(s5, Cm(0.8), Cm(2.6), Cm(18), Cm(15.7), [
    h('本論文 8 個 EMG 通道對應肌群'),
    b0('CH1：下頜肌（Jaw）── 控制開閉口，訊號最強'),
    b0('CH2：嘴輪匝肌（Orbicularis Oris）── 控制嘴形'),
    b1('Baseline 最重要通道', color=C_GRAY),
    b0('CH3：臉頰肌群'),
    b0('CH4：下頜骨肌群 ── 最重要！', bold=True, color=C_RED),
    b1('AuxCEMGR 最重要通道', color=C_GRAY),
    b0('CH5、CH6：喉部周圍'),
    b1('有聲時活躍；靜默模式幾乎靜止 → 效果最差', color=C_ORNG),
    b0('CH7、CH8：其他臉部輔助肌群'),
    t(),
    h('靜默 vs. 有聲的差異'),
    b0('整體：靜默 EMG 分佈與有聲相似，但能量幅度較低'),
    b0('CH5/6（喉部）：有聲明顯活躍，靜默幾乎靜止'),
])
if fig1_path:
    slide_img(s5, fig1_path, Cm(19.2), Cm(2.8), Cm(14))
    add_caption(s5, '圖 1：臉部 EMG 訊號 → 中文文字（論文 Fig.1）',
                Cm(19.2), Cm(14.5), Cm(14))

# 6 三大貢獻
mk('本論文的三大貢獻', [
    h('貢獻一：第一個中文 EMG-SSR 基準資料集', color=C_RED),
    b0('建立可公開申請的平行語料庫（EMG 訊號 ↔ 中文字元）'),
    b1('NBA 賽事報導語料，1,238 句 / 667 個獨特字元', color=C_GRAY),
    t(),
    h('貢獻二：第一個中文端對端神經網路 SSR 模型', color=C_RED),
    b0('架構：Transformer 編碼器 + CTC 解碼器'),
    b0('直接 EMG → 中文字元（end-to-end，無需中間步驟）'),
    t(),
    h('貢獻三：創新輔助任務 + 資料增強策略', color=C_RED),
    b0('輔助任務 1：拼音生成（利用漢語音節特性）'),
    b0('輔助任務 2：Session 分類（梯度反轉，消除電極差異干擾）'),
    b0('三種資料增強：頻譜相減、有聲 EMG 輔助、Mixup'),
    t(),
    b0('最終結果：測試集 CER = 38.0%（字元錯誤率，越低越好）',
       bold=True, color=C_MID),
])

# 7 資料集語料
mk('資料集建立 (1)：語料來源與受試者', [
    h('受試者'),
    b0('1 位女性志願者（母語普通話，無閱讀障礙，無語音障礙）'),
    b0('主要局限：僅單一受試者 → 無法反映個體差異', color=C_RED),
    t(),
    h('語料來源：新浪體育 NBA 賽事報導（2012 年起）'),
    b0('爬取工具：Selenium + BeautifulSoup'),
    b0('篩選：≤ 20 字、常見籃球術語、過濾 HTML'),
    b0('特性：主題集中、句子短、用字常見 → 適合初期研究'),
    t(),
    h('語料統計'),
    b0('總計：1,238 句 / 12,584 字 / 667 個獨特字元'),
    b0('平均每句：10 個字元'),
    b0('資料集切分：'),
    b1('訓練集 1,062 句（涵蓋全部 667 個字元）'),
    b1('驗證集 63 句  /  測試集 113 句'),
])

# 8 設備（+ Fig2 電極位置）
s8 = new_slide(); hdr(s8, '資料集建立 (2)：錄製設備與電極配置')
txtbox(s8, Cm(0.8), Cm(2.6), Cm(18), Cm(15.7), [
    h('EMG 設備'),
    b0('型號：Neuracle Technology NSW308M 雙極系統'),
    b0('電極：Ag/AgCl 差分電極'),
    b0('採樣率：1000 Hz，8 通道（16 個差分電極對）'),
    b0('參考電極：1 個，貼於左側鎖骨'),
    b1('差分電極 = 兩電極相減 → 消除共模雜訊', color=C_GRAY),
    t(),
    h('音訊設備（輔助）'),
    b0('麥克風，採樣率：44,100 Hz，單聲道'),
    t(),
    h('電極配置（各通道角色）'),
    b0('CH1：下頜  CH2：嘴輪匝肌  CH3：臉頰'),
    b0('CH4：下頜骨肌群（最重要）', bold=True, color=C_RED),
    b0('CH5、CH6：喉部  CH7、CH8：臉部輔助'),
])
if fig2_path:
    slide_img(s8, fig2_path, Cm(19.2), Cm(2.8), Cm(14))
    add_caption(s8, '圖 2：電極位置示意圖（論文 Fig.2）',
                Cm(19.2), Cm(14.5), Cm(14))

# 9 錄製流程
mk('資料集建立 (3)：錄製流程與 Session 設計', [
    h('兩種錄製模式'),
    b0('靜默模式（Silent）── 主要目標'),
    b1('受試者默唸，不發出聲音；輸出：EMG + 文字標籤'),
    b0('有聲模式（Audible）── 輔助資料'),
    b1('大聲朗讀相同句子；輸出：EMG + 音訊 + 文字標籤'),
    t(),
    h('Session 設計（每句話錄 5 次）'),
    b0('每個 session 之間：撕掉電極貼片，重新貼上再錄'),
    b0('目的：模擬現實中每次配戴位置不完全一樣'),
    b0('副作用：同一句話在不同 session 的 EMG 略有差異', color=C_ORNG),
    b1('→ 這是設計「Session 分類輔助任務」的動機', color=C_GRAY),
    t(),
    h('最終資料量'),
    b0('靜默 EMG：5,468 組（5.93 小時）'),
    b0('有聲 EMG：5,829 組（5.80 小時）'),
    b0('每段訊號平均長度：3.74 秒', color=C_ORNG),
])

# 10 靜默 vs. 有聲 EMG 波形（Fig3）
s10 = new_slide(); hdr(s10, '資料集建立 (4)：靜默 vs. 有聲 EMG 訊號比較')
txtbox(s10, Cm(0.8), Cm(2.6), Cm(16), Cm(15.7), [
    h('靜默 EMG vs. 有聲 EMG 的差異'),
    b0('分佈相似，但能量幅度不同'),
    b0('CH1（下頜）：有聲 ≈ 靜默（兩者都在動）'),
    b0('CH5、CH6（喉部）：',color=C_ORNG),
    b1('有聲時訊號明顯（聲帶振動）'),
    b1('靜默時幾乎為零（喉部不動）'),
    t(),
    h('為什麼這很重要？'),
    b0('→ 說明靜默辨識中喉部電極（CH5/6）貢獻極低'),
    b0('→ 模型需要從嘴部、下頜的動作推斷語音內容'),
    b0('→ 有聲 EMG 可以輔助訓練，但需要折扣係數 γ',
       color=C_ORNG),
])
if fig3_path:
    slide_img(s10, fig3_path, Cm(17.0), Cm(2.8), Cm(16.2))
    add_caption(s10,
        '圖 3：相同句子「最後時刻命中三分」的靜默 vs. 有聲 EMG（論文 Fig.3）',
        Cm(17.0), Cm(15.2), Cm(16.2))

# 11 前處理 濾波
mk('訊號前處理 (1)：濾波去雜訊', [
    h('原始 EMG 含有哪些雜訊？'),
    b0('生理雜訊：心跳（~1 Hz）、基線漂移（呼吸 / 汗水）'),
    b0('電源雜訊：交流電 50 Hz 及其諧波（150、250、350 Hz）'),
    b0('高頻雜訊：EMG 有用範圍以外的電磁干擾'),
    t(),
    h('步驟 1：帶通濾波器（Bandpass Filter）'),
    b0('類型：4 階 Butterworth 濾波器（通帶平坦、過渡帶陡峭）'),
    b0('通頻範圍：10 Hz – 400 Hz'),
    b1('< 10 Hz → 濾除（基線漂移）', color=C_GRAY),
    b1('> 400 Hz → 濾除（高頻雜訊）', color=C_GRAY),
    b1('實作：scipy.signal.butter() + filtfilt()', color=C_GRAY),
    t(),
    h('步驟 2：陷波濾波器（Notch Filter）'),
    b0('移除：50 Hz、150 Hz、250 Hz、350 Hz（電源基頻與諧波）'),
    b1('實作：scipy.signal.iirnotch()，對每個頻率分別套用', color=C_GRAY),
])

# 12 前處理 MFSC
mk('訊號前處理 (2)：MFSC 特徵提取', [
    h('為什麼不直接用原始 EMG 訊號？'),
    b0('維度過高（1000 Hz × 8 通道 × 數秒）'),
    b0('重要的是「頻率隨時間的變化」，而不是逐點數值'),
    t(),
    h('MFSC = log 梅爾頻率譜係數（源自語音辨識）'),
    b0('Step 1：Hanning 滑動窗口切分訊號（帶重疊）'),
    b0('Step 2：短時傅立葉轉換（STFT）→ 頻域能量分佈'),
    b0('Step 3：36 個梅爾濾波器（模仿人耳非線性感知）'),
    b0('Step 4：取對數（log）→ 壓縮動態範圍'),
    b0('結果：每個時間視窗 = 36 維特徵向量'),
    t(),
    h('實作工具'),
    b0('Librosa：librosa.feature.melspectrogram()'),
    b0('Scipy：濾波操作'),
    b0('MFSC 是 MFCC 的前驅版（不做 DCT），保留更多頻率資訊',
       color=C_GRAY, size=Pt(18)),
])

# 13 模型架構總覽（圖表D）
s13 = new_slide(); hdr(s13, '模型架構總覽：AuxCEMGR')
txtbox(s13, Cm(0.8), Cm(2.6), Cm(14.5), Cm(15.7), [
    h('AuxCEMGR = Auxiliary-enhanced Chinese EMG Recognition'),
    t(),
    h('三個組件'),
    b0('CNN（2層）：壓縮序列長度，提取局部特徵'),
    b0('Transformer 編碼器（6層）：全局時序依賴'),
    b0('CTC 解碼器：EMG → 中文字元（主任務）'),
    t(),
    h('兩個輔助任務'),
    b0('拼音生成 CTC：強迫編碼器學音素特徵'),
    b0('Session 分類 + GRL：讓特徵對電極位置不敏感'),
    t(),
    h('合併損失'),
    b0('L_all = L_CTC + η₁·L_pin + η₂·L_ses'),
    b1('η₁ = η₂ = 1.0（驗證集決定）', color=C_GRAY),
])
slide_img(s13, chart_arch, Cm(15.2), Cm(2.7), Cm(18.2))
add_caption(s13, '模型架構示意圖（自繪）', Cm(15.2), Cm(14.3), Cm(18.2))

# 14 CNN 層
mk('編碼器 Part 1：CNN 壓縮層', [
    h('為什麼先用 CNN？'),
    b0('Transformer 計算複雜度 O(T²)，序列長 T 越大越慢'),
    b0('原始 MFSC 序列：3.74 秒 × ~100 幀/秒 ≈ 374 幀'),
    b0('CNN 先壓縮序列長度（N → T），降低 Transformer 計算量'),
    t(),
    h('CNN 設定'),
    b0('層數：2 層，Kernel size 3 × 3（時間軸 × 特徵軸）'),
    b0('每層有 stride → 輸出序列縮短'),
    t(),
    h('CNN 的功能'),
    b0('提取局部時序模式：相鄰幾個時間步間的 EMG 波形變化'),
    b0('類比：圖像辨識的 CNN 學習「邊緣 / 紋理」，'),
    b1('這裡學的是「EMG 波形的局部變化模式」', color=C_GRAY),
    t(),
    h('公式'),
    b0('X = CNN×2(S)  →  輸出 X = x₁…xT（T < N）'),
])

# 15 Transformer（+ 論文原圖 Fig4）
s15 = new_slide(); hdr(s15, '編碼器 Part 2：Transformer 注意力機制')
txtbox(s15, Cm(0.8), Cm(2.6), Cm(17.5), Cm(15.7), [
    h('Transformer 的核心：自注意力（Self-Attention）'),
    b0('"Attention is All You Need"（Vaswani et al., 2017）'),
    b0('讓序列中每個位置都能「看見」所有其他位置'),
    b0('優勢：完全平行計算，能捕捉長距離依賴'),
    t(),
    h('自注意力的直觀理解'),
    b0('處理 EMG 時間步 t 時，同時參考所有其他時間步'),
    b0('學習：「說某個字時，應最關注哪些其他時刻的肌肉狀態」'),
    t(),
    h('本論文設定'),
    b0('6 層 Transformer，8 注意力頭，256 維'),
    b0('每層：Self-Attention → Add & Norm → FFN → Add & Norm'),
    t(),
    h('為什麼 Transformer >> LSTM？（數據）'),
    b0('Baseline（Transformer）：CER 44.5%', color=C_GRN, bold=True),
    b0('Wadkins 2019（LSTM）：CER 66.8%', color=C_RED, bold=True),
    b1('差距 22.3%，Transformer 大幅優勝', color=C_GRAY),
])
if fig4_path:
    slide_img(s15, fig4_path, Cm(18.3), Cm(2.8), Cm(15))
    add_caption(s15, '圖 4：AuxCEMGR 模型架構（論文 Fig.4）',
                Cm(18.3), Cm(14.0), Cm(15))

# 16 CTC 原理
mk('解碼器 (1)：CTC 的核心問題與解法', [
    h('根本問題：序列對齊（Alignment Problem）'),
    b0('EMG 序列有 T 個時間步（例如 374 個）'),
    b0('輸出文字只有 M 個字元（例如 10 個）'),
    b0('T >> M，且無法事先知道哪個時間步對應哪個字', color=C_RED),
    t(),
    h('CTC 的解法：引入「空白符號（─）」'),
    b0('每個時間步 t，從字元表選一個輸出（含空白符號）'),
    b0('解碼規則：合併相鄰重複 → 移除空白'),
    b0('範例：'),
    b1('原始輸出：猛─龍─已─是─東─部─前─五', color=C_GRAY),
    b1('合併重複後移除空白：猛龍已是東部前五 ✓', color=C_GRN),
    t(),
    h('訓練：最大化正確輸出的概率'),
    b0('L_CTC = −log P(Y_正確 | X)'),
    b0('P(Y|X) = 所有合法對齊路徑的概率之和（動態規劃）'),
])

# 17 Beam Search
mk('解碼器 (2)：Beam Search 推理', [
    h('為什麼需要 Beam Search？'),
    b0('暴力搜尋所有路徑：字元數^時間步數 → 不可行'),
    b0('貪婪搜尋（每步選最高機率）：選錯了無法回頭'),
    t(),
    h('Beam Search（束搜尋）'),
    b0('每個時間步保留前 k 個最可能的序列候選'),
    b0('本論文：Beam size = 5（每步保留 5 條路徑）'),
    b0('最後選機率最高的那條路徑作為輸出'),
    b0('類比：同時走 5 條路，最後選最好的', color=C_ORNG),
    t(),
    h('CTC 的局限（未來改進方向）'),
    b0('假設各時間步輸出相互獨立 → 忽略字與字之間的語言關係',
       color=C_ORNG),
    b0('→ 加入語言模型（LM）解碼可進一步降低 CER', color=C_GRN),
])

# 18 輔助任務1 拼音
mk('輔助任務 1：拼音生成（動機與原理）', [
    h('漢語的特殊性：音節 vs. 字元'),
    b0('普通話約 400 個音節（無聲調拼音），對應超過 10,000 個字元'),
    b0('同一音節對應多個字（同音字）'),
    b1('例：「ji」→ 雞、積、機、基、激…', color=C_GRAY),
    b0('EMG 記錄的是「發音動作」，與音節比與字更直接相關'),
    t(),
    h('實作細節'),
    b0('在 Transformer 編碼器輸出 H 上，加一個額外的 CTC 解碼器'),
    b0('預測無聲調拼音序列（例：「guo wang lian xu ming zhong」）'),
    b0('與主任務共用 Transformer 編碼器（參數共享）'),
    b0('→ 強迫編碼器學出更好的音素特徵', color=C_GRN),
    t(),
    h('損失 & 超參數'),
    b0('L_pin = −log P(拼音序列_正確 | X)'),
    b0('損失權重 η₁ = 1.0；η₁ > 0.5 時 CER 明顯改善',
       color=C_RED, bold=True),
])

# 19 輔助任務2 Session（Fig5 heat map）
s19 = new_slide(); hdr(s19, '輔助任務 2：Session 分類 + 梯度反轉層（GRL）')
txtbox(s19, Cm(0.8), Cm(2.6), Cm(18), Cm(15.7), [
    h('問題：電極重貼導致 Session 間差異'),
    b0('每個 session 電極位置略有偏移'),
    b0('模型可能「記住」session 特性，而非學到語音特徵',
       color=C_ORNG),
    t(),
    h('解法：梯度反轉層（GRL）'),
    b0('Session 分類器加在編碼器輸出上'),
    b0('GRL 正向傳播：直接通過'),
    b0('GRL 反向傳播：梯度正負號反轉！', color=C_RED, bold=True),
    b1('→ 分類器努力分類', color=C_GRAY),
    b1('→ 編碼器被迫學「讓分類器分不出來」的特徵', color=C_GRN),
    t(),
    h('損失'),
    b0('L_ses = log P_session(session | X)，η₂ = 1.0'),
])
if fig5_path:
    slide_img(s19, fig5_path, Cm(18.8), Cm(2.8), Cm(14.5))
    add_caption(s19, '圖 5：η₁ & η₂ 對 CER 的影響（論文 Fig.5）',
                Cm(18.8), Cm(14.2), Cm(14.5))

# 20 損失整合
mk('訓練損失函數整合', [
    h('三個損失的組合'),
    b0('主任務（CTC 文字生成）：L_CTC = −log P(Y | X)'),
    b0('輔助任務 1（拼音）：L_pin = −log P(拼音 | X)'),
    b0('輔助任務 2（Session）：L_ses = log P_session(session | X)'),
    t(),
    h('最終合併損失'),
    b0('L_all = L_CTC  +  η₁ × L_pin  +  η₂ × L_ses',
       bold=True, color=C_DARK, size=Pt(22)),
    b0('η₁ = 1.0，η₂ = 1.0（驗證集 grid search 決定）'),
    t(),
    h('訓練策略'),
    b0('三個損失同步 end-to-end 訓練'),
    b0('優化器：Adam + Noam warmup（初始 lr=0.1，warmup 2000 步）'),
    b0('Batch size：128'),
])

# 21 增強1 頻譜相減
mk('資料增強 (1)：頻譜相減（Spectral Subtraction）', [
    h('為什麼需要資料增強？'),
    b0('只有 1 位受試者、約 6 小時靜默 EMG → 訓練資料嚴重不足'),
    b0('EMG 訊號個體差異大，需要更多多樣性'),
    t(),
    h('頻譜相減原理'),
    b0('估算背景雜訊的頻譜，從原始訊號頻譜中減去'),
    b0('得到「降噪版」EMG 訊號'),
    b0('效果：資料集自動翻倍（原版 + 降噪版）', color=C_GRN, bold=True),
    t(),
    h('實驗效果'),
    b0('移除頻譜相減後：'),
    b1('Baseline CER 上升 9.4%（三種增強中最顯著）', color=C_RED),
    b1('AuxCEMGR CER 上升 7.8%', color=C_RED),
    b0('→ 三種增強中貢獻最大的一種', color=C_RED, bold=True),
])

# 22 增強2 有聲（Fig6）
s22 = new_slide(); hdr(s22, '資料增強 (2)：有聲 EMG 輔助訓練')
txtbox(s22, Cm(0.8), Cm(2.6), Cm(18), Cm(15.7), [
    h('概念'),
    b0('有聲模式 EMG 與靜默 EMG 不完全相同，但整體發音模式相似'),
    b0('加入有聲資料 → 補充訓練多樣性', color=C_GRN),
    t(),
    h('實作：加權損失'),
    b0('有聲資料損失乘上折扣係數 γ'),
    b0('最佳 γ：Baseline = 1.0，AuxCEMGR = 0.8'),
    b0('γ 在 0.6–1.0 間效果最佳', color=C_ORNG),
    t(),
    h('實驗效果'),
    b0('移除有聲資料後：Baseline +6.9%，AuxCEMGR +9.1%',
       color=C_RED),
    b0('額外發現：加入有聲訓練後，有聲 EMG 辨識也大幅改善',
       color=C_GRAY, size=Pt(18)),
])
if fig6_path:
    slide_img(s22, fig6_path, Cm(18.8), Cm(2.8), Cm(14.5))
    add_caption(s22, '圖 6：γ 值對 CER 的影響（論文 Fig.6）',
                Cm(18.8), Cm(14.2), Cm(14.5))

# 23 增強3 Mixup
mk('資料增強 (3)：Mixup（混合訓練樣本）', [
    h('核心概念'),
    b0('取兩筆訓練資料線性混合，讓模型學到更平滑的決策邊界'),
    b0('原本用於圖像分類，本論文改造成適用於 CTC 序列任務'),
    t(),
    h('做法'),
    b0('X₁₊₂ = λ × X₁ + (1−λ) × X₂'),
    b0('L ≈ λ × L(X₁₊₂, Y₁) + (1−λ) × L(X₁₊₂, Y₂)'),
    b0('λ ~ Beta(α, α)，α = 0.02（幾乎純樣本，偶爾輕微混合）',
       color=C_ORNG),
    t(),
    h('實驗效果'),
    b0('移除 Mixup：Baseline +5.1%，AuxCEMGR +6.5%', color=C_RED),
    t(),
    h('三種增強的綜合效果（CER 改善量）'),
    b0('頻譜相減 > 有聲資料輔助 > Mixup', bold=True, color=C_RED),
    b0('三種合用比任一單獨使用效果都好（互補）', color=C_GRN),
])

# 24 訓練設定
mk('訓練設定：照著做的完整超參數清單', [
    h('模型架構'),
    b0('CNN：2 層，Kernel 3×3'),
    b0('Transformer：6 層，8 注意力頭，256 維'),
    b0('CTC 主解碼器 + 拼音 CTC + GRL Session 分類器'),
    t(),
    h('訓練超參數'),
    b0('優化器：Adam  /  排程：Noam warmup（lr=0.1，warmup=2000）'),
    b0('Batch size：128  /  CTC Beam size：5'),
    t(),
    h('輔助任務 & 增強超參數'),
    b0('η₁ = 1.0（拼音）  /  η₂ = 1.0（Session）'),
    b0('有聲折扣 γ = 0.8（AuxCEMGR）/ 1.0（Baseline）'),
    b0('Mixup α = 0.02'),
    t(),
    h('硬體環境'),
    b0('6 × NVIDIA RTX 3090 GPU，24 核心 CPU'),
])

# 25 CER 指標
mk('評估指標：字元錯誤率（CER）詳解', [
    h('公式'),
    b0('CER = (替換數 + 刪除數 + 插入數) ÷ 正確答案字元總數 × 100%'),
    b0('用 Levenshtein 距離（動態規劃）計算'),
    t(),
    h('三種錯誤類型'),
    b0('替換：把錯的字改成對的字'),
    b1('例：「國主」→「國王」= 1 次替換', color=C_GRAY),
    b0('刪除：刪掉多預測的字'),
    b1('例：「籃球比賽場」→「籃球比賽」= 1 次刪除', color=C_GRAY),
    b0('插入：插入漏掉的字'),
    b1('例：「球比賽」→「籃球比賽」= 1 次插入', color=C_GRAY),
    t(),
    h('解讀'),
    b0('CER = 38.0% → 每 10 個字平均 ~3.8 個錯', color=C_ORNG),
    b0('實際應用目標：< 5–10%  /  越低越好', color=C_GRAY),
])

# 26 主要結果（圖表A）
s26 = new_slide(); hdr(s26, '主要實驗結果（測試集 CER %，越低越好）')
txtbox(s26, Cm(0.8), Cm(2.6), Cm(14.5), Cm(15.7), [
    h('三大關鍵發現'),
    b0('① 資料增強：平均貢獻 ~23.3% CER 改善',
       bold=True, color=C_RED),
    b0('② Transformer：比 LSTM 優勝 ~22%',
       bold=True, color=C_ORNG),
    b0('③ 雙輔助任務：再降 6.5%',
       bold=True, color=C_MID),
    t(),
    h('最佳模型'),
    b0('AuxCEMGR（完整增強）= 38.0%',
       bold=True, color=C_RED),
    t(),
    h('Baseline 完整增強 = 44.5%'),
    h('LSTM+CTC（Wadkins 2019）= 66.8%',
       color=C_GRAY, sb=4),
])
slide_img(s26, chart_main, Cm(15.2), Cm(2.7), Cm(18.2))
add_caption(s26, '主要實驗結果比較（自繪圖）', Cm(15.2), Cm(15.0), Cm(18.2))

# 27 資料增強效果（圖表B）
s27 = new_slide(); hdr(s27, '資料增強效果分析')
txtbox(s27, Cm(0.8), Cm(2.6), Cm(14.5), Cm(15.7), [
    h('各增強方法的 CER 貢獻（測試集）'),
    t(),
    b0('移除頻譜相減：', bold=True),
    b1('Baseline +9.4%  /  AuxCEMGR +7.8%', color=C_RED),
    t(sb=8),
    b0('移除 Mixup：', bold=True),
    b1('Baseline +5.1%  /  AuxCEMGR +6.5%', color=C_RED),
    t(sb=8),
    b0('移除有聲資料：', bold=True),
    b1('Baseline +6.9%  /  AuxCEMGR +9.1%', color=C_RED),
    t(sb=8),
    b0('三種全部移除：', bold=True),
    b1('Baseline +22.0%  /  AuxCEMGR +24.5%', color=C_RED),
    t(),
    b0('三種增強互補，合用效果最好', color=C_GRN, bold=True),
])
slide_img(s27, chart_aug, Cm(15.2), Cm(3.0), Cm(18.2))
add_caption(s27, '資料增強對 CER 的影響（自繪圖）', Cm(15.2), Cm(14.8), Cm(18.2))

# 28 電極分析（圖表C + Fig8/9）
s28 = new_slide(); hdr(s28, '深入分析 (1)：電極通道貢獻度')
txtbox(s28, Cm(0.8), Cm(2.6), Cm(14.5), Cm(15.7), [
    h('單一通道重要性'),
    b0('AuxCEMGR 最重要：CH4（下頜骨肌）', bold=True, color=C_RED),
    b0('Baseline 最重要：CH2（嘴輪匝肌）'),
    b0('效果最差：CH5、CH6（喉部）', color=C_ORNG),
    b1('靜默時喉部幾乎不動，訊號趨近於零', color=C_GRAY),
    t(),
    h('電極數量 vs. 效能'),
    b0('電極對越多 CER 越低（8 對最佳）'),
    b0('移除 1 對：影響極小'),
    b0('移除超過 2 對：準確率急遽下降', color=C_RED, bold=True),
    b0('建議：6–7 對已足夠', color=C_GRN),
])
slide_img(s28, chart_ch, Cm(15.2), Cm(3.0), Cm(18.2))
add_caption(s28, '各電極通道單獨使用的 CER（近似值，依論文 Fig.8 估讀）',
            Cm(15.2), Cm(13.8), Cm(18.2))

# 29 句長 & Session（Fig10/11）
s29 = new_slide(); hdr(s29, '深入分析 (2)：句長與 Session 穩定性')
txtbox(s29, Cm(0.8), Cm(2.6), Cm(15.5), Cm(15.7), [
    h('句長對 CER 的影響'),
    b0('字數 ≥ 12 時，CER 明顯上升', bold=True, color=C_ORNG),
    b0('AuxCEMGR（Final）在所有字數組別中均最佳'),
    t(),
    h('Session 穩定性'),
    b0('AuxCEMGR（Base）在各 session 間最穩定',
       bold=True, color=C_GRN),
    b1('→ GRL 梯度反轉確實有效消除電極偏差', color=C_GRAY),
    b0('加增強後，兩個模型均更一致'),
])
if fig10_path:
    slide_img(s29, fig10_path, Cm(16.3), Cm(2.8), Cm(10))
    add_caption(s29, '圖 10：句長 vs. CER（論文 Fig.10）',
                Cm(16.3), Cm(12.2), Cm(10))
if fig11_path:
    slide_img(s29, fig11_path, Cm(26.5), Cm(2.8), Cm(7))
    add_caption(s29, '圖 11：Session vs. CER（論文 Fig.11）',
                Cm(26.5), Cm(10.8), Cm(7))

# 30 案例研究（Fig12 + Table6）
s30 = new_slide(); hdr(s30, '深入分析 (3)：案例研究──模型輸出對比')
txtbox(s30, Cm(0.8), Cm(2.6), Cm(18.5), Cm(15.7), [
    h('範例句：「國王此後連續命中三分」'),
    t(),
    b0('Baseline（無增強）：「快船此後不命中三分」',
       color=C_RED),
    b1('× 主詞辨識錯誤（快船 vs. 國王）', color=C_GRAY),
    b0('Baseline（完整增強）：「國王此後連度命中三分」',
       color=C_ORNG),
    b1('✓ 主詞正確；「續」→「度」（發音部位相近）', color=C_GRAY),
    b0('AuxCEMGR（無增強）：「國王後頻繁投命中三分」',
       color=C_ORNG),
    b1('✓ 主詞正確；仍有省略與錯字', color=C_GRAY),
    b0('AuxCEMGR（完整增強）：「國王此後連續命中三分」',
       bold=True, color=C_GRN),
    b1('✓ 完全正確！', color=C_GRN),
    t(),
    h('三個觀察'),
    b0('「分」「命」特別容易辨識 → 肌肉動作幅度較大'),
    b0('AuxCEMGR 輸出更合乎語法 → Session 無關特徵更語意合理'),
    b0('資料增強讓固定搭配詞被完整學習'),
])
if fig12_path:
    slide_img(s30, fig12_path, Cm(19.5), Cm(5.0), Cm(13.8))
    add_caption(s30, '圖 12：有聲 vs. 靜默 EMG 效能（論文 Fig.12）',
                Cm(19.5), Cm(13.5), Cm(13.8))

# 31 如何複製
mk('如何複製實驗：Step-by-Step', [
    h('Step 1：取得程式碼與資料'),
    b0('程式碼：git clone https://github.com/bluishwhite/EMG_ASR'),
    b0('資料：Email 申請（yinerwei1985@gmail.com）'),
    b1('包含所有 EMG 訊號（silent + audible）+ 文字標籤', color=C_GRAY),
    t(),
    h('Step 2：環境安裝'),
    b0('pip install librosa scipy numpy torch'),
    t(),
    h('Step 3：前處理'),
    b0('Butterworth 帶通（10–400 Hz）：scipy.signal.butter() + filtfilt()'),
    b0('陷波（50/150/250/350 Hz）：scipy.signal.iirnotch()'),
    b0('MFSC 提取（36 Mel filters）：librosa.feature.melspectrogram()'),
    t(),
    h('Step 4：訓練 & 評估'),
    b0('框架：PyTorch，按投影片 24 的超參數設定'),
    b0('建議順序：Baseline → 加增強 → 加輔助任務，逐步驗證'),
    b0('目標：Baseline ≈ 44.5%，AuxCEMGR ≈ 38.0%（測試集 CER）'),
])

# 32 限制
mk('論文的四大局限（= 未來研究機會）', [
    h('局限 1：單一受試者（最關鍵）', color=C_RED),
    b0('只有 1 位女性志願者，EMG 訊號因人而異'),
    b0('→ 模型無法直接用於其他人（Cross-subject 泛化問題）', color=C_RED),
    t(),
    h('局限 2：封閉詞彙（OOV 問題）', color=C_RED),
    b0('只有 667 個字，出現訓練集以外的字 → 完全無法處理'),
    t(),
    h('局限 3：封閉語料庫（NBA 主題）', color=C_ORNG),
    b0('語料全是 NBA 報導，無法反映日常語言多樣性'),
    t(),
    h('局限 4：CER 仍偏高（38%）', color=C_ORNG),
    b0('每 10 字平均 ~3.8 個錯；實際應用通常需 < 5–10%'),
    b0('未來改進：加語言模型、更多受試者、更大語料庫', color=C_GRN),
])

# 33 對你的意義
mk('對你研究的意義與改進方向', [
    h('直接比較基準'),
    b0('同樣是：臉部 sEMG + Transformer + 中文 SSR'),
    b0('你的目標：找到改善 38% CER 的方法'),
    t(),
    h('你的差異化機會'),
    b0('多受試者訓練 → 解決跨人泛化（最重要！）',
       color=C_GRN, bold=True),
    b0('聲調辨識輔助任務 → 利用普通話 4 聲（本文未處理）',
       color=C_GRN, bold=True),
    b0('語言模型整合解碼 → 改善 CTC 語意合理性',
       color=C_GRN, bold=True),
    b0('使用 CSL-EMG 資料集 → 資料集差異影響的比較',
       color=C_GRN, bold=True),
    t(),
    h('可直接參考的實作細節'),
    b0('電極重點：CH4（下頜骨肌）→ MyoWare 感測器放置參考'),
    b0('前處理標準流程：10–400 Hz + 50/150/250/350 Hz 陷波'),
    b0('超參數起始點：η₁=η₂=1.0，α=0.02，γ=0.8'),
])

# 34 總結
mk('總結', [
    h('本論文做了什麼'),
    b0('第一個中文 EMG-SSR 基準資料集（1,238 句，8 通道）'),
    b0('AuxCEMGR：Transformer + CTC + 拼音 + GRL Session 分類'),
    b0('三種資料增強：頻譜相減 + 有聲 EMG + Mixup'),
    b0('最佳 CER：38.0%（中文 SSR 第一條基準線）'),
    t(),
    h('技術貢獻排名'),
    b0('① 資料增強 ~23.3%  ② Transformer ~22%  ③ 輔助任務 6.5%',
       bold=True, color=C_RED),
    t(),
    h('開放資源'),
    b0('程式碼：github.com/bluishwhite/EMG_ASR'),
    b0('資料集：yinerwei1985@gmail.com（Email 申請）'),
    t(),
    h('一句話摘要'),
    b0('臉部 sEMG + Transformer + CTC + 拼音 / Session 輔助任務'
       ' → CER 38%，開啟中文無聲語音辨識研究之門',
       bold=True, color=C_DARK),
])

# ── 儲存 ────────────────────────────────────────────
prs.save(OUT_PATH)
print(f'\n✓ Saved: {OUT_PATH}')
print(f'  Total slides: {len(prs.slides)}')
