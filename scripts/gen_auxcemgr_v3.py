#!/usr/bin/env python3
"""
AuxCEMGR 論文解析 PPT — Python-pptx 版（40 頁，含目錄、頁碼、精準裁圖）
"""
import os
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIG_DIR  = '/Users/rayopenclaw/ray-agent/papers/auxcemgr_figs'
OUT_PATH = '/Users/rayopenclaw/ray-agent/papers/AuxCEMGR_v3.pptx'

def fig(name): return f'{FIG_DIR}/{name}.png'

# ── 色彩常數 ─────────────────────────────────────────────────────
C_DARK  = RGBColor(0x1F, 0x38, 0x64)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_ORNG  = RGBColor(0xC5, 0x60, 0x00)
C_GRN   = RGBColor(0x00, 0x6B, 0x00)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LIGHT = RGBColor(0xD6, 0xE4, 0xF7)

# ── 簡報設定 ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(33.87)   # 16:9 寬螢幕
prs.slide_height = Cm(19.05)

W = prs.slide_width
H = prs.slide_height
TOTAL = 40

FONT_TC = '標楷體'
FONT_EN = 'Times New Roman'

# ── 輔助函數 ─────────────────────────────────────────────────────
def blank():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def run_fmt(run, bold=False, size=None, color=None, font=None):
    run.font.bold = bold
    if size:  run.font.size = Pt(size)
    if color: run.font.color.rgb = color
    f = font or FONT_TC
    run.font.name = f
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run.font._element
        for tag in [qn('a:latin'), qn('a:ea')]:
            el = rPr.find(tag)
            if el is None:
                el = etree.SubElement(rPr, tag)
            el.set('typeface', f)
    except Exception:
        pass

def set_font(tf, bold=False, size=20, color=C_BLACK, font=FONT_TC):
    for para in tf.paragraphs:
        for run in para.runs:
            run_fmt(run, bold=bold, size=size, color=color, font=font)

def add_text_box(slide, text, l, t, w, h,
                 bold=False, size=20, color=C_BLACK,
                 align=PP_ALIGN.LEFT, font=FONT_TC,
                 word_wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run_fmt(run, bold=bold, size=size, color=color, font=font)
    return txb

def add_header(slide, title, slide_no):
    """深藍標題列 + 頁碼"""
    # 標題背景
    bar = slide.shapes.add_shape(
        1, Cm(0), Cm(0), W, Cm(2.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_DARK
    bar.line.fill.background()
    # 標題文字
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = title
    run_fmt(run, bold=True, size=24, color=C_WHITE, font=FONT_TC)
    # 頁碼（右下角）
    pn = slide.shapes.add_textbox(W - Cm(3.5), H - Cm(1.0), Cm(3.2), Cm(0.8))
    tf2 = pn.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f'{slide_no} / {TOTAL}'
    run_fmt(r2, size=14, color=C_GRAY, font=FONT_EN)

def add_image(slide, path, l, t, w, h):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)

def bullets(slide, items, l, t, w, h, size=20):
    """多層次條列（items: list of (indent, text)）"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for indent, text in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.level = indent
        p.space_before = Pt(4)
        run = p.add_run(); run.text = text
        fsz = size - indent * 2
        run_fmt(run, size=max(fsz, 14), color=C_BLACK, font=FONT_TC)

# ── 左文右圖 layout helper ────────────────────────────────────────
TL = Cm(0.6)   # 文字左邊
TT = Cm(2.3)   # 文字頂部
TW = Cm(16.5)  # 文字寬（左半）
TH = Cm(15.8)  # 文字高
IL = Cm(17.5)  # 圖片左邊
IT = Cm(2.0)   # 圖片頂部
IW = Cm(15.8)  # 圖片寬
IH = Cm(16.0)  # 圖片高

# ── 建立投影片 ────────────────────────────────────────────────────
n = 0  # 頁碼計數

# ══ 1. 封面 ══════════════════════════════════════════════════════
n += 1; sl = blank()
bg = sl.shapes.add_shape(1, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = C_DARK; bg.line.fill.background()
add_text_box(sl, '用臉部肌電圖做中文無聲語音辨識', Cm(2), Cm(5.5), W-Cm(4), Cm(3),
             bold=True, size=36, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, 'Neural Chinese Silent Speech Recognition with Facial EMG\nXie et al., 2025 ｜ Speech Communication',
             Cm(2), Cm(9.0), W-Cm(4), Cm(2),
             size=20, color=C_LIGHT, align=PP_ALIGN.CENTER, font=FONT_EN)
add_text_box(sl, '整理：陳睿莆　國立虎尾科技大學', Cm(2), Cm(16.5), W-Cm(4), Cm(1.2),
             size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
# 頁碼
pn = sl.shapes.add_textbox(W-Cm(3.5), H-Cm(1.0), Cm(3.2), Cm(0.8))
tf2 = pn.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
r2 = p2.add_run(); r2.text = f'{n} / {TOTAL}'
run_fmt(r2, size=14, color=C_LIGHT, font=FONT_EN)

# ══ 2. 目錄 ══════════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '目錄', n)
toc = [
    (0, '一、研究背景與動機　（頁 3–5）'),
    (0, '二、資料集與硬體設備　（頁 6–8）'),
    (0, '三、訊號前處理　（頁 9–11）'),
    (0, '四、特徵萃取 MFSC　（頁 12–14）'),
    (0, '五、資料增強　（頁 15–18）'),
    (0, '六、AuxCEMGR 模型架構　（頁 19–26）'),
    (0, '七、訓練設定　（頁 27–29）'),
    (0, '八、評估指標 CER　（頁 30）'),
    (0, '九、實驗結果分析　（頁 31–37）'),
    (0, '十、結論與未來方向　（頁 38–39）'),
    (0, '十一、實驗復現快速指南　（頁 40）'),
]
bullets(sl, toc, Cm(1.5), Cm(2.5), W-Cm(3), Cm(16.0), size=22)

# ══ 3. 什麼是無聲語音介面 ════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '一、研究背景 | 什麼是無聲語音介面（SSI）？', n)
bullets(sl, [
    (0, '定義：不發出聲音，靠感測器辨識說話意圖 → 轉成文字'),
    (1, '關鍵詞：Silent Speech Interface（SSI）'),
    (0, '應用場景'),
    (1, '喉嚨受傷者、氣管切開術後無法發聲的病人'),
    (1, '高噪音環境（工廠、戰場）'),
    (1, '隱密通訊（不希望他人聽到內容）'),
    (0, '為什麼難？'),
    (1, '靜默時 EMG 訊號比有聲弱 3–5 倍，訊雜比（SNR）低'),
    (1, '中文字元龐大（本研究 667 個唯一字元）'),
    (1, '每次貼電極位置略不同 → 訊號分布偏移（Session Drift）'),
], TL, TT, TW, TH)
add_image(sl, fig('fig1'), IL, IT, IW, IH)

# ══ 4. 為什麼用臉部 sEMG ══════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '一、研究背景 | 為什麼用臉部 sEMG？', n)
bullets(sl, [
    (0, 'sEMG（表面肌電圖）'),
    (1, '貼在皮膚表面的電極，偵測肌肉收縮時的微弱電訊號（0.1–5 mV）'),
    (1, '非侵入式（不需手術）、即時、可穿戴'),
    (0, '為什麼臉部？'),
    (1, '說話時嘴唇、下巴、臉頰、喉部肌肉都會動'),
    (1, '靜默說話時，這些肌肉仍有微小收縮，可被偵測'),
    (0, '與其他 SSI 技術的比較'),
    (1, '喉部麥克風：需靠聲帶振動，完全靜默時無效'),
    (1, '腦電圖（EEG）：解析度低、個體差異大、不穩定'),
    (1, '臉部 sEMG：輕便、準確、可長時間佩戴'),
], TL, TT, TW, TH)
add_image(sl, fig('fig2'), IL, IT, IW, IH)

# ══ 5. 研究動機與挑戰 ══════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '一、研究背景 | 研究動機與挑戰', n)
bullets(sl, [
    (0, '現有研究的空白'),
    (1, '現有 sEMG 語音辨識大多針對英文'),
    (1, '中文 SSI 研究幾乎空白，本文為首篇'),
    (0, '中文的特殊挑戰'),
    (1, '字元空間龐大：667 個唯一字元（英文 26 字母）'),
    (1, '聲調系統複雜：4 個聲調 + 輕聲'),
    (1, '同音字多：「ji」可以是「雞」、「機」、「積」…'),
    (0, 'Session 偏移問題（Domain Shift）'),
    (1, '每次錄音重新貼電極，位置略不同'),
    (1, '導致同一個字的 EMG 特徵在不同 Session 分布不一樣'),
    (1, '直接跨 Session 測試：準確率大幅下降'),
    (0, '本研究目標：設計能克服以上挑戰的第一個中文臉部 EMG SSI 系統'),
], TL, TT, TW, TH)

# ══ 6. 資料集介紹 ══════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '二、資料集與硬體 | 資料集介紹', n)
bullets(sl, [
    (0, '受試者：1 名成年女性志願者'),
    (0, '語料庫：NBA 籃球播報語料（涵蓋日常用語與專有名詞）'),
    (0, '規模'),
    (1, '1,238 個句子'),
    (1, '12,584 個總字元'),
    (1, '667 個唯一字元'),
    (0, '資料分割'),
    (1, '訓練：驗證：測試 = 7 : 1 : 2'),
    (0, 'Session 設計'),
    (1, '共 5 個 Session（不同日期，每次重新佩戴電極）'),
    (1, '每 Session 錄 1,238 句 → 共 6,190 次錄音'),
    (1, '訓練 Session 1-3，測試 Session 4-5（跨 Session）'),
], TL, TT, TW, TH)

# ══ 7. 硬體設備 ════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '二、資料集與硬體 | 硬體設備 Neuracle NSW308M', n)
bullets(sl, [
    (0, '放大器型號：Neuracle NSW308M（醫療級生物訊號放大器）'),
    (0, '採樣率：1000 Hz（每秒採集 1000 個數據點）'),
    (0, '通道數：8 通道（同時錄 8 個電極的訊號）'),
    (0, '電極材質：Ag/AgCl（銀/氯化銀）'),
    (1, '導電性佳、生物相容性高、電極-皮膚阻抗低'),
    (0, '訊號特性'),
    (1, '原始訊號振幅：0.1–5 mV（極微弱）'),
    (1, '需放大 1000–10000 倍才能進行數位處理'),
    (1, '採用電池供電 + 藍牙傳輸，減少 60 Hz 電源雜訊'),
], TL, TT, TW, TH)

# ══ 8. 電極位置 ════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '二、資料集與硬體 | 8 個電極的擺放位置', n)
bullets(sl, [
    (0, 'CH1–CH2：上嘴唇兩側（口輪匝肌上半）'),
    (1, '控制上唇開合，發「ㄚ」、「ㄛ」時明顯收縮'),
    (0, 'CH3–CH4：下巴左右（下頷肌 / 頦肌）'),
    (1, 'CH4 實驗證實最重要的通道（CER 最低）'),
    (0, 'CH5–CH6：喉部（環甲肌 / 甲狀舌骨肌）'),
    (1, '靜默說話時幾乎不動 → 貢獻最小的通道'),
    (0, 'CH7–CH8：臉頰（顴大肌）'),
    (1, '發「ㄧ」等前元音時收縮'),
    (0, '參考電極：耳後乳突骨（訊號基準點，消除共模雜訊）'),
], TL, TT, TW, TH)
add_image(sl, fig('fig2'), IL, IT, IW, IH)

# ══ 9. 訊號前處理 — 為什麼 ═══════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '三、訊號前處理 | 為什麼需要濾波？', n)
bullets(sl, [
    (0, '原始 EMG 訊號充滿雜訊'),
    (1, '工頻電磁干擾（50 Hz 及諧波 150/250/350 Hz）'),
    (1, '電極接觸不穩定的基線漂移（< 10 Hz）'),
    (1, '高頻肌肉噪訊（> 400 Hz，非說話相關）'),
    (0, '如果不濾波的後果'),
    (1, '模型學到雜訊特徵而非說話特徵 → 準確率大幅下降'),
    (1, '訓練集和測試集雜訊不同 → 泛化能力差'),
    (0, '本研究的前處理流程'),
    (1, 'Step 1：Butterworth 帶通濾波（保留 10–400 Hz）'),
    (1, 'Step 2：Notch 陷波濾波（去除 50/150/250/350 Hz）'),
    (1, 'Step 3：分段（切成句子長度的片段）'),
], TL, TT, TW, TH)

# ══ 10. Butterworth 濾波器 ═══════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '三、訊號前處理 | Butterworth 帶通濾波器', n)
bullets(sl, [
    (0, '名稱由來：英國工程師 Stephen Butterworth（1930 年）'),
    (0, '作用：只讓 10–400 Hz 的訊號通過，其餘截除'),
    (1, '10 Hz 以下：慢漂移（電極膠乾掉、頭部移動造成）'),
    (1, '400 Hz 以上：高頻電氣雜訊，非 EMG 有效成分'),
    (0, '什麼是「4 階」？'),
    (1, '截止斜率：階數越高，截止越銳利（每倍頻衰減 -80 dB）'),
    (1, '4 階是 EMG 處理的常用選擇（效果好、運算量合理）'),
    (0, '直覺比喻：像是只開中間那道閘門，太高太低的水都不讓流過'),
    (0, '程式碼：'),
    (1, 'from scipy.signal import butter, filtfilt'),
    (1, 'b, a = butter(4, [10, 400], btype="band", fs=1000)'),
    (1, 'filtered = filtfilt(b, a, raw_signal)'),
], TL, TT, TW, TH)

# ══ 11. Notch 濾波器 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '三、訊號前處理 | Notch（陷波）濾波器', n)
bullets(sl, [
    (0, '作用：精準去除特定頻率（像是在頻譜上挖洞）'),
    (0, '去除頻率：50, 150, 250, 350 Hz'),
    (1, '50 Hz：台灣市電工頻干擾（電源線輻射電磁場）'),
    (1, '150, 250, 350 Hz：50 Hz 的 2/3/4 次諧波'),
    (1, '諧波：非線性電路產生的整數倍頻率干擾，不去除會殘留明顯峰值'),
    (0, '品質因子 Q = 30 的意義'),
    (1, 'Q 越高 → 陷波越窄 → 只精準去除目標頻率，不影響旁邊訊號'),
    (1, '本研究 Q=30：每個陷波帶寬約 ±0.83 Hz，足夠精準'),
    (0, '程式碼：'),
    (1, 'from scipy.signal import iirnotch'),
    (1, 'for f0 in [50, 150, 250, 350]:'),
    (1, '    b, a = iirnotch(f0, Q=30, fs=1000)'),
    (1, '    signal = filtfilt(b, a, signal)'),
], TL, TT, TW, TH)

# ══ 12. 什麼是 MFSC ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '四、特徵萃取 | 什麼是 MFSC？', n)
bullets(sl, [
    (0, 'MFSC = Mel Frequency Spectral Coefficients（Mel 頻率濾波器組係數）'),
    (0, '為什麼不直接用原始訊號？'),
    (1, '原始訊號：每秒 1000 個點 × 8 通道 = 8000 維/秒，維度太高'),
    (1, '大量點是冗餘的：相鄰時刻的訊號幾乎一樣'),
    (0, '什麼是 Mel 尺度？'),
    (1, '人耳對低頻分辨力強、對高頻分辨力弱'),
    (1, 'Mel 尺度把頻率壓縮，低頻區解析度高，高頻區解析度低'),
    (1, '對應人耳感知，減少無效資訊'),
    (0, '本研究設定'),
    (1, '36 個 Mel 濾波器（覆蓋 0–400 Hz）'),
    (1, '每 10 ms 計算一次 → 每個時步輸出 8×36 維特徵'),
    (1, '輸入維度：從 8000 維/秒 → 壓縮成 2880 維/秒，且更有意義'),
], TL, TT, TW, TH)

# ══ 13. MFSC 計算步驟 ═══════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '四、特徵萃取 | MFSC 計算步驟（逐步說明）', n)
bullets(sl, [
    (0, 'Step 1：分窗（Framing）'),
    (1, '用 Hanning 窗把訊號切成小段'),
    (1, '每段 25 ms（25 點），每次移動 10 ms（10 點，重疊 60%）'),
    (1, 'Hanning 窗作用：避免「截斷效應」，讓邊緣平滑過渡'),
    (0, 'Step 2：快速傅立葉變換（FFT）'),
    (1, '把每段時域訊號轉換成頻域 → 看這段訊號有哪些頻率成分'),
    (1, '輸出：功率頻譜（Power Spectrum）'),
    (0, 'Step 3：套用 Mel 濾波器組'),
    (1, '在 Mel 尺度均勻放 36 個三角形濾波器'),
    (1, '計算每個濾波器覆蓋範圍內的總能量'),
    (0, 'Step 4：取對數（Log）'),
    (1, '模仿人耳對音量的對數感知（dB 單位）'),
    (1, '也讓數值範圍更適合神經網路訓練'),
    (0, '結果：8 通道 × T 時步 × 36 維 的特徵矩陣'),
], TL, TT, TW, TH)

# ══ 14. 靜默 vs 有聲 EMG ═════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '四、特徵萃取 | 靜默 vs 有聲 EMG 波形差異', n)
bullets(sl, [
    (0, '觀察論文 Fig.3 的關鍵發現'),
    (1, '有聲說話 EMG：振幅大，波形明顯，特徵清晰'),
    (1, '靜默說話 EMG：振幅約為有聲的 1/3–1/5，波形微弱'),
    (0, '為什麼靜默更難辨識？'),
    (1, '訊雜比（SNR）低：有用訊號和雜訊的比例惡化'),
    (1, '特徵不明顯：模型需要更強的表徵學習能力'),
    (0, '這個觀察直接啟發了資料增強②'),
    (1, '如果把有聲 EMG 按比例混入靜默 EMG，可以增強訊號'),
    (1, '讓模型學到更穩健的特徵（詳見下一章節）'),
], TL, TT, TW, TH)
add_image(sl, fig('fig3'), IL, IT, IW, IH)

# ══ 15. 資料增強 — 為什麼 ════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '五、資料增強 | 為什麼需要資料增強？', n)
bullets(sl, [
    (0, '問題①：資料量少，容易過擬合'),
    (1, '只有 1 名受試者，1,238 句訓練資料'),
    (1, '模型看到的樣本有限 → 記住訓練資料而非學到規律'),
    (0, '問題②：Session 間訊號分布偏移'),
    (1, '每次重貼電極，位置略不同 → Domain Shift'),
    (1, '在一個 Session 訓練的模型，直接用在另一個 Session 效果差'),
    (0, '問題③：靜默 EMG 訊號微弱，噪訊比例高'),
    (1, '單純靠原始訊號，模型不易學到穩健特徵'),
    (0, '解決策略：三種資料增強組合'),
    (1, '① 頻譜相減（Spectral Subtraction）→ 降噪'),
    (1, '② 有聲 EMG 混合（Audible EMG Mixing）→ 增強訊號'),
    (1, '③ Mixup（樣本插值）→ 提升泛化能力'),
], TL, TT, TW, TH)

# ══ 16. 頻譜相減 ════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '五、資料增強① | 頻譜相減（Spectral Subtraction）', n)
bullets(sl, [
    (0, '原理：估計背景雜訊的頻譜，從訊號中減去'),
    (0, '步驟'),
    (1, 'Step 1：在受試者靜止不說話時，錄製背景 EMG（約 5 秒）'),
    (1, 'Step 2：計算背景 EMG 的平均功率頻譜（N̂(f)）'),
    (1, 'Step 3：對說話段的頻譜，每個頻率減去估計噪訊'),
    (1, 'Step 4：若結果為負，設為 0（Half-wave Rectification）'),
    (0, '公式：Ŝ(f) = max( X(f) − N̂(f) , 0 )'),
    (1, 'X(f)：含噪訊的原始頻譜'),
    (1, 'N̂(f)：估計背景噪訊頻譜'),
    (1, 'Ŝ(f)：去噪後的乾淨訊號頻譜'),
    (0, '效果：Baseline CER 從 66.5% → 53.9%（改善 12.6 個百分點）'),
    (1, '所有增強方法中，頻譜相減的貢獻最大'),
], TL, TT, TW, TH)

# ══ 17. 有聲 EMG 混合 ═══════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '五、資料增強② | 有聲 EMG 混合（Audible EMG Mixing）', n)
bullets(sl, [
    (0, '動機：有聲說話的 EMG 訊號較強，可幫助訓練'),
    (0, '方法：把有聲說話的 EMG 和靜默 EMG 加權混合'),
    (0, '公式：x_aug = (1−γ)·x_silent + γ·x_audible'),
    (1, 'γ：混合比例，本研究 γ = 0.8'),
    (1, 'γ=0：純靜默，γ=1：純有聲'),
    (1, 'γ=0.8 表示 80% 有聲成分 + 20% 靜默成分'),
    (0, '為什麼 γ=0.8 最好？（見論文 Fig.6）'),
    (1, 'γ 太小：增強效果不足，靜默特徵仍微弱'),
    (1, 'γ 太大（=1.0）：訓練和測試分布差太多，泛化差'),
    (1, 'γ=0.8：甜蜜點，讓訓練資料更豐富，同時不失真'),
    (0, '效果：AuxCEMGR 加入此增強後 CER 再降 6.5 個百分點'),
], TL, TT, TW, TH)
add_image(sl, fig('fig6'), IL, IT, IW, IH)

# ══ 18. Mixup ═══════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '五、資料增強③ | Mixup（樣本插值）', n)
bullets(sl, [
    (0, '原理：隨機取兩個訓練樣本，按比例混合特徵和標籤'),
    (0, '公式'),
    (1, 'x̃ = λ·xᵢ + (1−λ)·xⱼ　（特徵混合）'),
    (1, 'ỹ = λ·yᵢ + (1−λ)·yⱼ　（標籤混合，CTC 中為兩個序列）'),
    (0, 'λ 從 Beta(α, α) 分布取樣，本研究 α = 0.02'),
    (1, 'Beta 分布：連續機率分布，值域 [0, 1]'),
    (1, 'α=0.02 → Beta(0.02,0.02) 極度偏向 0 或 1'),
    (1, '效果：大多數時候幾乎只用一個樣本（λ≈0 或 λ≈1），偶爾才真正混合'),
    (1, '為什麼？這樣可以在不破壞原始樣本的前提下，引入少量擾動'),
    (0, '訓練效果'),
    (1, '使模型對輸入的微小變化更穩健（Regularization 效果）'),
    (1, '減少對特定訓練樣本的過度記憶'),
    (0, '效果：三種增強合用，AuxCEMGR CER 從 44.5% → 38.0%'),
], TL, TT, TW, TH)
add_image(sl, fig('fig7'), IL, IT, IW, IH)

# ══ 19. 模型架構總覽 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | 模型總覽', n)
bullets(sl, [
    (0, '全名：Auxiliary-task Cascaded EMG Recognizer'),
    (0, '四大模組（由左至右串聯）'),
    (1, '① CNN 特徵萃取層 — 提取局部肌肉動作模式'),
    (1, '② Transformer 編碼器（6層）— 捕捉全局時序依賴'),
    (1, '③ CTC 解碼器（主任務）— 輸出中文字元序列'),
    (0, '兩個輔助任務（並行分支）'),
    (1, '④ 拼音 CTC — 輔助任務一，幫助學音韻結構'),
    (1, '⑤ GRL + Session 分類 — 輔助任務二，消除 Session 偏移'),
    (0, '訓練思路：多任務學習'),
    (1, '主任務（字元辨識）+ 輔助任務 → 共同優化'),
    (1, '輔助任務提供額外監督訊號 → 特徵學得更通用'),
], TL, TT, TW, TH)
add_image(sl, fig('chart_arch'), IL, IT, IW, IH)

# ══ 20. CNN 特徵萃取 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | ① CNN 特徵萃取層', n)
bullets(sl, [
    (0, '輸入：MFSC 特徵，形狀 (8通道, 36維, T時步)'),
    (0, 'CNN 架構'),
    (1, '2 層 2D 卷積（Conv2D）'),
    (1, '卷積核大小：3×3，步長（stride）= 2'),
    (1, '每層同時在時間軸和頻率軸上尋找局部模式'),
    (0, '步長=2 的效果'),
    (1, '每層將時序長度壓縮一半（T → T/2 → T/4）'),
    (1, '同時提取更抽象的特徵（從低階頻率到高階動作模式）'),
    (0, '輸出'),
    (1, '維度：256（隱藏維度）'),
    (1, '長度：T/4（時序壓縮 4 倍）'),
    (0, '為什麼用 CNN 而非直接用 Transformer？'),
    (1, 'CNN 擅長捕捉局部時空模式（短時間的肌肉動作）'),
    (1, 'Transformer 後再處理 CNN 輸出 → 效率更高、效果更好'),
], TL, TT, TW, TH)
add_image(sl, fig('fig4'), IL, IT, IW, IH)

# ══ 21. Transformer 編碼器 ══════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | ② Transformer 編碼器', n)
bullets(sl, [
    (0, 'Transformer 是什麼？'),
    (1, '2017 年 Google「Attention is All You Need」論文提出'),
    (1, '核心機制：自注意力（Self-Attention）— 讓每個時步都能「看」到整個序列'),
    (0, '本研究規格：6 層、8 個注意力頭、隱藏維度 256'),
    (0, '什麼是「注意力頭（Head）」？'),
    (1, '把同一段序列從 8 個不同角度（子空間）同時分析'),
    (1, '例如：頭1 看短距依賴，頭2 看長距依賴，頭3 看特定頻率模式…'),
    (1, '最後把 8 個角度的結果拼接合併'),
    (0, '位置編碼（Positional Encoding）'),
    (1, 'Transformer 本身不知道時間順序'),
    (1, '需額外加入正弦/餘弦位置編碼，告訴模型「這是第幾個時步」'),
    (0, 'Transformer vs RNN 的差異'),
    (1, 'RNN：序列處理，慢、難捕捉長距依賴'),
    (1, 'Transformer：並行計算，快、能捕捉任意長度的依賴'),
], TL, TT, TW, TH)

# ══ 22. 什麼是 CTC ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | 什麼是 CTC？為什麼需要？', n)
bullets(sl, [
    (0, 'CTC = Connectionist Temporal Classification（連結時序分類）'),
    (0, '核心問題：輸入輸出長度不一樣'),
    (1, 'EMG 序列長度：約 100 個時步'),
    (1, '文字序列長度：約 20 個字'),
    (1, '而且我們不知道每個字對應哪幾個時步（沒有對齊標註）'),
    (0, 'CTC 的解法：引入「空白符號 ε」'),
    (1, '讓模型在每個時步輸出：某個字元 或 ε（無輸出）'),
    (1, '折疊規則：合併連續重複字元，移除所有 ε'),
    (0, '範例'),
    (1, '模型輸出：你-你-好-ε-ε-世-界'),
    (1, '折疊後：你好世界 ✓'),
    (0, 'CTC 損失函數'),
    (1, '最大化正確字元序列出現的所有可能對齊路徑的總概率'),
    (1, '用前向-後向演算法（Forward-Backward Algorithm）高效計算'),
], TL, TT, TW, TH)

# ══ 23. Beam Search ═════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | ③ Beam Search 解碼（size=5）', n)
bullets(sl, [
    (0, '解碼問題：CTC 輸出每個時步的字元概率，如何找最佳句子？'),
    (0, '方法一：Greedy 解碼（貪婪）'),
    (1, '每步選最高概率的字元 → 快，但不一定全局最優'),
    (1, '容易陷入局部最優（一步走錯，後面都錯）'),
    (0, '方法二：Beam Search（集束搜索）'),
    (1, '同時維護「最可能的 k 條路徑」（k=Beam Size=5）'),
    (1, '每步擴展這 5 條路徑，再保留最好的 5 條'),
    (1, '最終選擇總概率最高的路徑'),
    (0, 'Beam Size = 5 的意涵'),
    (1, 'Beam=1：等同 Greedy，最快最差'),
    (1, 'Beam=5：在準確率和計算量之間取得平衡'),
    (1, 'Beam=50：更準但慢很多，邊際效益遞減'),
    (0, '比喻：5 個偵探同時從現場出發，最後選出最合理的那條路'),
], TL, TT, TW, TH)

# ══ 24. 輔助任務一：拼音 CTC ════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | ④ 輔助任務①：拼音 CTC（η₁=1.0）', n)
bullets(sl, [
    (0, '動機：中文字元空間太大（667 個），直接辨識難'),
    (1, '拼音只有 63 種音節，學習難度小很多'),
    (1, '讓模型先學「怎麼發音」，再學「對應哪個字」'),
    (0, '做法'),
    (1, 'Transformer 輸出後，分出一個拼音 CTC 分支'),
    (1, '同時訓練：主任務（字元）+ 輔助任務（拼音）'),
    (0, 'η₁（eta₁）= 1.0 的意義'),
    (1, '總損失 = L_char + η₁ × L_pinyin + η₂ × L_session'),
    (1, 'η₁=1.0 → 拼音輔助任務的損失與主任務等權重'),
    (1, '若 η₁=0.5 → 拼音輔助任務影響減半'),
    (0, '超參數搜索結果（論文 Fig.5 熱力圖）'),
    (1, 'η₁=1.0, η₂=1.0 組合在驗證集上效果最好'),
], TL, TT, TW, TH)
add_image(sl, fig('fig5'), IL, IT, IW, IH)

# ══ 25. 輔助任務二：Session + GRL ═══════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | ⑤ 輔助任務②：Session 分類 + GRL', n)
bullets(sl, [
    (0, '問題：5 個 Session 的電極位置略不同 → 特徵分布偏移'),
    (1, '若不處理：跨 Session 測試 CER 大幅下降'),
    (0, '目標：讓 CNN 學到「Session 無關」的通用特徵'),
    (0, '做法'),
    (1, 'CNN 輸出後，分出一個 Session 分類器（判斷是哪個 Session）'),
    (1, '在 CNN 和 Session 分類器之間，插入梯度反轉層（GRL）'),
    (0, '直覺理解'),
    (1, '沒有 GRL：CNN 反而會學到「Session 特有」的特徵'),
    (1, '有 GRL：CNN 被迫學到「讓 Session 分類器混淆」的特徵'),
    (1, '→ CNN 的特徵對所有 Session 看起來都一樣（Domain Invariant）'),
    (0, 'η₂ = 1.0：Session 分類損失與主任務等權重'),
], TL, TT, TW, TH)

# ══ 26. 什麼是 GRL ══════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '六、AuxCEMGR 架構 | 什麼是梯度反轉層（GRL）？', n)
bullets(sl, [
    (0, '正向傳播（Forward Pass）'),
    (1, 'GRL 不做任何改變，資料原樣通過'),
    (1, '效果：Session 分類器正常運作，嘗試分辨不同 Session'),
    (0, '反向傳播（Backward Pass / 梯度更新）'),
    (1, 'GRL 把梯度乘以 −λ（反轉方向，λ=1）'),
    (1, '效果：CNN 收到「讓 Session 分類器做錯」的更新訊號'),
    (0, '最終結果'),
    (1, 'Session 分類器越努力分辨 → CNN 越努力讓特徵無法被分辨'),
    (1, '達到對抗（Adversarial）訓練的效果'),
    (1, 'CNN 特徵：對主任務（字元辨識）有用 + 對 Session 辨識無用'),
    (0, '比喻：像「反間諜特訓」'),
    (1, '每次偵探（Session 分類器）學會識破你'),
    (1, 'GRL 讓你反過來學會讓偵探更困惑'),
], TL, TT, TW, TH)

# ══ 27. 總損失函數 ══════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '七、訓練設定 | 總損失函數', n)
bullets(sl, [
    (0, '公式'),
    (1, 'L_total = L_CTC_char + η₁ × L_CTC_pinyin + η₂ × L_session'),
    (0, '各項解釋'),
    (1, 'L_CTC_char：主任務，中文字元 CTC 損失'),
    (1, 'L_CTC_pinyin：輔助任務①，拼音 CTC 損失（η₁=1.0）'),
    (1, 'L_session：輔助任務②，Session 分類交叉熵損失（η₂=1.0）'),
    (0, '超參數搜索'),
    (1, 'η₁ 和 η₂ 各在 {0.1, 0.5, 1.0, 1.5, 2.0} 中搜索'),
    (1, '最佳組合：η₁=1.0, η₂=1.0'),
    (1, '結果見論文 Fig.5 熱力圖（顏色越深 CER 越低）'),
    (0, '訓練時三個損失同時反向傳播'),
    (1, 'GRL 讓 Session 損失的梯度方向反轉（再次強調）'),
], TL, TT, TW, TH)

# ══ 28. Noam 暖機排程器 ═════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '七、訓練設定 | Noam 暖機學習率排程器', n)
bullets(sl, [
    (0, '問題：為什麼不用固定學習率？'),
    (1, '訓練初期：模型參數隨機，梯度方向不穩，學習率太大 → 發散'),
    (1, '訓練後期：接近最優解，學習率太大 → 振盪，無法收斂'),
    (0, 'Noam 排程器的策略'),
    (1, '前 warmup_steps 步：學習率線性從 0 增加到峰值'),
    (1, '之後：按照 step^(−0.5) 緩慢下降'),
    (0, '公式：lr = d_model^(−0.5) × min(step^(−0.5), step × warmup^(−1.5))'),
    (1, 'd_model=256（模型維度），warmup=4000（本研究設定）'),
    (0, '直覺比喻'),
    (1, '像爬山：先慢慢走到山腰，確認方向後再加速'),
    (1, '接近山頂時再放慢腳步，避免越過最高點'),
    (0, '實作：PyTorch 中用 LambdaLR + 自訂 lambda 函數'),
], TL, TT, TW, TH)

# ══ 29. 其他訓練設定 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '七、訓練設定 | 其他超參數總覽', n)
bullets(sl, [
    (0, '優化器：Adam（Adaptive Moment Estimation）'),
    (1, '結合 Momentum（加速收斂）和 RMSProp（自適應學習率）優點'),
    (1, '超參數：β₁=0.9, β₂=0.98, ε=1e-9（Transformer 標準設定）'),
    (0, '批次大小（Batch Size）：128'),
    (1, '每次參數更新使用 128 個訓練句子'),
    (1, '較大批次 → 梯度估計更穩定，但記憶體需求高'),
    (0, '提前停止（Early Stopping）'),
    (1, '監控驗證集 CER，連續 10 個 Epoch 不改善就停止'),
    (1, '防止過擬合，節省訓練時間'),
    (0, 'Dropout'),
    (1, 'Transformer 層間加入 Dropout（rate=0.1）'),
    (1, '訓練時隨機關閉 10% 的神經元 → 防止過擬合'),
    (0, '框架：PyTorch（論文未明確說明，依架構推測）'),
], TL, TT, TW, TH)

# ══ 30. CER 評估指標 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '八、評估指標 | 字元錯誤率（CER）', n)
bullets(sl, [
    (0, '公式：CER = (S + D + I) / N'),
    (1, 'S（Substitution）：替換錯誤，「好」→「號」'),
    (1, 'D（Deletion）：刪除錯誤，「你好世界」→「你世界」'),
    (1, 'I（Insertion）：插入錯誤，「你好」→「你啊好」'),
    (1, 'N：參考文字總字元數'),
    (0, '計算工具：Levenshtein 距離（編輯距離）'),
    (1, '動態規劃演算法，找出最少操作次數把辨識結果變成正確答案'),
    (1, 'Python：from jiwer import cer，或手動實作'),
    (0, '範例計算'),
    (1, '參考：「你好世界」（N=4）'),
    (1, '辨識：「你世界」→ 刪除「好」：D=1'),
    (1, 'CER = 1/4 = 25%'),
    (0, 'CER 越低越好，0% 表示完全正確'),
], TL, TT, TW, TH)
add_image(sl, fig('chart_cer'), IL, IT, IW, IH)

# ══ 31. 主要實驗結果 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | 主要實驗結果（測試集）', n)
bullets(sl, [
    (0, '對比模型'),
    (1, 'Wadkins 2019（LSTM+CTC）：66.8%（SSI 研究的早期基準）'),
    (1, 'Baseline 無增強：66.5%'),
    (1, 'Baseline 完整增強：44.5%'),
    (1, 'AuxCEMGR 無增強：62.5%'),
    (1, '★ AuxCEMGR 完整增強：38.0%（最佳）'),
    (0, '關鍵結論'),
    (1, '單純加輔助任務（AuxCEMGR 無增強）：62.5%，比 Baseline 差'),
    (1, '→ 輔助任務需要搭配增強才能發揮作用'),
    (1, '完整方法（模型 + 增強）：38.0%，大幅領先'),
    (0, '改善幅度：比 Wadkins 2019 基準降低 28.8 個百分點'),
], TL, TT, TW, TH)
add_image(sl, fig('chart_main'), IL, IT, IW, IH)

# ══ 32. 資料增強消融 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | 資料增強消融實驗', n)
bullets(sl, [
    (0, '消融實驗（Ablation Study）：逐一移除每種增強，觀察影響'),
    (0, '結果'),
    (1, '移除頻譜相減：CER 上升最多（最重要的增強）'),
    (1, '移除有聲 EMG 混合：CER 明顯上升'),
    (1, '移除 Mixup：CER 小幅上升'),
    (0, '三種增強的相對重要性'),
    (1, '頻譜相減 > 有聲 EMG 混合 > Mixup'),
    (0, '結論'),
    (1, '三種增強相互補充，缺少任何一種都會影響最終效果'),
    (1, '組合使用比單獨使用效果更好（協同效應）'),
], TL, TT, TW, TH)
add_image(sl, fig('chart_aug'), IL, IT, IW, IH)

# ══ 33. η 超參數熱力圖 ══════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | η₁/η₂ 超參數熱力圖（Fig.5）', n)
bullets(sl, [
    (0, '搜索範圍：η₁, η₂ ∈ {0.1, 0.5, 1.0, 1.5, 2.0}'),
    (0, '讀圖方式：顏色越深（藍）→ CER 越低 → 效果越好'),
    (0, '最佳組合：η₁=1.0, η₂=1.0（CER=38.0%）'),
    (0, '觀察'),
    (1, 'η₁ 太小（0.1）：拼音輔助任務影響不足，效果差'),
    (1, 'η₁ 太大（2.0）：過度強調拼音，可能干擾主任務'),
    (1, '同樣的規律適用於 η₂'),
    (0, '結論：兩個輔助任務等權重（各為 1.0）效果最佳'),
], TL, TT, TW, TH)
add_image(sl, fig('fig5'), IL, IT, IW, IH)

# ══ 34. γ 參數分析 ══════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | γ 參數敏感度分析（Fig.6）', n)
bullets(sl, [
    (0, '搜索範圍：γ ∈ {0, 0.2, 0.4, 0.6, 0.8, 1.0}'),
    (0, '結果'),
    (1, 'γ=0（純靜默）：CER 最高，效果最差'),
    (1, 'γ=0.8：CER 最低，最佳值'),
    (1, 'γ=1.0（純有聲）：CER 略微上升'),
    (0, '解讀'),
    (1, 'γ 太小：有聲訊號增強不足，靜默特徵仍微弱'),
    (1, 'γ=0.8：最佳平衡，有足夠的訊號強度增強，但不過度'),
    (1, 'γ=1.0：訓練資料（有聲）和測試資料（靜默）分布差異最大，泛化差'),
    (0, '結論：γ=0.8 是甜蜜點，接近完全有聲但保留少量靜默成分'),
], TL, TT, TW, TH)
add_image(sl, fig('fig6'), IL, IT, IW, IH)

# ══ 35. 電極通道分析 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | 電極通道重要性分析（Fig.8）', n)
bullets(sl, [
    (0, '實驗方式：每次只用 1 個通道，測 CER（單通道消融）'),
    (0, '結果排名（CER 從低到高 = 重要性從高到低）'),
    (1, 'CH4（下巴右側）：CER 最低 → 最重要通道'),
    (1, 'CH3（下巴左側）：第二重要'),
    (1, 'CH1, CH2（上嘴唇）：中等重要'),
    (1, 'CH7, CH8（臉頰）：較低重要'),
    (1, 'CH5, CH6（喉部）：CER 最高 → 最不重要'),
    (0, '為什麼喉部通道最不重要？'),
    (1, '靜默說話時，聲帶不振動 → 喉部肌肉幾乎不收縮'),
    (1, '喉部通道主要在有聲說話時才有強烈訊號'),
    (0, '應用：未來設計輕量版系統時，可優先選 CH3/CH4（下巴）'),
], TL, TT, TW, TH)
add_image(sl, fig('fig8'), IL, IT, IW, IH)

# ══ 36. 句子長度分析 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | 句子長度 vs CER（Fig.10）', n)
bullets(sl, [
    (0, '觀察：句子越長，CER 越高'),
    (1, '短句（5–10 字）：CER 相對低'),
    (1, '中句（10–20 字）：CER 明顯上升'),
    (1, '長句（20+ 字）：CER 大幅上升'),
    (0, '原因分析'),
    (1, 'CTC 的對齊難度隨句子長度增加'),
    (1, 'Beam Search 路徑數量爆炸，找不到最優路徑'),
    (1, '長句中，前半段的錯誤會影響後半段的預測'),
    (0, '改善方向'),
    (1, '加入語言模型（Language Model）輔助解碼'),
    (1, '語言模型能根據上下文修正不合理的字元序列'),
    (1, '例如：出現「籃球運動」時，下一個字更可能是「員」而非「魚」'),
], TL, TT, TW, TH)
add_image(sl, fig('fig10'), IL, IT, IW, IH)

# ══ 37. 跨 Session 分析 ═════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '九、實驗結果 | 跨 Session 分析（Fig.11/12）', n)
bullets(sl, [
    (0, '跨 Session 挑戰（Fig.11）'),
    (1, '5 個 Session 之間的 CER 差異最高達 ±10%'),
    (1, '原因：電極位置差異 → 訊號分布偏移'),
    (1, 'AuxCEMGR 的 GRL 有效縮小 Session 間差距'),
    (0, '有聲 vs 靜默 EMG 性能比較（Fig.12）'),
    (1, '有聲 EMG 訓練 → 靜默 EMG 測試：CER 很高（分布差異大）'),
    (1, '靜默 EMG 訓練 → 靜默 EMG 測試：最佳（本研究方法）'),
    (1, '有聲 EMG 資料增強：縮小此差距，讓訓練分布更接近測試分布'),
    (0, '結論：Session 偏移是實際部署的最大挑戰，GRL 是解決關鍵'),
], TL, TT, TW, TH)
add_image(sl, fig('fig11'), IL, IT, IW, IH)

# ══ 38. 結論 ════════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '十、結論與未來方向 | 結論', n)
bullets(sl, [
    (0, '本研究貢獻'),
    (1, '建立第一個中文臉部 sEMG 無聲語音辨識系統（AuxCEMGR）'),
    (1, '最佳測試集 CER = 38.0%，優於所有基準'),
    (1, '證明三種資料增強 + 兩個輔助任務各自都有貢獻'),
    (0, '關鍵技術貢獻'),
    (1, '拼音 CTC 輔助任務：幫助學習中文音韻結構'),
    (1, 'GRL + Session 分類：有效消除電極位置偏移'),
    (1, '頻譜相減 + 有聲 EMG 混合 + Mixup：三種增強協同作用'),
    (0, '研究限制'),
    (1, '只有 1 名受試者，跨人泛化能力未知'),
    (1, 'NBA 語料庫特定性強，不一定代表日常對話'),
    (1, 'CER 38% 仍然很高，實用化還有距離'),
], TL, TT, TW, TH)

# ══ 39. 未來方向 ════════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '十、結論與未來方向 | 未來研究方向', n)
bullets(sl, [
    (0, '①多受試者跨人辨識（Cross-subject）'),
    (1, '目前只有 1 人，需增加多人資料集'),
    (1, '挑戰：不同人的臉部肌肉解剖結構不同'),
    (0, '②加入語言模型（Language Model）解碼'),
    (1, '改善長句辨識效果（解決 CER 隨句長增加的問題）'),
    (0, '③擴充語料庫'),
    (1, '不同領域、方言、說話速度的訓練資料'),
    (0, '④硬體微型化'),
    (1, '從桌面設備→可穿戴式（眼鏡型、面罩型）'),
    (0, '⑤多模態融合'),
    (1, '結合嘴型影像（視訊）+ sEMG → 互補資訊，降低 CER'),
    (0, '⑥即時性優化'),
    (1, '目前為離線（Offline）處理，未來需要即時（Online）推論'),
], TL, TT, TW, TH)

# ══ 40. 實驗復現指南 ════════════════════════════════════════════
n += 1; sl = blank()
add_header(sl, '十一、實驗復現指南 | 完整流程', n)
bullets(sl, [
    (0, '環境安裝'),
    (1, 'pip install torch librosa scipy numpy jiwer'),
    (0, 'Step 1：訊號前處理'),
    (1, 'butter(4,[10,400],"band",1000) + iirnotch([50,150,250,350])'),
    (0, 'Step 2：特徵萃取 MFSC'),
    (1, '36 個 Mel 濾波器，Hanning 窗，10 ms 移位'),
    (1, 'librosa.filters.mel(sr=1000, n_fft=25, n_mels=36)'),
    (0, 'Step 3：資料增強'),
    (1, '頻譜相減（γ_noise=1.0），有聲混合（γ=0.8），Mixup（α=0.02）'),
    (0, 'Step 4：模型建立'),
    (1, 'nn.Conv2d(8,256,3,stride=2) × 2 → TransformerEncoder(6層,8頭,d=256)'),
    (1, '→ CTC Linear 輸出層（字元 + 拼音 + Session 三個頭）'),
    (0, 'Step 5：訓練'),
    (1, 'Adam + Noam(warmup=4000) + Batch=128 + EarlyStopping'),
    (0, 'Step 6：解碼評估'),
    (1, 'Beam Search(size=5) → CER 計算（jiwer.cer）'),
], TL, TT, TW, TH)

prs.save(OUT_PATH)
print(f'✓ Python-pptx 版完成：{OUT_PATH}（{n} 頁）')
