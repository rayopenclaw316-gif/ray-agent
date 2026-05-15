#!/usr/bin/env python3
"""
PPT Generator: Neural Chinese Silent Speech Recognition with Facial Electromyography
Xie et al., Speech Communication 2025
"""
import os
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Cm(33.87)
prs.slide_height = Cm(19.05)

# ── 顏色 ──────────────────────────────────────────────
C_DARK  = RGBColor(0x1F, 0x38, 0x64)
C_MID   = RGBColor(0x2E, 0x75, 0xB6)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED   = RGBColor(0xC0, 0x00, 0x00)
C_ORNG  = RGBColor(0xC5, 0x60, 0x00)
C_GRN   = RGBColor(0x00, 0x6B, 0x00)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LIGHT = RGBColor(0xD6, 0xE4, 0xF7)

FONT = '標楷體'
blank = prs.slide_layouts[6]

# ── 基礎工具函數 ───────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(blank)

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txtbox(slide, l, t, w, h, items):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        if item.get('sb', 0):
            p.space_before = Pt(item['sb'])
        run = p.add_run()
        indent = '    ' * item.get('indent', 0)
        blt    = item.get('blt', '')
        run.text = indent + (blt + ' ' if blt else '') + item.get('text', '')
        run.font.name  = item.get('font', FONT)
        run.font.size  = item.get('size', Pt(20))
        run.font.bold  = item.get('bold', False)
        run.font.color.rgb = item.get('color', C_BLACK)
    return box

def hdr(slide, title):
    rect(slide, Cm(0), Cm(0), Cm(33.87), Cm(2.3), C_DARK)
    txtbox(slide, Cm(0.5), Cm(0.25), Cm(33), Cm(1.85), [
        {'text': title, 'size': Pt(26), 'bold': True, 'color': C_WHITE}
    ])

def mk(title, items):
    s = new_slide()
    hdr(s, title)
    txtbox(s, Cm(0.8), Cm(2.6), Cm(32.2), Cm(15.7), items)
    return s

# ── 條目快捷函數 ──────────────────────────────────────
def b0(text, bold=False, color=C_BLACK, size=Pt(20), sb=8):
    return {'text': text, 'blt': '•', 'indent': 0, 'bold': bold,
            'color': color, 'size': size, 'sb': sb}

def b1(text, bold=False, color=C_GRAY, size=Pt(18), sb=3):
    return {'text': text, 'blt': '─', 'indent': 1, 'bold': bold,
            'color': color, 'size': size, 'sb': sb}

def b2(text, bold=False, color=C_GRAY, size=Pt(17), sb=2):
    return {'text': text, 'blt': '·', 'indent': 2, 'bold': bold,
            'color': color, 'size': size, 'sb': sb}

def h(text, color=C_MID, size=Pt(21), sb=12):
    return {'text': text, 'blt': '', 'indent': 0, 'bold': True,
            'color': color, 'size': size, 'sb': sb}

def t(text='', bold=False, color=C_BLACK, size=Pt(20), sb=0,
      align=PP_ALIGN.LEFT):
    return {'text': text, 'blt': '', 'indent': 0, 'bold': bold,
            'color': color, 'size': size, 'sb': sb, 'align': align}

# ══════════════════════════════════════════════════════
# 投影片 1：封面
# ══════════════════════════════════════════════════════
s1 = new_slide()
rect(s1, Cm(0), Cm(0),   Cm(33.87), Cm(19.05), C_DARK)
rect(s1, Cm(0), Cm(5.5), Cm(33.87), Cm(9),     C_MID)
txtbox(s1, Cm(1), Cm(6.0), Cm(32), Cm(3.5), [
    {'text': '從臉部肌電訊號到中文文字', 'size': Pt(38), 'bold': True,
     'color': C_WHITE, 'align': PP_ALIGN.CENTER},
])
txtbox(s1, Cm(1), Cm(9.5), Cm(32), Cm(2.5), [
    {'text': 'AuxCEMGR 模型深度解析', 'size': Pt(28), 'bold': True,
     'color': C_LIGHT, 'align': PP_ALIGN.CENTER},
])
txtbox(s1, Cm(1), Cm(12.5), Cm(32), Cm(3.5), [
    {'text': 'Neural Chinese Silent Speech Recognition with Facial Electromyography',
     'size': Pt(17), 'color': C_LIGHT, 'align': PP_ALIGN.CENTER},
    {'text': 'Xie et al.  |  Speech Communication 171 (2025) 103230',
     'size': Pt(17), 'color': C_LIGHT, 'align': PP_ALIGN.CENTER, 'sb': 8},
    {'text': '國防創新研究院 / 北京大學 / 哈工大（深圳）',
     'size': Pt(16), 'color': C_LIGHT, 'align': PP_ALIGN.CENTER, 'sb': 6},
])

# ══════════════════════════════════════════════════════
# 投影片 2：研究背景 (1)
# ══════════════════════════════════════════════════════
mk('研究背景 (1)：語音辨識的成功與限制', [
    h('自動語音辨識（ASR）已非常成熟'),
    b0('智慧音箱（Siri、小愛）、影片自動字幕、客服機器人、文字聽打'),
    b0('關鍵突破：2017 年 Transformer 架構讓 ASR 性能大幅提升'),
    b0('現代 ASR 在英文 LibriSpeech 資料集上已接近人類水準'),
    t(),
    h('但 ASR 有一個根本限制'),
    b0('必須有「聲音輸入」才能運作'),
    b0('沒有聲音？→ 傳統 ASR 完全失效', color=C_RED, bold=True),
    t(),
    h('→ 無聲語音辨識（Silent Speech Recognition, SSR）因此誕生'),
    b0('目標：不需要聲音輸出，直接從肌肉動作還原語意'),
    b0('本論文使用的訊號：臉部表面肌電圖（Facial sEMG）'),
])

# ══════════════════════════════════════════════════════
# 投影片 3：研究背景 (2) — 無聲語音的應用場景
# ══════════════════════════════════════════════════════
mk('研究背景 (2)：什麼時候需要無聲語音辨識？', [
    h('場景 A：醫療需求'),
    b0('喉切除術患者（laryngectomy）：聲帶移除，無法發聲'),
    b0('語音復健：協助患者學習無聲溝通'),
    t(),
    h('場景 B：噪音 / 保密需求'),
    b0('工廠、戰場等高噪音環境，麥克風無法正常使用'),
    b0('軍事 / 特殊任務中需要靜默下達指令（Kapur et al., 2018 — AlterEgo）'),
    t(),
    h('場景 C：多模態語音強化'),
    b0('在嘈雜環境下，用 EMG 輔助聲學訊號提升辨識準確率'),
    b0('即使有聲音，EMG 仍可作為補充資訊'),
    t(),
    h('現有研究的空缺'),
    b0('幾乎所有 EMG-SSR 研究都聚焦英文', color=C_ORNG),
    b0('中文（普通話）是聲調語言，問題更複雜', color=C_ORNG),
    b0('中文端對端 EMG-SSR → 本論文是全球第一篇！', color=C_RED, bold=True),
])

# ══════════════════════════════════════════════════════
# 投影片 4：sEMG 基本原理
# ══════════════════════════════════════════════════════
mk('sEMG 是什麼？（從零開始的外行說明）', [
    h('EMG = 表面肌電圖（Surface Electromyography）'),
    b0('肌肉收縮時，神經細胞向肌纖維發出電訊號（動作電位）'),
    b0('動作電位傳導到皮膚表面 → 產生微弱電壓變化（單位：微伏特 μV）'),
    b0('「表面」電極貼在皮膚外側即可偵測，不需穿刺 → 非侵入式', color=C_GRN),
    t(),
    h('與語音辨識的關聯'),
    b0('說話 = 大腦控制臉部 / 喉部肌肉 → 產生聲音'),
    b0('默唸時：這些肌肉仍然有微弱的收縮動作（肌電訊號存在！）', color=C_GRN, bold=True),
    b0('結論：即使完全靜默，臉部 EMG 仍帶有「你在說什麼」的資訊'),
    t(),
    h('直觀類比'),
    b0('就像想像自己在跑步，腿部肌肉也有微弱的電訊號'),
    b0('默唸「籃球」時，下頜 / 嘴唇的肌肉仍然做出對應動作'),
    b1('但幅度遠比發聲時小 → 這是靜默辨識困難的核心原因', color=C_ORNG),
])

# ══════════════════════════════════════════════════════
# 投影片 5：臉部說話肌群
# ══════════════════════════════════════════════════════
mk('說話時臉部肌肉的運動分析', [
    h('本論文使用的 8 個 EMG 通道對應肌群'),
    b0('CH1：下頜肌（Jaw）── 控制開閉口，訊號能量最強'),
    b0('CH2：嘴輪匝肌（Orbicularis Oris）── 控制嘴形（圓嘴 / 扁嘴）'),
    b1('Baseline 中此通道貢獻最大', color=C_GRAY),
    b0('CH3：臉頰肌群'),
    b0('CH4：下頜骨肌群（Mandibular Muscles）── 最重要！', bold=True, color=C_RED),
    b1('AuxCEMGR 中此通道貢獻最大', color=C_GRAY),
    b0('CH5、CH6：喉部周圍（Larynx 附近）'),
    b1('有聲模式訊號明顯；靜默模式幾乎靜止 → 靜默辨識中效果最差', color=C_ORNG),
    b0('CH7、CH8：其他臉部輔助肌群'),
    t(),
    h('靜默 vs. 有聲的關鍵差異'),
    b0('下頜（CH1）：有聲 ≈ 靜默（兩者都在動）'),
    b0('喉部（CH5/6）：有聲時明顯活躍，靜默時幾乎沒有動作'),
    b0('整體：靜默 EMG 分佈與有聲相似，但能量幅度較低', color=C_ORNG),
])

# ══════════════════════════════════════════════════════
# 投影片 6：三大貢獻
# ══════════════════════════════════════════════════════
mk('本論文的三大貢獻', [
    h('貢獻一：第一個中文 EMG-SSR 基準資料集', color=C_RED),
    b0('過去研究幾乎全是英語；中文是聲調語言，研究空白'),
    b0('建立可公開申請的平行語料庫（EMG 訊號 ↔ 中文字元）'),
    b1('語料：NBA 賽事報導，1,238 句 / 667 個獨特字元', color=C_GRAY),
    t(),
    h('貢獻二：第一個中文端對端神經網路 SSR 模型', color=C_RED),
    b0('架構：Transformer 編碼器 + CTC 解碼器'),
    b0('直接 EMG → 中文字元（無需中間步驟，end-to-end）'),
    b0('優於傳統 GMM-HMM / LSTM 方法'),
    t(),
    h('貢獻三：創新輔助任務 + 資料增強策略', color=C_RED),
    b0('輔助任務 1：拼音生成（利用漢語音節特性）'),
    b0('輔助任務 2：Session 分類（梯度反轉，消除電極差異干擾）'),
    b0('三種資料增強：頻譜相減、有聲 EMG 輔助、Mixup'),
    t(),
    b0('最終結果：測試集 CER = 38.0%（字元錯誤率，越低越好）',
       bold=True, color=C_MID),
])

# ══════════════════════════════════════════════════════
# 投影片 7：資料集 — 語料與受試者
# ══════════════════════════════════════════════════════
mk('資料集建立 (1)：語料來源與受試者', [
    h('受試者'),
    b0('1 位女性志願者（母語普通話，無閱讀障礙，無語音障礙）'),
    b0('接受簡單任務訓練，給予合理報酬確保資料品質'),
    b0('主要局限：僅單一受試者 → 無法反映個體差異', color=C_RED),
    t(),
    h('語料來源：新浪體育 NBA 賽事報導（2012 年起）'),
    b0('爬取工具：Selenium + BeautifulSoup'),
    b0('篩選條件：不超過 20 字、包含常見籃球術語、過濾 HTML 標籤'),
    b0('語料特性：主題集中（NBA）、句子短、用字常見，適合初期研究'),
    t(),
    h('語料統計'),
    b0('總計：1,238 句 / 12,584 字 / 667 個獨特字元'),
    b0('平均每句：10 個字元'),
    b0('資料集切分：'),
    b1('訓練集 1,062 句（涵蓋全部 667 個字元）'),
    b1('驗證集 63 句'),
    b1('測試集 113 句'),
])

# ══════════════════════════════════════════════════════
# 投影片 8：資料集 — 硬體設備
# ══════════════════════════════════════════════════════
mk('資料集建立 (2)：錄製設備規格', [
    h('EMG 設備（主要）'),
    b0('型號：Neuracle Technology NSW308M 雙極系統'),
    b0('電極材質：Ag/AgCl（氯化銀）差分電極'),
    b1('差分電極 = 兩個電極相減 → 消除共模雜訊', color=C_GRAY),
    b0('採樣率：1000 Hz（每秒採集 1000 個資料點）'),
    b0('通道數：8 通道（16 個差分電極對）'),
    b0('參考電極：1 個，貼於左側鎖骨（接地用）'),
    t(),
    h('音訊設備（輔助收音）'),
    b0('一般播放用麥克風'),
    b0('採樣率：44,100 Hz，單聲道（1 Channel）'),
    t(),
    h('錄製環境'),
    b0('安靜房間，減少環境電磁干擾'),
    b0('每次錄製前：重新撕掉電極貼片、重新貼上，確保接觸良好'),
    b1('→ 這個步驟導致每次位置略有偏移，是 Session 差異的來源', color=C_ORNG),
])

# ══════════════════════════════════════════════════════
# 投影片 9：資料集 — 電極位置
# ══════════════════════════════════════════════════════
mk('資料集建立 (3)：電極位置與各通道角色', [
    h('電極配置（參考 Diener et al., 2015 的標準位置）'),
    b0('共 8 組差分電極對（CH1–CH8），貼於臉部不同肌肉群上'),
    b0('另有 1 個接地 / 參考電極貼於左側鎖骨'),
    t(),
    h('各通道位置與功能'),
    b0('CH1：下頜肌 → 控制開閉口（訊號能量最強）'),
    b0('CH2：嘴輪匝肌 → 控制嘴形（Baseline 最重要通道）'),
    b0('CH3：臉頰肌群（顴骨附近）'),
    b0('CH4：下頜骨肌群 → AuxCEMGR 最重要通道', bold=True, color=C_RED),
    b0('CH5、CH6：喉部周圍 → 靜默模式效果最差（幾乎不動）'),
    b0('CH7、CH8：其他臉部輔助肌群'),
    t(),
    h('電極數量的影響（實驗結論）'),
    b0('電極對越多 CER 越低；但只要保留 6–7 對，效能差距不大'),
    b0('移除超過 2 對電極 → 準確率急遽下降', color=C_RED),
    b0('→ 實際應用可縮減至 6 對，兼顧舒適性', color=C_GRN),
])

# ══════════════════════════════════════════════════════
# 投影片 10：資料集 — 錄製流程
# ══════════════════════════════════════════════════════
mk('資料集建立 (4)：錄製流程與 Session 設計', [
    h('兩種錄製模式'),
    b0('靜默模式（Silent Mode）── 主要資料'),
    b1('受試者默唸每個句子，不發出任何聲音'),
    b1('輸出：EMG 訊號 + 文字標籤（EMG-文字平行語料）'),
    b0('有聲模式（Audible Mode）── 輔助資料'),
    b1('受試者大聲朗讀相同句子'),
    b1('輸出：EMG + 音訊 + 文字標籤（三元組）'),
    t(),
    h('Session 設計：為什麼要錄 5 次？'),
    b0('每個句子在每種模式下各錄製 5 個 session'),
    b0('每個 session 之間：撕掉電極貼片，重新貼上再錄'),
    b0('目的：模擬現實使用情境（每次配戴位置不完全一樣）'),
    b0('副作用：同一句話在不同 session 間 EMG 訊號略有差異', color=C_ORNG),
    b1('→ 這是設計「Session 分類輔助任務」的核心動機', color=C_GRAY),
    t(),
    h('最終資料量'),
    b0('靜默 EMG：5,468 組 EMG-文字對（共 5.93 小時）'),
    b0('有聲 EMG：5,829 組三元組（共 5.80 小時）'),
    b0('每段訊號平均長度：3.74 秒', color=C_ORNG),
])

# ══════════════════════════════════════════════════════
# 投影片 11：前處理 — 濾波
# ══════════════════════════════════════════════════════
mk('訊號前處理 (1)：濾波去雜訊', [
    h('原始 EMG 含有哪些雜訊？'),
    b0('生理雜訊：心跳（~1 Hz）、基線漂移（呼吸 / 汗水）'),
    b0('電源雜訊：交流電 50 Hz 及其諧波（150、250、350 Hz）'),
    b0('高頻雜訊：EMG 有用範圍以外的電磁干擾'),
    t(),
    h('步驟 1：帶通濾波器（Bandpass Filter）── 保留有效頻率'),
    b0('類型：4 階 Butterworth 濾波器（特性：通帶平坦、過渡帶陡峭）'),
    b0('通頻範圍：10 Hz – 400 Hz'),
    b1('< 10 Hz → 濾除（基線漂移、心跳）', color=C_GRAY),
    b1('> 400 Hz → 濾除（高頻無用雜訊）', color=C_GRAY),
    b1('實作：scipy.signal.butter() + filtfilt()（零相位）', color=C_GRAY),
    t(),
    h('步驟 2：陷波濾波器（Notch Filter）── 移除電源雜訊'),
    b0('目標：精確移除 50 Hz、150 Hz、250 Hz、350 Hz'),
    b0('原理：在特定頻率「挖個洞」，只讓那個頻率的訊號歸零'),
    b1('實作：scipy.signal.iirnotch()，對每個頻率分別套用', color=C_GRAY),
])

# ══════════════════════════════════════════════════════
# 投影片 12：前處理 — MFSC
# ══════════════════════════════════════════════════════
mk('訊號前處理 (2)：MFSC 特徵提取', [
    h('為什麼不直接用原始 EMG 訊號？'),
    b0('原始訊號維度太高（1000 Hz × 8 通道 × 數秒 ≈ 數千個數字）'),
    b0('其中有大量冗餘，重要的是「頻率隨時間的變化」，不是逐點數值'),
    t(),
    h('MFSC = log 梅爾頻率譜係數（源自語音辨識領域）'),
    b0('Step 1：用 Hanning 窗口切分訊號（滑動視窗，帶重疊）'),
    b0('Step 2：短時傅立葉轉換（STFT）→ 每個視窗的頻域能量分佈'),
    b0('Step 3：套用 36 個梅爾濾波器'),
    b1('梅爾尺度模仿人耳：低頻解析度高、高頻解析度低', color=C_GRAY),
    b0('Step 4：取對數（log）→ 壓縮數值範圍，讓訓練更穩定'),
    b0('結果：每個時間視窗 → 1 個 36 維特徵向量'),
    t(),
    h('實作工具'),
    b0('Librosa：librosa.feature.melspectrogram()'),
    b0('Scipy：濾波操作'),
    b0('MFSC 是 MFCC 的前驅版（不做離散餘弦轉換），保留更多資訊', color=C_GRAY, size=Pt(18)),
])

# ══════════════════════════════════════════════════════
# 投影片 13：模型架構總覽
# ══════════════════════════════════════════════════════
mk('模型架構總覽：AuxCEMGR', [
    h('AuxCEMGR = Auxiliary-enhanced Chinese EMG Recognition'),
    t(),
    h('整體流程（EMG → 中文文字）'),
    b0('輸入：臉部 sEMG（8 通道 × 36 維 MFSC 特徵）'),
    b0('Step 1：2 層 CNN → 壓縮序列長度，提取局部特徵'),
    b0('Step 2：6 層 Transformer 編碼器 → 學習全局時序依賴'),
    b0('Step 3：CTC 解碼器 → 輸出中文字元序列（主任務）'),
    t(),
    h('兩個輔助任務（與主任務同步訓練）'),
    b0('輔助任務 1：拼音生成 CTC 解碼器'),
    b1('強迫編碼器學到更好的音素特徵', color=C_GRAY),
    b0('輔助任務 2：Session 分類器 + 梯度反轉層（GRL）'),
    b1('讓編碼器特徵對電極位置差異不敏感', color=C_GRAY),
    t(),
    h('架構示意（文字版）'),
    t('sEMG → [CNN] → [Transformer] ──→ [CTC] → 中文文字',
      color=C_MID, size=Pt(19)),
    t('                     │                 ',
      color=C_MID, size=Pt(19)),
    t('                     ├──→ [拼音 CTC]（輔助任務 1）',
      color=C_MID, size=Pt(19)),
    t('                     └──→ [GRL + Session 分類]（輔助任務 2）',
      color=C_MID, size=Pt(19)),
])

# ══════════════════════════════════════════════════════
# 投影片 14：CNN 層
# ══════════════════════════════════════════════════════
mk('編碼器 Part 1：CNN 壓縮層', [
    h('為什麼先用 CNN？'),
    b0('Transformer 計算複雜度為 O(T²)（T = 序列長度）'),
    b0('原始 MFSC 序列：3.74 秒 × 約 100 幀/秒 ≈ 374 幀 → T 很大'),
    b0('CNN 先壓縮序列長度（N → T，T < N），大幅降低 Transformer 的計算量'),
    t(),
    h('CNN 設定'),
    b0('層數：2 層 Convolutional Neural Network'),
    b0('每層 Kernel size：3 × 3（時間軸 × 特徵軸各延伸 3 個單位）'),
    b0('每層有 stride（步長）→ 輸出序列長度縮短'),
    t(),
    h('CNN 在這裡的功能'),
    b0('提取局部時序模式：相鄰幾個時間步之間的肌肉動作變化'),
    b0('類比：就像圖像辨識的 CNN 學習「邊緣 / 紋理」，'),
    b1('這裡學的是「EMG 波形的局部變化模式」', color=C_GRAY),
    t(),
    h('公式'),
    b0('X = CNN×2(S)  →  輸出 X = x₁...xT（T < N）'),
    b1('第一層 Transformer 含位置嵌入（Positional Encoding），其餘層不含',
       color=C_GRAY),
])

# ══════════════════════════════════════════════════════
# 投影片 15：Transformer 編碼器
# ══════════════════════════════════════════════════════
mk('編碼器 Part 2：Transformer 自注意力機制（原理）', [
    h('Transformer 是什麼？'),
    b0('2017 年 Google 提出（"Attention is All You Need"）'),
    b0('核心：自注意力（Self-Attention）── 讓每個位置都能「看見」序列中所有其他位置'),
    b0('優勢：完全平行計算（不像 RNN 需要按順序），能捕捉長距離依賴'),
    t(),
    h('自注意力的直觀理解'),
    b0('處理 EMG 的每個時間步 t 時，同時參考所有其他時間步的資訊'),
    b0('學習：「說這個字的某時刻，應該最關注哪些其他時刻的肌肉狀態」'),
    b1('例：說「籃球」時，下頜動作和嘴唇動作的時序關係', color=C_GRAY),
    t(),
    h('本論文設定'),
    b0('層數：6 層 Transformer（堆疊 6 個相同結構）'),
    b0('注意力頭數：8 個頭（Multi-head Attention，從 8 個角度看序列）'),
    b0('特徵維度：256 維'),
    b0('每層結構：Self-Attention → Add & Norm → FFN → Add & Norm'),
    t(),
    h('為什麼 Transformer >> LSTM？（實驗數據）'),
    b0('Baseline（Transformer）CER：44.5%',
       color=C_GRN, bold=True),
    b0('Wadkins 2019（LSTM + CTC）CER：66.8%',
       color=C_RED, bold=True),
    b1('差距 22.3%，Transformer 大幅優勝', color=C_GRAY),
])

# ══════════════════════════════════════════════════════
# 投影片 16：CTC 解碼器原理
# ══════════════════════════════════════════════════════
mk('解碼器 (1)：CTC 的核心問題與解法', [
    h('根本問題：序列對齊（Alignment Problem）'),
    b0('EMG 序列有 T 個時間步（例如 374 個）'),
    b0('輸出文字只有 M 個字元（例如 10 個）'),
    b0('問題：T >> M，且無法事先知道哪個時間步對應哪個字', color=C_RED),
    b1('傳統 HMM-GMM 需要手動對齊 → 費時費力', color=C_GRAY),
    t(),
    h('CTC 的解法：引入「空白符號（─）」'),
    b0('在每個時間步 t，模型從字元表中選一個輸出（含特殊空白符號 ─）'),
    b0('解碼規則：'),
    b1('先合併「相鄰重複字元」'),
    b1('再移除所有「空白符號 ─」'),
    b0('範例：'),
    b1('CTC 原始輸出：猛─龍─已─是─東─部─前─五', color=C_GRAY),
    b1('步驟一（合併重複）：猛龍已是東部前五',     color=C_ORNG),
    b1('步驟二（移除空白）：猛龍已是東部前五 ✓', color=C_GRN),
    t(),
    h('優點'),
    b0('完全不需要手動標注哪個時間步對應哪個字'),
    b0('同一個輸出字串可對應無數種不同的對齊方式 → 訓練訊號豐富'),
])

# ══════════════════════════════════════════════════════
# 投影片 17：CTC 訓練與推理
# ══════════════════════════════════════════════════════
mk('解碼器 (2)：CTC 的訓練與 Beam Search 推理', [
    h('訓練：最大化正確輸出的概率'),
    b0('某個輸出 Y 的概率 = 所有能產生 Y 的合法對齊路徑 C 的概率之和'),
    b0('P(Y|X) = Σ_{C ∈ Ω(C,Y)} ∏ᵀₜ₌₁ πₜ[cₜ]'),
    b0('損失函數：L_CTC = −log P(Y_正確 | X)'),
    b1('Ω(C,Y) 的計算用動態規劃（Forward-Backward Algorithm）', color=C_GRAY),
    t(),
    h('推理：Beam Search（束搜尋）'),
    b0('貪婪搜尋：每步選最高機率字元 → 簡單但容易選錯無法回頭'),
    b0('Beam Search：每步保留前 k 個最可能的序列候選'),
    b1('本論文設定：Beam size = 5（每步保留 5 條路徑）', color=C_GRAY),
    b0('最後選機率最高的那條路徑作為輸出'),
    t(),
    h('CTC 的局限（也是未來改進方向）'),
    b0('假設各時間步輸出相互獨立 → 忽略字與字之間的語言關係', color=C_ORNG),
    b0('→ 未來加入語言模型（LM）解碼可以進一步降低 CER', color=C_GRN),
])

# ══════════════════════════════════════════════════════
# 投影片 18：輔助任務 1 — 拼音
# ══════════════════════════════════════════════════════
mk('輔助任務 1：拼音生成（動機與原理）', [
    h('漢語的特殊性：音節 vs. 字元'),
    b0('普通話約有 400 個音節（無聲調拼音），對應超過 10,000 個字元'),
    b0('一個音節 = 多個字（同音字）'),
    b1('例：音節「ji」可能是「雞、積、機、基、激、籍、跡...」', color=C_GRAY),
    b0('EMG 訊號本質上記錄的是「發音動作」，與音節比與字更直接相關'),
    t(),
    h('為什麼加拼音生成有助於辨識？'),
    b0('強迫 Transformer 編碼器學到更好的「發音特徵」（音素級表示）'),
    b0('讓模型先學容易的（音節 400 個），再映射到難的（字元 667 個）'),
    b0('兩個 CTC 解碼器共用同一個 Transformer 編碼器 → 知識互相強化'),
    t(),
    h('實作細節'),
    b0('在編碼器輸出 H 上，加一個額外的 CTC 解碼器預測無聲調拼音序列'),
    b0('例：「國王連續命中」→「guo wang lian xu ming zhong」'),
    b0('損失函數：L_pin = −log P(拼音序列_正確 | X)'),
    b0('損失權重 η₁ = 1.0（實驗發現 η₁ > 0.5 時 CER 明顯改善）',
       color=C_RED, bold=True),
])

# ══════════════════════════════════════════════════════
# 投影片 19：輔助任務 2 — Session 分類
# ══════════════════════════════════════════════════════
mk('輔助任務 2：Session 分類 + 梯度反轉層（GRL）', [
    h('問題：電極重貼導致 Session 間 EMG 差異'),
    b0('每個 session 電極重新黏貼 → 位置略微偏移'),
    b0('同一句話在不同 session 的 EMG 訊號有微小差異'),
    b0('模型可能「記住」某 session 的電極位置特性，而非真正學到語音特徵',
       color=C_ORNG),
    t(),
    h('解法：梯度反轉層（Gradient Reversal Layer, GRL）'),
    b0('在編碼器輸出 H 上接一個 Session 分類器（Linear → Softmax）'),
    b0('GRL 夾在編碼器輸出和分類器之間，做一件特殊的事：'),
    b1('正向傳播（Forward）：直接通過，不改變數值', color=C_GRN),
    b1('反向傳播（Backward，梯度更新時）：把梯度正負號反轉！', color=C_RED),
    b0('效果：'),
    b1('分類器努力學習區分 session', color=C_GRAY),
    b1('編碼器被迫朝「讓分類器更難分」的方向更新', color=C_GRAY),
    b1('→ 編碼器輸出變得 session-invariant（對電極位置不敏感）', color=C_GRN),
    t(),
    h('損失 & 超參數'),
    b0('L_ses = log P_session(session_正確 | X)'),
    b0('損失權重 η₂ = 1.0'),
])

# ══════════════════════════════════════════════════════
# 投影片 20：損失函數整合
# ══════════════════════════════════════════════════════
mk('訓練損失函數整合', [
    h('三個損失的組合'),
    b0('主任務（CTC 文字生成）：'),
    b1('L_CTC = −log P(Y_正確 | X)', color=C_MID),
    b0('輔助任務 1（拼音生成）：'),
    b1('L_pin = −log P(拼音序列_正確 | X)', color=C_MID),
    b0('輔助任務 2（Session 分類）：'),
    b1('L_ses = log P_session(session_正確 | X)  ← 注意是正號',
       color=C_MID),
    b2('分類器最大化 P_ses；GRL 讓編碼器反向優化', color=C_GRAY),
    t(),
    h('最終合併損失函數'),
    b0('L_all = L_CTC  +  η₁ × L_pin  +  η₂ × L_ses',
       bold=True, color=C_DARK, size=Pt(21)),
    b0('η₁ = 1.0（由驗證集 grid search 決定）'),
    b0('η₂ = 1.0（由驗證集 grid search 決定）'),
    t(),
    h('訓練策略'),
    b0('三個損失同步 end-to-end 訓練（單一反向傳播流程）'),
    b0('優化器：Adam + Noam 排程（warmup 2000 步後學習率衰減）'),
    b0('Batch size：128'),
])

# ══════════════════════════════════════════════════════
# 投影片 21：資料增強 1 — 頻譜相減
# ══════════════════════════════════════════════════════
mk('資料增強 (1)：頻譜相減（Spectral Subtraction）', [
    h('為什麼需要資料增強？'),
    b0('只有 1 位受試者、約 6 小時靜默 EMG → 訓練資料嚴重不足'),
    b0('EMG 訊號個體差異大，電極每次位置略有不同 → 需要更多多樣性'),
    t(),
    h('頻譜相減的原理'),
    b0('頻譜相減是語音增強（Speech Enhancement）的經典演算法'),
    b0('核心思路：估算背景雜訊的頻譜，從原始訊號頻譜中減去'),
    b0('步驟：'),
    b1('Step 1：找出無語音段（靜音段）→ 估算背景雜訊頻譜'),
    b1('Step 2：從原始 EMG 的頻譜中減去雜訊頻譜'),
    b1('Step 3：取得「降噪版」EMG 訊號'),
    b0('效果：每筆訓練資料 → 原版 + 降噪版，資料集自動翻倍！',
       color=C_GRN, bold=True),
    t(),
    h('實驗效果'),
    b0('移除頻譜相減後：'),
    b1('Baseline CER 上升 9.4%（最顯著）', color=C_RED),
    b1('AuxCEMGR CER 上升 7.8%', color=C_RED),
    b0('→ 三種增強中貢獻最大的一種', color=C_RED, bold=True),
])

# ══════════════════════════════════════════════════════
# 投影片 22：資料增強 2 — 有聲 EMG
# ══════════════════════════════════════════════════════
mk('資料增強 (2)：有聲 EMG 輔助訓練', [
    h('概念'),
    b0('錄製時同步收集了有聲模式的 EMG（5,829 組，5.80 小時）'),
    b0('有聲 EMG ≠ 靜默 EMG（喉部差異明顯），但整體發音動作模式相似'),
    b0('→ 有聲 EMG 可作為額外訓練資料，幫助模型泛化', color=C_GRN),
    t(),
    h('實作：加權損失'),
    b0('有聲訓練資料的損失乘上折扣係數 γ（< 1 表示降低可信度）'),
    b0('最佳 γ 值（驗證集實驗決定）：'),
    b1('Baseline：γ = 1.0（等權重）'),
    b1('AuxCEMGR：γ = 0.8（略微折扣）'),
    b0('γ 在 0.6–1.0 之間效果最佳', color=C_ORNG),
    t(),
    h('實驗效果'),
    b0('移除有聲資料後：'),
    b1('Baseline CER 上升 6.9%', color=C_RED),
    b1('AuxCEMGR CER 上升 9.1%（影響更大）', color=C_RED),
    b0('額外發現：加入有聲訓練後，模型在有聲 EMG 上的辨識也大幅改善',
       color=C_GRAY, size=Pt(18)),
])

# ══════════════════════════════════════════════════════
# 投影片 23：資料增強 3 — Mixup
# ══════════════════════════════════════════════════════
mk('資料增強 (3)：Mixup（混合訓練樣本）', [
    h('原始 Mixup 概念（Zhang et al., 2018）'),
    b0('用於圖像分類：把兩張圖線性混合，標籤也按同比例混合'),
    b0('效果：讓模型學到更平滑的決策邊界，避免過擬合'),
    t(),
    h('本論文的改造版：適用於 CTC 序列任務'),
    b0('取兩筆訓練資料 (X₁, Y₁) 和 (X₂, Y₂)'),
    b0('產生混合輸入：X₁₊₂ = λ × X₁ + (1−λ) × X₂'),
    b0('損失近似計算（無需知道確切混合標籤）：'),
    b1('L(X₁₊₂, Y₁₊₂) ≈ λ × L(X₁₊₂, Y₁) + (1−λ) × L(X₁₊₂, Y₂)',
       color=C_MID),
    t(),
    h('超參數設定'),
    b0('λ ~ Beta(α, α) 分佈，α = 0.02'),
    b1('α 非常小 → λ 大多接近 0 或 1（偶爾輕微混合）', color=C_GRAY),
    b1('這是為了保持訓練樣本大致合理，不產生太離奇的混合', color=C_GRAY),
    t(),
    h('實驗效果'),
    b0('移除 Mixup 後：Baseline +5.1%，AuxCEMGR +6.5%', color=C_RED),
])

# ══════════════════════════════════════════════════════
# 投影片 24：訓練超參數
# ══════════════════════════════════════════════════════
mk('訓練設定：照著做的完整超參數清單', [
    h('模型架構'),
    b0('CNN：2 層，Kernel size 3 × 3'),
    b0('Transformer：6 層，8 注意力頭，256 維'),
    b0('CTC 主解碼器：Linear(256 → 字元數) → Softmax'),
    b0('拼音 CTC 解碼器：Linear(256 → 拼音數) → Softmax'),
    b0('Session 分類器：GRL → Linear(256 → 5) → Softmax（5 個 session）'),
    t(),
    h('訓練超參數'),
    b0('優化器：Adam'),
    b0('學習率排程：Noam warmup（初始 lr = 0.1，warmup steps = 2000）'),
    b0('Batch size：128'),
    b0('CTC Beam Search：beam size = 5'),
    t(),
    h('輔助任務 & 增強超參數'),
    b0('η₁（拼音損失權重）= 1.0'),
    b0('η₂（Session 損失權重）= 1.0'),
    b0('有聲資料折扣 γ = 0.8（AuxCEMGR） / 1.0（Baseline）'),
    b0('Mixup α = 0.02'),
    t(),
    h('硬體環境（論文使用）'),
    b0('6 × NVIDIA RTX 3090 GPU，24 核心 CPU'),
])

# ══════════════════════════════════════════════════════
# 投影片 25：CER 評估指標
# ══════════════════════════════════════════════════════
mk('評估指標：字元錯誤率（CER）詳解', [
    h('CER 的定義'),
    b0('CER = （替換數 + 刪除數 + 插入數）÷ 正確答案字元總數 × 100%'),
    b0('概念：讓預測序列「變成」正確答案所需的最少編輯步驟數'),
    b0('計算方法：Levenshtein 距離（動態規劃）'),
    t(),
    h('三種錯誤類型（舉例說明）'),
    b0('替換（Substitution）：把錯的字改成對的字'),
    b1('預測：「國主此後」→ 正確：「國王此後」→ 1 次替換', color=C_GRAY),
    b0('刪除（Deletion）：刪掉多預測的字'),
    b1('預測：「籃球比賽場」→ 正確：「籃球比賽」→ 1 次刪除', color=C_GRAY),
    b0('插入（Insertion）：插入漏掉的字'),
    b1('預測：「球比賽」→ 正確：「籃球比賽」→ 1 次插入', color=C_GRAY),
    t(),
    h('解讀標準'),
    b0('CER = 38.0% → 每 10 個字平均有約 3.8 個錯', color=C_ORNG),
    b0('實際語音辨識系統目標：< 5–10%', color=C_GRAY),
    b0('越低越好；0% 表示完全正確', color=C_GRN),
])

# ══════════════════════════════════════════════════════
# 投影片 26：主要實驗結果
# ══════════════════════════════════════════════════════
mk('主要實驗結果（測試集 CER %，越低越好）', [
    h('Baseline（Transformer + CTC）'),
    b0('無任何增強（Base）：66.5%',                color=C_RED),
    b0('+ 頻譜相減：53.9%（↓9.4%）',               color=C_ORNG),
    b0('+ Mixup：49.6%（↓5.1%）',                  color=C_ORNG),
    b0('+ 有聲資料：51.4%（↓6.9%）',               color=C_ORNG),
    b0('完整增強（Final）：44.5%（共↓22.0%）',     color=C_GRN),
    t(),
    h('AuxCEMGR（+ 雙輔助任務）'),
    b0('無任何增強（Base）：62.5%',                color=C_ORNG),
    b0('完整增強（Final）：38.0% ← 最佳結果',
       bold=True, color=C_RED),
    b1('相比 Baseline Final 再降 6.5%（雙輔助任務的貢獻）', color=C_GRAY),
    t(),
    h('比較基準'),
    b0('Wadkins 2019（LSTM + CTC，完整增強）：66.8%', color=C_GRAY),
    b1('本論文 Baseline Final 已優於此 22.3%', color=C_GRAY),
    t(),
    h('三大關鍵貢獻排名（CER 改善量）'),
    b0('① 資料增強：~23.3%  ② Transformer 架構：~22%  ③ 雙輔助任務：6.5%',
       bold=True, color=C_RED),
])

# ══════════════════════════════════════════════════════
# 投影片 27：電極分析
# ══════════════════════════════════════════════════════
mk('深入分析 (1)：電極通道貢獻度', [
    h('單一通道貢獻度實驗'),
    b0('方法：每次只使用 1 個通道訓練 / 測試，觀察 CER'),
    b0('AuxCEMGR 最重要通道：CH4（下頜骨肌群）', bold=True, color=C_RED),
    b0('Baseline 最重要通道：CH2（嘴輪匝肌）'),
    b0('CH5、CH6（喉部）：兩個模型中效果均最差', color=C_ORNG),
    b1('原因：靜默模式下喉部幾乎不動，訊號趨近於零', color=C_GRAY),
    b1('對比：有聲 EMG 研究中 CH5 幾乎最重要 → 靜默/有聲差異顯著',
       color=C_GRAY),
    t(),
    h('電極數量 vs. 效能'),
    b0('電極對越多，CER 越低（8 對最佳）'),
    b0('移除任意 1 對：影響極小'),
    b0('移除超過 2 對：準確率急遽下降', bold=True, color=C_RED),
    b0('建議：實際應用可縮減至 6–7 對，兼顧舒適性', color=C_GRN),
    t(),
    h('整體結論'),
    b0('CH1–CH4（嘴部 / 下頜）> CH5–CH8（喉部 / 臉側）'),
    b0('AuxCEMGR 對電極數量的穩健性優於 Baseline'),
])

# ══════════════════════════════════════════════════════
# 投影片 28：句長 & Session 分析
# ══════════════════════════════════════════════════════
mk('深入分析 (2)：句長與 Session 穩定性', [
    h('句長對 CER 的影響'),
    b0('測試集按字數分成 6 組：≤9、10、11、12、13、14+ 字'),
    b0('趨勢：字數 ≥ 12 時，CER 明顯上升', bold=True, color=C_ORNG),
    b0('AuxCEMGR（Final）在所有字數組別中均優於其他設定'),
    b1('直觀解釋：句子越長，連續肌肉動作越複雜，模型負擔越重',
       color=C_GRAY),
    t(),
    h('Session 穩定性分析（5 個 session 的 CER 差異）'),
    b0('同一模型在不同 session 的 CER 有差異 → 電極位置影響確實存在'),
    b0('AuxCEMGR（Base，無增強）在各 session 間最穩定',
       bold=True, color=C_GRN),
    b1('→ 梯度反轉的 Session 分類有效消除電極位置偏差', color=C_GRAY),
    b0('AuxCEMGR（Final）和 Baseline（Final）均一致 → 資料增強也提升穩健性'),
    t(),
    h('靜默 vs. 有聲辨識比較'),
    b0('加入有聲訓練後，有聲 EMG 的辨識 CER 大幅下降'),
    b0('AuxCEMGR 在有聲 EMG 上依然優於 Baseline'),
])

# ══════════════════════════════════════════════════════
# 投影片 29：案例研究
# ══════════════════════════════════════════════════════
mk('深入分析 (3)：案例研究──模型輸出對比', [
    h('範例句子'),
    b0('正確答案：「國王此後連續（xù）命中三分」',
       bold=True, color=C_DARK),
    b1('（The King then hits back-to-back 3-pointers）', color=C_GRAY),
    t(),
    h('各模型輸出'),
    b0('Baseline（無增強）：「快船此後不命中三分」',
       color=C_RED),
    b1('問題：主詞完全辨識錯誤（快船 vs. 國王）', color=C_GRAY),
    b0('Baseline（完整增強）：「國王此後連度（dù）命中三分」',
       color=C_ORNG),
    b1('改善：主詞正確；「續」被認成「度」（發音部位相近）', color=C_GRAY),
    b0('AuxCEMGR（無增強）：「國王後頻繁投命中三分」',
       color=C_ORNG),
    b1('改善：主詞正確；仍有省略與錯字', color=C_GRAY),
    b0('AuxCEMGR（完整增強）：「國王此後連續命中三分」',
       bold=True, color=C_GRN),
    b1('完全正確！', color=C_GRN),
    t(),
    h('三個觀察'),
    b0('「分」「命」等字特別容易辨識：這些字的肌肉動作幅度較大'),
    b0('AuxCEMGR 輸出更合乎語法：session 無關特徵讓語意更合理'),
    b0('資料增強讓「連續命中三分」這類固定搭配被正確學習'),
])

# ══════════════════════════════════════════════════════
# 投影片 30：如何複製實驗
# ══════════════════════════════════════════════════════
mk('如何複製實驗：Step-by-Step', [
    h('Step 1：取得程式碼與資料'),
    b0('程式碼：git clone https://github.com/bluishwhite/EMG_ASR'),
    b0('資料：Email 申請（通訊作者：yinerwei1985@gmail.com）'),
    b1('申請內容：包含所有 EMG 訊號（silent + audible）+ 文字標籤', color=C_GRAY),
    t(),
    h('Step 2：環境安裝與前處理'),
    b0('pip install librosa scipy numpy torch'),
    b0('前處理流程：'),
    b1('Butterworth 帶通（10–400 Hz）：scipy.signal.butter() + filtfilt()'),
    b1('陷波濾波（50/150/250/350 Hz）：scipy.signal.iirnotch()'),
    b1('MFSC 提取（36 Mel filters）：librosa.feature.melspectrogram()'),
    t(),
    h('Step 3：訓練模型'),
    b0('框架：PyTorch'),
    b0('按投影片 24 的超參數配置（CNN + Transformer + CTC + 雙輔助任務）'),
    b0('建議順序：先訓練 Baseline → 加增強 → 加輔助任務，逐步驗證'),
    t(),
    h('Step 4：評估'),
    b0('計算測試集 CER（Levenshtein 距離）'),
    b0('對齊目標：Baseline Final ≈ 44.5%，AuxCEMGR Final ≈ 38.0%'),
])

# ══════════════════════════════════════════════════════
# 投影片 31：論文限制
# ══════════════════════════════════════════════════════
mk('論文的四大局限（= 未來研究的機會）', [
    h('局限 1：單一受試者（最關鍵）', color=C_RED),
    b0('整個資料集只有 1 位女性志願者錄製'),
    b0('EMG 訊號因人而異（電極阻抗、肌肉厚度、說話習慣）'),
    b0('→ 模型無法直接用於其他人（Cross-subject 泛化問題）', color=C_RED),
    t(),
    h('局限 2：封閉詞彙（OOV 問題）', color=C_RED),
    b0('訓練集只有 667 個字，測試集也在同一封閉集合內'),
    b0('出現訓練集沒有的字（Out-of-Vocabulary）→ 模型完全無法處理'),
    t(),
    h('局限 3：封閉語料庫（NBA 主題）', color=C_ORNG),
    b0('全部語料都是 NBA 賽事報導，用詞非常集中'),
    b0('→ 無法反映日常語言的多樣性'),
    t(),
    h('局限 4：CER 仍高（38%）', color=C_ORNG),
    b0('平均每 10 個字就有約 4 個錯誤'),
    b0('實際應用通常要求 < 5–10% CER，距離仍遠'),
    b0('未來改進：加入語言模型解碼、更多受試者、更大語料庫',
       color=C_GRN),
])

# ══════════════════════════════════════════════════════
# 投影片 32：對你研究的意義
# ══════════════════════════════════════════════════════
mk('對你研究的意義與改進方向', [
    h('這篇論文是你的直接比較基準（Baseline）'),
    b0('同樣是：臉部 sEMG + Transformer + 中文 SSR'),
    b0('你的目標：在類似架構上找到可改善 38% CER 的方法'),
    t(),
    h('你的差異化機會（論文局限 = 你的研究空間）'),
    b0('多受試者訓練 → 解決跨人泛化問題（最重要！）',
       color=C_GRN, bold=True),
    b0('聲調辨識輔助任務 → 利用普通話 4 聲特性（本文完全未處理）',
       color=C_GRN, bold=True),
    b0('語言模型整合解碼 → 改善 CTC 的語意合理性',
       color=C_GRN, bold=True),
    b0('使用 CSL-EMG 資料集 → 比較不同資料集對模型的影響',
       color=C_GRN, bold=True),
    t(),
    h('可直接參考的實作細節'),
    b0('電極配置：CH4（下頜骨肌）最重要 → 你的 MyoWare 感測器放置參考'),
    b0('前處理標準流程：10–400 Hz 帶通 + 50/150/250/350 Hz 陷波'),
    b0('超參數起始點：η₁ = η₂ = 1.0，α = 0.02，γ = 0.8'),
    b0('6–7 對電極已足夠，不必堅持 8 對（舒適度的考量）'),
])

# ══════════════════════════════════════════════════════
# 投影片 33：總結
# ══════════════════════════════════════════════════════
mk('總結', [
    h('本論文做了什麼'),
    b0('第一個中文 EMG-SSR 基準資料集（1,238 句，8 通道臉部 sEMG）'),
    b0('AuxCEMGR：Transformer + CTC + 拼音生成 + Session 分類（梯度反轉）'),
    b0('三種資料增強：頻譜相減 + 有聲 EMG 輔助 + Mixup'),
    b0('最佳 CER：38.0%（建立中文 SSR 的第一條基準線）'),
    t(),
    h('關鍵技術貢獻排名（CER 改善量）'),
    b0('① 資料增強（~23.3% CER）', bold=True, color=C_RED),
    b0('② Transformer 架構（~22% vs. LSTM）', bold=True, color=C_ORNG),
    b0('③ 雙輔助任務（6.5%）', bold=True, color=C_MID),
    t(),
    h('開放資源'),
    b0('程式碼：github.com/bluishwhite/EMG_ASR（公開）'),
    b0('資料集：向通訊作者 Email 申請（yinerwei1985@gmail.com）'),
    t(),
    h('一句話摘要'),
    b0('臉部 sEMG + Transformer + CTC + 拼音 / Session 輔助任務'
       ' → CER 38%，開啟中文無聲語音辨識研究之門',
       bold=True, color=C_DARK, size=Pt(20)),
])

# ══════════════════════════════════════════════════════
# 儲存
# ══════════════════════════════════════════════════════
out_dir = '/Users/rayopenclaw/ray-agent/papers'
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, 'AuxCEMGR_解析.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Total slides: {len(prs.slides)}')
