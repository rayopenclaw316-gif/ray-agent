from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import re, os

# ── constants ─────────────────────────────────────────────────────────────────
TITLE_COLOR = RGBColor(0x0E, 0x28, 0x41)
BLACK       = RGBColor(0x00, 0x00, 0x00)
RED         = RGBColor(0xFF, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

RED_TERMS = sorted([
    "sEMG", "EMG", "EEG", "SNN", "1D-CNN", "CNN", "BLE", "SSR", "SSI",
    "Siamese", "Transformer", "GRU", "LSTM", "Few-shot", "Trigram",
    "Moving Window", "PMLSF", "Sensor Fusion", "Language Model",
    "CDAN", "MMD", "Domain Adaptation",
], key=len, reverse=True)

IMG_DIR = "/Users/rayopenclaw/Downloads/ray-agent/papers/2603.11715_imgs"
OUT     = "/Users/rayopenclaw/Downloads/ray-agent/affect_decoding_emg_ppt.pptx"

# ── helpers ───────────────────────────────────────────────────────────────────
_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def _set_fonts(run, font_name):
    rPr = run._r.get_or_add_rPr()
    for tag in ["latin", "ea", "cs"]:
        el = rPr.find(f"{{{_NS}}}{tag}")
        if el is None:
            el = etree.SubElement(rPr, f"{{{_NS}}}{tag}")
        el.set("typeface", font_name)

def _split_lang(text):
    parts, last = [], 0
    for m in re.finditer(r'[一-鿿　-〿＀-￯]+', text):
        if m.start() > last:
            parts.append((False, text[last:m.start()]))
        parts.append((True, m.group()))
        last = m.end()
    if last < len(text):
        parts.append((False, text[last:]))
    return parts

def add_text(para, text, bold=False, color=BLACK, size_pt=17):
    segments = []
    remaining = text
    while remaining:
        ep, et = len(remaining), None
        for term in RED_TERMS:
            p = remaining.find(term)
            if p != -1 and p < ep:
                ep, et = p, term
        if et is None:
            segments.append((remaining, False)); break
        if ep > 0:
            segments.append((remaining[:ep], False))
        segments.append((et, True))
        remaining = remaining[ep + len(et):]
    for seg, is_red in segments:
        c = RED if is_red else color
        for is_cjk, chunk in _split_lang(seg):
            if not chunk: continue
            run = para.add_run()
            run.text = chunk
            run.font.bold = bold
            run.font.color.rgb = c
            run.font.size = Pt(size_pt)
            _set_fonts(run, "標楷體" if is_cjk else "Times New Roman")

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide

def title_box(slide, text, size_pt=28):
    tb = slide.shapes.add_textbox(Cm(2.3), Cm(1.0), Cm(29.2), Cm(3.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_text(p, text, bold=True, color=TITLE_COLOR, size_pt=size_pt)
    return tf

def content_box(slide):
    tb = slide.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(29.2), Cm(12.1))
    tf = tb.text_frame; tf.word_wrap = True
    return tf

def pagenum(slide, n):
    tb = slide.shapes.add_textbox(Cm(23.9), Cm(17.7), Cm(3.0), Cm(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(n)
    run.font.size = Pt(14); run.font.color.rgb = BLACK
    _set_fonts(run, "Times New Roman")

def para(tf, text, bold=False, color=BLACK, size_pt=17, sp=Pt(6), indent=0):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = sp
    p.level = indent
    add_text(p, text, bold=bold, color=color, size_pt=size_pt)
    return p

def img(slide, fname, l, t, w, h=None):
    path = os.path.join(IMG_DIR, fname)
    if h:
        slide.shapes.add_picture(path, Cm(l), Cm(t), Cm(w), Cm(h))
    else:
        slide.shapes.add_picture(path, Cm(l), Cm(t), Cm(w))

def divider_line(slide):
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    spTree = slide.shapes._spTree
    sp = etree.SubElement(spTree, qn('p:sp'))
    # minimal connector – use a thin rectangle as a horizontal rule
    ln = slide.shapes.add_textbox(Cm(2.3), Cm(4.6), Cm(29.2), Cm(0.05))
    ln.fill.solid()
    ln.fill.fore_color.rgb = TITLE_COLOR

# ── presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
# English title
tb = s.shapes.add_textbox(Cm(2.3), Cm(1.5), Cm(29.2), Cm(3.0))
tf = tb.text_frame; tf.word_wrap = True
p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
add_text(p0, "Affect Decoding in Phonated and Silent Speech Production from Surface EMG",
         bold=True, color=TITLE_COLOR, size_pt=25)

# Chinese title
tb2 = s.shapes.add_textbox(Cm(2.3), Cm(4.8), Cm(29.2), Cm(1.5))
tf2 = tb2.text_frame; p1 = tf2.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
add_text(p1, "從表面肌電圖解碼發音與無聲語音中的情感", bold=True, color=TITLE_COLOR, size_pt=21)

# Venue
tb3 = s.shapes.add_textbox(Cm(2.3), Cm(6.6), Cm(29.2), Cm(0.9))
tf3 = tb3.text_frame; p2 = tf3.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
add_text(p2, "arXiv:2603.11715 [eess.AS]  ·  Submitted March 2026", color=BLACK, size_pt=17)

# Authors
tb4 = s.shapes.add_textbox(Cm(2.3), Cm(7.9), Cm(29.2), Cm(2.6))
tf4 = tb4.text_frame; tf4.word_wrap = True
pa = tf4.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
add_text(pa, "Simon Pistrosch*, Kleanthis Avramidis*, Zhao Ren, Tiantian Feng, Jihwan Lee,", color=BLACK, size_pt=16)
pb = tf4.add_paragraph(); pb.alignment = PP_ALIGN.CENTER
add_text(pb, "Monica Gonzalez-Machorro, Anton Batliner, Tanja Schultz, Shrikanth Narayanan, Björn W. Schuller", color=BLACK, size_pt=16)
pc = tf4.add_paragraph(); pc.space_before = Pt(4); pc.alignment = PP_ALIGN.CENTER
add_text(pc, "(*Equal contribution)", color=BLACK, size_pt=14)

# Institutions
tb5 = s.shapes.add_textbox(Cm(2.3), Cm(11.0), Cm(29.2), Cm(2.5))
tf5 = tb5.text_frame; tf5.word_wrap = True
pi = tf5.paragraphs[0]; pi.alignment = PP_ALIGN.CENTER
add_text(pi, "TUM University Hospital · University of Southern California · University of Bremen", color=BLACK, size_pt=15)
pi2 = tf5.add_paragraph(); pi2.space_before = Pt(3); pi2.alignment = PP_ALIGN.CENTER
add_text(pi2, "Imperial College London · MCML Munich Center for Machine Learning", color=BLACK, size_pt=15)

# Date
tb6 = s.shapes.add_textbox(Cm(2.3), Cm(13.9), Cm(29.2), Cm(0.9))
tf6 = tb6.text_frame; pd = tf6.paragraphs[0]; pd.alignment = PP_ALIGN.CENTER
add_text(pd, "2026年4月24日", color=BLACK, size_pt=16)

pagenum(s, 1)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "目錄")
tf = content_box(s)

toc = [
    ("1. 研究背景與動機", "3"),
    ("2. 三大研究問題（RQ）", "4"),
    ("3. 資料集：ST-CASE", "5"),
    ("4. 實驗設計：三大任務", "6"),
    ("5. 硬體設備與電極配置", "7"),
    ("6. 訊號前處理流程", "8"),
    ("7. 特徵提取方法", "9"),
    ("8. 機器學習模型與評估指標", "10"),
    ("9. 結果 RQ1：sEMG 情感解碼能力", "11"),
    ("10. 結果 RQ2：發音 vs 無聲語音", "12"),
    ("11. 結果 RQ3：自發語音泛化", "13"),
    ("12. 消融分析：電極空間分布", "14"),
    ("13. 結論與研究意義", "15"),
]
for i, (title, pg) in enumerate(toc):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.space_before = Pt(5)
    p.alignment = PP_ALIGN.LEFT
    # left part
    add_text(p, f"{title}", size_pt=17)
    # right-aligned page number via tab
    run_tab = p.add_run(); run_tab.text = "\t"
    run_tab.font.size = Pt(17); _set_fonts(run_tab, "Times New Roman")
    run_pg = p.add_run(); run_pg.text = f"p.{pg}"
    run_pg.font.size = Pt(17); run_pg.font.color.rgb = BLACK
    _set_fonts(run_pg, "Times New Roman")

pagenum(s, 2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Background & Motivation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "研究背景與動機")
tf = content_box(s)

items = [
    ("情感（Affect）在語音溝通中的核心地位", True),
    ("除詞彙意義外，語音同時傳遞態度、禮貌、挫折感等副語言資訊", False),
    ("這些情感線索透過臉部、喉部、呼吸系統的協調肌肉活動來表達", False),
    ("聲學分析的局限性", True),
    ("傳統情感辨識依賴聲學訊號（韻律、頻譜）；然而在無聲語音介面（SSI）、喉切除術後助聲裝置、低音量環境中，聲學資訊不可用或受干擾", False),
    ("sEMG 的優勢與研究缺口", True),
    ("sEMG（表面肌電圖）：非侵入式記錄語音生產時的肌肉電訊號，已被用於解碼音素、無聲詞彙辨識（SSR）及重建語音", False),
    ("然而，現有 sEMG 研究多聚焦於語言內容（linguistic content），極少探索情感調變（affective modulation）是否編碼於周邊肌肉活動中", False),
    ("本研究的切入點", True),
    ("探索情感資訊是否嵌入於發音執行的神經肌肉過程，且在無聲發音時是否仍可保留，為情感感知型無聲語音介面（affect-aware SSI）奠基", False),
]
for i, (text, bold) in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3) if not bold else Pt(8)
    p.level = 0 if bold else 1
    add_text(p, ("▶ " if bold else "• ") + text, bold=bold, size_pt=17 if bold else 16)

pagenum(s, 3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Research Questions
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "三大研究問題（Research Questions）")
tf = content_box(s)

rqs = [
    ("RQ1：sEMG 情感解碼可行性",
     "臉部與頸部 sEMG 在語音生產期間，能否有效解碼情感狀態（挫折 vs 禮貌 vs 中性）？\n→ 評估受試者內（intra-subject）與跨受試者（inter-subject）兩種情境"),
    ("RQ2：發音模式的影響",
     "情感解碼在發音語音（phonated speech）與無聲語音（silent speech）之間有何差異？\n→ 比較兩種發音模式的解碼性能，並測試跨模式遷移（cross-mode transfer）"),
    ("RQ3：實驗情境的影響",
     "從受控朗讀任務（prompted/scripted）到自發對話（spontaneous）情境，情感肌肉訊號的泛化能力如何？\n→ 評估模型從受控任務訓練後在自發語音測試集上的表現"),
]
for i, (hd, body) in enumerate(rqs):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(14)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=18)
    for line in body.split("\n"):
        p2 = tf.add_paragraph(); p2.space_before = Pt(4); p2.level = 1
        add_text(p2, line, size_pt=16)

pagenum(s, 4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Dataset ST-CASE (with conceptual figure)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "資料集：ST-CASE（SAIL-TUM Corpus on Affective Speech & EMG）")

# Text on left
tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(12.5))
tf = tb.text_frame; tf.word_wrap = True

para_data = [
    ("基本資訊", True),
    ("受試者：N = 12（9 位女性），平均年齡 26.2 歲（SD=5.2，範圍 20-36 歲）", False),
    ("語言：母語英語 5 位，其餘具 C2 英語程度", False),
    ("錄音地點：德國慕尼黑 TUM 大學醫院隔音室", False),
    ("資料規模", True),
    ("總發話數：2,780 則（發音 1,588 則 / 無聲 1,192 則）", False),
    ("情感標籤：挫折 1,143 / 中性 479 / 禮貌 1,158", False),
    ("總錄音時長：約 194 分鐘（平均每則 4.19 ± 2.81 秒）", False),
    ("多模態記錄", True),
    ("8 通道臉部與頸部 sEMG + 音訊（Rode NT1-A 麥克風）", False),
    ("資料不公開（GDPR 隱私保護）；可透過 Article 46 條款申請", False),
]
for i, (text, bold) in enumerate(para_data):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3) if not bold else Pt(8)
    p.level = 0 if bold else 1
    add_text(p, ("▶ " if bold else "• ") + text, bold=bold, size_pt=16 if bold else 15)

# Image on right
img(s, "page1_img1.png", 20.3, 4.9, 12.5, 12.0)
pagenum(s, 5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Experimental Design (with annotation figure)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "實驗設計：三大任務")

# Left text
tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.8), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

tasks = [
    ("Task 1：受控朗讀任務（Prompted Reading）", True, [
        "50 句（公寓搜尋主題）：中性×10、禮貌措辭×10、挫折措辭×10、加重複的中性句×20",
        "每句先發音朗讀，再無聲發音 → 100 次試驗/受試者",
        "情感標籤：中性 / 禮貌 / 挫折",
    ]),
    ("Task 2：自發對話（Wizard-of-Oz）", True, [
        "場景：汽車保險電話對話，受試者不知對方為 AI 控制",
        "子任務 2A（禮貌誘導）：申請折扣、登記車輛、申請綠卡",
        "子任務 2B（挫折誘導）：追討未付報價單，代理人故意難溝通",
        "3 位標注者以 5 點 Likert 量表評定每則發話的挫折/禮貌程度",
    ]),
    ("Task 3：重複朗讀任務（約 30 分鐘後）", True, [
        "與 Task 1 完全相同，用於評估跨時間穩定性與學習效應",
    ]),
]
first = True
for hd, bold, bullets in tasks:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(10)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=17)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(3); pb.level = 1
        add_text(pb, "• " + b, size_pt=15)

# Right: annotation figure
img(s, "page3_img1.png", 20.5, 5.1, 12.0, 7.5)

# Caption
tb_cap = s.shapes.add_textbox(Cm(20.5), Cm(12.9), Cm(12.0), Cm(1.2))
tfc = tb_cap.text_frame
pc = tfc.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
add_text(pc, "Task 2A（禮貌）/2B（挫折）標注分布\nKrippendorff's α：挫折=0.50、禮貌=0.55",
         color=BLACK, size_pt=13)

pagenum(s, 6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Sensor Setup & Electrode Placement
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "硬體設備與電極配置")

tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

hw_items = [
    ("放大器設備", True, [
        "actiCHamp Plus（Brain Products GmbH），8 通道雙極 Ag/AgCl 電極",
        "採樣率 10 kHz → 後處理降至 1 kHz；接地電極置於鼻骨末端",
        "電極阻抗 < 100 kΩ，使用 Ten20 導電膠",
    ]),
    ("音訊設備", True, [
        "Rode NT1-A 麥克風，48 kHz 採樣，Focusrite Scarlett 2i2 介面",
        "軟體觸發同步 sEMG 與音訊起始/結束標記",
    ]),
    ("電極位置（8 通道）", True, [
        "E1 Infrahyoid（舌骨下肌）：語音相關，頸部甲狀軟骨外側",
        "E2 Suprahyoid（舌骨上肌）：語音相關，下顎與舌骨中間",
        "E3 Mylohyoid（下頜舌骨肌）：語音相關，頦部下方",
        "E4 Mentalis（頦肌）：語音 + 情感相關，頦部中央外側",
        "E5 Orbicularis Oris Superioris（上唇口輪匝肌）：語音相關",
        "E6 Depressor Supercilii（降眉間肌）：情感相關，眉間",
        "E7/E8 Zygomaticus Major（顴大肌）左右：語音 + 情感相關",
    ]),
]
first = True
for hd, bold, bullets in hw_items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=17)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(3); pb.level = 1
        add_text(pb, "• " + b, size_pt=15)

# Right: use RQ1 electrode map as illustration
img(s, "page7_img1.png", 20.3, 5.1, 12.5, 10.5)
tb_cap2 = s.shapes.add_textbox(Cm(20.3), Cm(15.8), Cm(12.5), Cm(1.0))
tfc2 = tb_cap2.text_frame
pc2 = tfc2.paragraphs[0]; pc2.alignment = PP_ALIGN.CENTER
add_text(pc2, "電極通道 AUC 熱圖（受試者內/跨受試者）", color=BLACK, size_pt=13)

pagenum(s, 7)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Signal Preprocessing
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "訊號前處理流程")
tf = content_box(s)

steps = [
    ("sEMG 前處理管線", True, [
        "高通濾波：4 階 Butterworth，截止頻率 100 Hz（去除低頻漂移）",
        "陷波濾波：50 Hz 工頻干擾 + 多達 8 個諧波迭代去除",
        "降採樣：抗鋸齒低通濾波後 IIR 降至 1 kHz",
        "異常值截斷：±10 個標準差以外數值截除",
        "基準正規化：每個任務均有靜息基準錄音，計算中位數與四分位距（IQR），對訊號進行穩健縮放（robust scaling）",
        "時間裁切：每則發話起始 0.56 秒、結尾 0.54 秒設為保守性裁切邊界",
    ]),
    ("音訊前處理管線", True, [
        "降噪：noisereduce Python 套件去除環境雜音",
        "正規化：峰值 RMS 正規化（目標 RMS = 0.5）",
        "降採樣：至 16 kHz",
        "Task 2 自動轉錄：Whisper-small 取得發話文字（本研究未使用轉錄結果）",
    ]),
]
first = True
for hd, bold, bullets in steps:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=18)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(4); pb.level = 1
        add_text(pb, "• " + b, size_pt=16)

pagenum(s, 8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Feature Extraction
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "特徵提取方法")
tf = content_box(s)

feat_sections = [
    ("手工統計特徵（Structural Features，92 維）", True, [
        "每通道整流訊號 x̃_c = |x_c| 的均值、標準差、變異係數（CoV）",
        "峰值振幅、均方根值（RMS）：√(Σx²_c / T)",
        "功率頻譜密度（Welch 法）→ 中位頻率 f_med,c = ½∫PSD_c(f)df",
        "頻譜熵 H_spec,c = -Σ p_{c,k} log p_{c,k}（量化頻率分布均勻度）",
        "跨通道 Pearson 相關（整流訊號），最終編碼為 92 維向量",
    ]),
    ("TD-0 時域特徵（Time-Domain Features）", True, [
        "三角形濾波器（截止 134 Hz）分解低頻 / 高頻分量（雙移動平均）",
        "分幀：27 ms 矩形窗，10 ms 步長",
        "每幀計算：低頻均方、低頻均值、高頻均方、高頻絕對均值、過零率（ZCR）",
        "跨幀統計：均值、標準差、第 0/25/75/100 百分位數",
    ]),
    ("BioCodec 嵌入特徵（深度學習基礎模型）", True, [
        "BioCodec：開源 sEMG 神經編解碼器，原為腕部手勢預訓練",
        "提取編碼器輸出：通道逐一 128 維嵌入（量化模組前一層）",
        "零樣本（zero-shot）使用於臉部/頸部 sEMG，驗證重建品質後採用",
    ]),
    ("語音特徵（對照組）", True, [
        "eGeMAPSv02（openSMILE）：88 個韻律/聲學描述子",
        "Vox-Profile（256 維嵌入）：Whisper-Large 微調於 MSP-Podcast，輸出情感維度（喚醒度、效價、支配度）",
    ]),
]
first = True
for hd, bold, bullets in feat_sections:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(7)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=17)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(3); pb.level = 1
        add_text(pb, "• " + b, size_pt=15)

pagenum(s, 9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ML Models & Evaluation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "機器學習模型與評估指標")
tf = content_box(s)

ml_items = [
    ("分類器", True, [
        "SVM（支持向量機）RBF 核：用於手工特徵（訓練前 z-score 正規化）",
        "線性探針分類器（L2 正則化）：用於嵌入特徵（訓練前 robust scaler 正規化）",
        "二元分類：挫折 vs 禮貌（排除中性樣本以降低模糊性）",
    ]),
    ("評估策略", True, [
        "受試者內（Intra-subject）：5-fold 交叉驗證（以句子為單位分割，防止同一句子出現在訓練/測試集）",
        "跨受試者（Inter-subject）：留一法（LOSO）外迴圈 + 5-fold 內迴圈",
        "RQ3 自發語音：以 Task 1 & 3 訓練，在 Task 2 測試（完全排除測試受試者）",
    ]),
    ("評估指標", True, [
        "AUC（ROC 曲線下面積）：主要指標，不受類別不平衡影響，0.5=隨機，1.0=完美",
        "BAC（平衡準確率）：各類別召回率的算術平均，適用類別不平衡情境",
        "報告方式：各受試者之均值 ± 標準差",
        "顯著性檢定：McNemar 檢定（跨模態比較）、Wilcoxon 符號秩檢定（任務間比較，FDR 校正）",
    ]),
]
first = True
for hd, bold, bullets in ml_items:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=18)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(4); pb.level = 1
        add_text(pb, "• " + b, size_pt=16)

pagenum(s, 10)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RQ1 Results
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "結果 RQ1：sEMG 情感解碼能力")

# Left text
tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
add_text(p, "受試者內（Intra-subject）結果", bold=True, color=TITLE_COLOR, size_pt=17)

table_rows = [
    ("模態", "特徵", "BAC", "AUC"),
    ("sEMG", "Structural", "0.749±0.075", "0.820±0.081"),
    ("sEMG", "TD-0", "0.762±0.063", "0.845±0.058 ★"),
    ("sEMG", "BioCodec", "0.721±0.053", "0.792±0.075"),
    ("語音", "eGeMAPS", "0.610±0.125", "0.644±0.168"),
    ("語音", "Vox-Profile", "0.659±0.097", "0.732±0.104"),
]
for i, row in enumerate(table_rows):
    pb = tf.add_paragraph(); pb.space_before = Pt(3 if i > 0 else 6)
    txt = "  ".join(f"{v:<18}" for v in row)
    bold_row = i == 0 or i == 2
    add_text(pb, txt, bold=bold_row, size_pt=14 if i == 0 else 15)

pb2 = tf.add_paragraph(); pb2.space_before = Pt(10)
add_text(pb2, "★ TD-0 達最高 AUC=0.845（受試者內）", bold=False, color=RED, size_pt=16)

pb3 = tf.add_paragraph(); pb3.space_before = Pt(8)
add_text(pb3, "跨受試者（Inter-subject）結果", bold=True, color=TITLE_COLOR, size_pt=17)

interrows = [
    ("sEMG TD-0", "AUC=0.567（接近隨機）"),
    ("語音 Vox-Profile", "AUC=0.657（最佳跨受試者）"),
]
for ir in interrows:
    pb4 = tf.add_paragraph(); pb4.space_before = Pt(4); pb4.level = 1
    add_text(pb4, f"• {ir[0]}：{ir[1]}", size_pt=16)

pb5 = tf.add_paragraph(); pb5.space_before = Pt(8)
add_text(pb5, "關鍵發現：受試者內情感調變主要透過肌肉活動表現，而非聲學差異（McNemar p<0.001）", bold=True, size_pt=16)

pagenum(s, 11)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Lexical Confound Control + Task comparison (with fig)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "詞彙混淆控制：Unique vs Repeated 句子分析")

tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
add_text(p, "設計動機", bold=True, color=TITLE_COLOR, size_pt=17)
pb = tf.add_paragraph(); pb.space_before = Pt(4); pb.level = 1
add_text(pb, "• 「相同句子以不同情感朗讀」：若模型依靠詞彙線索（lexical cues），重複句測試時 AUC 應顯著下滑", size_pt=15)

pb2 = tf.add_paragraph(); pb2.space_before = Pt(10)
add_text(pb2, "受試者內 AUC（Unique vs Repeated）", bold=True, color=TITLE_COLOR, size_pt=17)

rows2 = [
    ("特徵", "Unique AUC", "Repeated AUC"),
    ("sEMG Structural", "0.856±0.069", "0.720±0.159"),
    ("sEMG TD-0", "0.824±0.126", "0.751±0.145"),
    ("sEMG BioCodec", "0.799±0.058", "0.747±0.145"),
    ("語音 eGeMAPS", "0.643±0.060", "0.559±0.132"),
    ("語音 Vox-Profile", "0.889±0.075", "0.469±0.202 ★"),
]
for i, row in enumerate(rows2):
    pb3 = tf.add_paragraph(); pb3.space_before = Pt(3)
    add_text(pb3, "  ".join(f"{v:<22}" for v in row), bold=(i == 0), size_pt=14 if i == 0 else 15)

pb4 = tf.add_paragraph(); pb4.space_before = Pt(8)
add_text(pb4, "★ Vox-Profile 在重複句崩潰至隨機水準（AUC=0.469）→ 主要依賴詞彙資訊", color=RED, size_pt=15)
pb5 = tf.add_paragraph(); pb5.space_before = Pt(4)
add_text(pb5, "sEMG 在重複句仍保持 AUC>0.7，顯示肌肉活動確實編碼情感而非詞彙", bold=True, size_pt=16)

img(s, "page6_img1.png", 20.3, 5.1, 12.5, 10.5)
tb_cap = s.shapes.add_textbox(Cm(20.3), Cm(15.8), Cm(12.5), Cm(1.0))
pc = tb_cap.text_frame.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
add_text(pc, "Task1 vs Task3 受試者內 AUC（左：全句；右：重複句）", color=BLACK, size_pt=13)

pagenum(s, 12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RQ2: Phonated vs Silent
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "結果 RQ2：發音語音 vs 無聲語音的情感解碼")

tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
add_text(p, "模式內（Intra-mode）解碼性能（受試者內 AUC）", bold=True, color=TITLE_COLOR, size_pt=17)

rows3 = [
    ("特徵", "Phonated→Phonated", "Silent→Silent"),
    ("Structural", "0.815±0.110", "0.829±0.056"),
    ("TD-0", "0.806±0.113", "0.811±0.100"),
    ("BioCodec", "0.758±0.066", "0.792±0.102"),
]
for i, row in enumerate(rows3):
    pb = tf.add_paragraph(); pb.space_before = Pt(3)
    add_text(pb, "  ".join(f"{v:<26}" for v in row), bold=(i == 0), size_pt=14 if i == 0 else 15)

pb4 = tf.add_paragraph(); pb4.space_before = Pt(6)
add_text(pb4, "→ 無聲語音解碼性能與發音語音相當（差距 < 2%），無顯著統計差異", bold=True, size_pt=16)

pb5 = tf.add_paragraph(); pb5.space_before = Pt(12)
add_text(pb5, "跨模式（Cross-mode）遷移（受試者內 AUC）", bold=True, color=TITLE_COLOR, size_pt=17)

rows4 = [
    ("特徵", "Phonated→Silent", "Silent→Phonated"),
    ("Structural", "0.707±0.158", "0.663±0.168"),
    ("TD-0", "0.705±0.118", "0.626±0.145"),
    ("BioCodec", "0.763±0.094", "0.745±0.075 ★"),
]
for i, row in enumerate(rows4):
    pb = tf.add_paragraph(); pb.space_before = Pt(3)
    add_text(pb, "  ".join(f"{v:<28}" for v in row), bold=(i == 0), size_pt=14 if i == 0 else 15)

pb6 = tf.add_paragraph(); pb6.space_before = Pt(6)
add_text(pb6, "★ BioCodec 雙向遷移最佳；Structural/TD-0 從 Phonated→Silent 效果優於反向", color=RED, size_pt=15)
pb7 = tf.add_paragraph(); pb7.space_before = Pt(4)
add_text(pb7, "→ 可以僅用發音語音訓練，部署於無聲語音介面（SSI）", bold=True, size_pt=16)

img(s, "page7_img2.png", 20.3, 5.1, 12.5, 10.5)
tb_cap = s.shapes.add_textbox(Cm(20.3), Cm(15.8), Cm(12.5), Cm(1.0))
pc = tb_cap.text_frame.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
add_text(pc, "通道 AUC 熱圖：發音（左） vs 無聲（右）", color=BLACK, size_pt=13)

pagenum(s, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — RQ3: Spontaneous Speech
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "結果 RQ3：自發語音的跨受試者泛化")

tb = s.shapes.add_textbox(Cm(2.3), Cm(5.1), Cm(17.5), Cm(13.0))
tf = tb.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
add_text(p, "Task 2 跨受試者挫折偵測（Inter-subject AUC）", bold=True, color=TITLE_COLOR, size_pt=17)

rows5 = [
    ("模態", "特徵", "BAC", "AUC"),
    ("sEMG", "Structural", "0.616±0.003", "0.623±0.005"),
    ("sEMG", "TD-0", "0.527±0.006", "0.518±0.005"),
    ("sEMG", "BioCodec", "0.595±0.014", "0.630±0.009"),
    ("語音", "eGeMAPS", "0.607±0.002", "0.679±0.003"),
    ("語音", "Vox-Profile", "0.670±0.006", "0.743±0.004 ★"),
]
for i, row in enumerate(rows5):
    pb = tf.add_paragraph(); pb.space_before = Pt(3)
    add_text(pb, "  ".join(f"{v:<20}" for v in row), bold=(i == 0), size_pt=14 if i == 0 else 15)

pb6 = tf.add_paragraph(); pb6.space_before = Pt(8)
add_text(pb6, "★ 自發語音中語音模態（Vox-Profile AUC=0.743）優於 sEMG", color=RED, size_pt=15)
pb7 = tf.add_paragraph(); pb7.space_before = Pt(4)
add_text(pb7, "→ 自發語音提供更豐富的聲學情感線索（韻律變化更自然）", size_pt=16)

pb8 = tf.add_paragraph(); pb8.space_before = Pt(8)
add_text(pb8, "空間分布的變化", bold=True, color=TITLE_COLOR, size_pt=17)
for bullet in [
    "受控任務：臉部通道（E6 眉間、E7/E8 顴骨）分辨力最高",
    "自發對話：額面通道性能下降，下顎與頸部通道相對較好",
    "→ 自發語音情感的神經肌肉模式與受控任務存在空間重分布（spatial redistribution）",
]:
    pb9 = tf.add_paragraph(); pb9.space_before = Pt(4); pb9.level = 1
    add_text(pb9, "• " + bullet, size_pt=15)

img(s, "page8_img1.png", 20.3, 5.1, 12.5, 10.5)
tb_cap = s.shapes.add_textbox(Cm(20.3), Cm(15.8), Cm(12.5), Cm(1.0))
pc = tb_cap.text_frame.paragraphs[0]; pc.alignment = PP_ALIGN.CENTER
add_text(pc, "通道 AUC 熱圖：受控（左） vs 自發（右）", color=BLACK, size_pt=13)

pagenum(s, 14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
title_box(s, "結論與研究意義")
tf = content_box(s)

conclusions = [
    ("主要發現", True, [
        "sEMG 可有效解碼情感：受試者內 AUC 高達 0.845（TD-0），優於聲學特徵（McNemar p<0.001）",
        "情感訊號在無聲語音中持續存在：無聲模式解碼性能與發音模式相當（差距 <2%），且 Task 3 學習效應顯著改善無聲解碼",
        "詞彙無關性確認：在重複句（相同句子、不同情感）上，sEMG AUC 仍維持 >0.7，而語音 Vox-Profile 崩潰至隨機水準",
        "E6（降眉間肌）為最具分辨力的通道：在所有設定中始終貢獻最高 AUC，情感相關肌肉活動具空間特異性",
        "BioCodec 嵌入在跨域場景表現最穩健（重複句、跨模式、跨受試者）",
    ]),
    ("對無聲語音介面（SSI）的意義", True, [
        "可行性驗證：affect-aware SSI 在技術上可行，甚至不需要發音語音錄製（cross-mode 遷移成立）",
        "挑戰：跨受試者泛化仍是主要瓶頸，12 人樣本無法建立全局情感標記；需更大規模資料集",
        "未來方向：更大樣本、更生態效度的實驗設計；探索情感調變與特定韻律變化的耦合機制",
    ]),
    ("研究限制", True, [
        "樣本量小（N=12）且性別不平衡；情感為人工誘導而非自然發生",
        "無法分離發音運動調變與伴隨語音的臉部表情",
    ]),
]
first = True
for hd, bold, bullets in conclusions:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(8)
    add_text(p, hd, bold=True, color=TITLE_COLOR, size_pt=18)
    for b in bullets:
        pb = tf.add_paragraph(); pb.space_before = Pt(4); pb.level = 1
        add_text(pb, "• " + b, size_pt=15)

pagenum(s, 15)

# ── save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
